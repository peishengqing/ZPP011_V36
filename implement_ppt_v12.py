#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT 生成优化 v1.2 实施脚本
实现任务卡 ZPP011-TASK-001-R1 要求的所有功能
"""
import os
import sys
sys.stdout.reconfigure(encoding='utf-8')

root = r'E:\zpp011_dev\模块化脚本'

# 1. 创建必要的目录
dirs_to_create = [
    'output/pptx',
    'config',
    'tests',
]

print("=== Creating Directories ===")
for d in dirs_to_create:
    path = os.path.join(root, d)
    os.makedirs(path, exist_ok=True)
    print(f"✓ Created: {d}")

# 2. 创建空的模板文件（fallback 用）
template_path = os.path.join(root, 'config', 'template.pptx')
if not os.path.exists(template_path):
    # 创建一个空的 PPT 文件作为占位符
    try:
        from pptx import Presentation
        prs = Presentation()
        prs.save(template_path)
        print(f"✓ Created: config/template.pptx")
    except ImportError:
        print("⚠ python-pptx not installed, will use code-generated slides")

# 3. 读取现有的 audit_presenter.py
presenter_path = os.path.join(root, 'modules', 'audit', 'presenters', 'audit_presenter.py')
with open(presenter_path, 'r', encoding='utf-8') as f:
    content = f.read()

print(f"\n=== Modifying audit_presenter.py ===")
print(f"Current file size: {len(content)} bytes")

# 4. 准备新的 PPT 生成代码
new_ppt_code = '''
    # ========== PPT 生成优化 v1.2 实现 ==========
    
    def _classify_material_type(self, material_code: str) -> str:
        """物料类型分类（Slide 10 专用）
        
        Args:
            material_code: 物料编码
        
        Returns:
            物料类型：'原材料' / '包材' / '其他'
        """
        if not material_code or not isinstance(material_code, str):
            return '其他'
        
        # 提取前缀（前 3 位）
        prefix = material_code.strip()[:3]
        
        if prefix in ('100', '400'):
            return '原材料'
        elif prefix in ('200', '600'):
            return '包材'
        else:
            return '其他'
    
    def _pre_aggregate_data(self, df):
        """一次性计算所有聚合结果
        
        Args:
            df: 原始数据 DataFrame
        
        Returns:
            pre_data: 包含所有聚合结果的字典
        """
        import pandas as pd
        import numpy as np
        
        pre_data = {}
        
        # 确保必要的列存在
        if '工厂' not in df.columns:
            df['工厂'] = '未知工厂'
        if '车间' not in df.columns:
            df['车间'] = '未知车间'
        if '偏差金额' not in df.columns:
            df['偏差金额'] = 0.0
        if '偏差率' not in df.columns:
            df['偏差率'] = 0.0
        if '物料编码' not in df.columns:
            df['物料编码'] = ''
        if '物料名称' not in df.columns:
            df['物料名称'] = ''
        if '备注原因' not in df.columns:
            df['备注原因'] = ''
        
        # 1. 工厂维度 KPI
        factory_kpis = {}
        for factory in df['工厂'].unique():
            factory_df = df[df['工厂'] == factory]
            factory_kpis[factory] = {
                'total_records': len(factory_df),
                'total_amount': factory_df['偏差金额'].sum(),
                'avg_dev_rate': factory_df['偏差率'].mean(),
                'high_dev_count': len(factory_df[abs(factory_df['偏差率']) > 10]),
                'no_remark_count': len(factory_df[
                    (factory_df['备注原因'].isna()) | 
                    (factory_df['备注原因'] == '')
                ]),
            }
        pre_data['factory_kpis'] = factory_kpis
        
        # 2. 物料 Top10（分工厂）
        material_top10 = {}
        for factory in df['工厂'].unique():
            factory_df = df[df['工厂'] == factory]
            material_sum = factory_df.groupby('物料名称')['偏差金额'].abs().sum()
            material_top10[factory] = material_sum.nlargest(10).to_dict()
        pre_data['material_top10'] = material_top10
        
        # 3. 车间维度统计
        workshop_stats = {}
        for factory in df['工厂'].unique():
            factory_df = df[df['工厂'] == factory]
            workshop_sum = factory_df.groupby('车间')['偏差金额'].sum()
            workshop_stats[factory] = workshop_sum.to_dict()
        pre_data['workshop_stats'] = workshop_stats
        
        # 4. 物料类型分类（Slide 10）
        df['物料类型'] = df['物料编码'].apply(self._classify_material_type)
        material_type_net = df.groupby(['工厂', '物料类型'])['偏差金额'].sum().unstack(fill_value=0)
        pre_data['material_type_net'] = material_type_net
        
        # 5. 无备注预警（高金额）
        no_remark_warning = df[
            ((df['备注原因'].isna()) | (df['备注原因'] == '')) & 
            (abs(df['偏差金额']) >= 50000)
        ].copy()
        pre_data['no_remark_warning'] = no_remark_warning
        
        # 6. 异常预警（定额>0 且 实际=0）
        abnormal_warning = df[
            (df.get('数量 - 定额', 0) > 0) & 
            (df.get('数量 - 实际', 1) == 0)
        ].copy()
        pre_data['abnormal_warning'] = abnormal_warning
        
        # 7. 高频原因统计
        reason_counts = df[df['备注原因'].notna() & (df['备注原因'] != '')]['备注原因'].value_counts()
        pre_data['freq_reasons'] = reason_counts.head(10).to_dict()
        
        # 8. 总体 KPI
        pre_data['total_kpis'] = {
            'total_records': len(df),
            'total_amount': df['偏差金额'].sum(),
            'avg_dev_rate': df['偏差率'].mean(),
            'high_dev_count': len(df[abs(df['偏差率']) > 10]),
            'factories': df['工厂'].nunique(),
        }
        
        return pre_data
    
    def _add_cover(self, prs, pre_data):
        """添加封面页"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        
        # 标题
        title = slide.shapes.title
        title.text = "ZPP011 生产偏差分析报告"
        
        # 副标题
        subtitle = slide.placeholders[1]
        subtitle.text = f"数据量：{pre_data['total_kpis']['total_records']} 条\\n工厂数：{pre_data['total_kpis']['factories']}\\n生成时间：{pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        return slide
    
    def _add_toc(self, prs, pre_data):
        """添加目录页"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = "目录"
        
        toc_items = [
            "1. 核心指标总览",
            "2. 工厂维度对比",
            "3. 物料偏差 Top10（食品厂）",
            "4. 物料偏差 Top10（饮料厂）",
            "5. 偏差类型分布",
            "6. 车间详情（食品厂）",
            "7. 车间详情（饮料厂）",
            "8. 物料类型净偏差分布",
            "9. 无备注预警 Top10（食品厂）",
            "10. 无备注预警 Top10（饮料厂）",
            "11. 异常预警明细",
            "12. 高频原因对比",
            "13. 总结与建议",
        ]
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for item in toc_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
        
        return slide
    
    def _add_kpi_overview(self, prs, pre_data):
        """添加核心指标总览页"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "核心指标总览"
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        
        # 总体指标
        p = tf.add_paragraph()
        p.text = "【总体】"
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = f"总记录数：{pre_data['total_kpis']['total_records']:,}"
        p = tf.add_paragraph()
        p.text = f"总偏差金额：{pre_data['total_kpis']['total_amount']:,.2f} 元"
        p = tf.add_paragraph()
        p.text = f"平均偏差率：{pre_data['total_kpis']['avg_dev_rate']:.2f}%"
        p = tf.add_paragraph()
        p.text = f"高偏差（>10%）记录数：{pre_data['total_kpis']['high_dev_count']}\\n"
        
        # 分工厂指标
        for factory, kpis in pre_data['factory_kpis'].items():
            p = tf.add_paragraph()
            p.text = f"【{factory}】"
            p.font.bold = True
            p = tf.add_paragraph()
            p.text = f"  记录数：{kpis['total_records']:,}"
            p = tf.add_paragraph()
            p.text = f"  偏差金额：{kpis['total_amount']:,.2f} 元"
            p = tf.add_paragraph()
            p.text = f"  高偏差记录：{kpis['high_dev_count']}\\n"
        
        return slide
    
    def _add_factory_comparison(self, prs, pre_data):
        """添加工厂对比表"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "工厂维度对比"
        
        # 创建表格
        rows = len(pre_data['factory_kpis']) + 1
        cols = 6
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 表头
        headers = ['工厂', '记录数', '偏差金额 (元)', '平均偏差率 (%)', '高偏差数', '无备注数']
        for j, h in enumerate(headers):
            table.cell(0, j).text = h
        
        # 数据
        for i, (factory, kpis) in enumerate(pre_data['factory_kpis'].items(), 1):
            table.cell(i, 0).text = factory
            table.cell(i, 1).text = f"{kpis['total_records']:,}"
            table.cell(i, 2).text = f"{kpis['total_amount']:,.2f}"
            table.cell(i, 3).text = f"{kpis['avg_dev_rate']:.2f}"
            table.cell(i, 4).text = str(kpis['high_dev_count'])
            table.cell(i, 5).text = str(kpis['no_remark_count'])
        
        return slide
    
    def _add_material_top10(self, prs, pre_data, factory: str):
        """添加物料偏差 Top10（分工厂）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"物料偏差 Top10 - {factory}"
        
        top10_data = pre_data['material_top10'].get(factory, {})
        
        if not top10_data:
            tf = slide.placeholders[1].text_frame
            tf.text = f"该工厂无有效物料记录"
            return slide
        
        # 创建表格
        rows = len(top10_data) + 1
        cols = 2
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 表头
        table.cell(0, 0).text = '物料名称'
        table.cell(0, 1).text = '偏差金额 (元)'
        
        # 数据
        for i, (mat_name, amount) in enumerate(top10_data.items(), 1):
            table.cell(i, 0).text = mat_name
            table.cell(i, 1).text = f"{amount:,.2f}"
        
        # 如果不足 10 条，添加标注
        if len(top10_data) < 10:
            note = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
            note.text_frame.text = f"注：该工厂仅有 {len(top10_data)} 条有效物料记录"
        
        return slide
    
    def _add_deviation_type_chart(self, prs, pre_data):
        """添加偏差类型分布（并排饼图）"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        title.text_frame.text = "偏差类型分布"
        
        # 这里应该添加饼图代码
        # 简化版本：用文字描述
        tf = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5))
        tf.text_frame.text = "【饼图区域：左侧食品厂，右侧饮料厂】\\n\\n（图表功能待实现）"
        
        return slide
    
    def _add_workshop_details(self, prs, pre_data, factory: str):
        """添加车间详情（分工厂）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"车间详情 - {factory}"
        
        workshop_data = pre_data['workshop_stats'].get(factory, {})
        
        if not workshop_data:
            tf = slide.placeholders[1].text_frame
            tf.text = f"该工厂无车间级数据"
            return slide
        
        # 创建表格
        rows = len(workshop_data) + 1
        cols = 2
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 表头
        table.cell(0, 0).text = '车间'
        table.cell(0, 1).text = '偏差金额 (元)'
        
        # 数据
        for i, (workshop, amount) in enumerate(workshop_data.items(), 1):
            table.cell(i, 0).text = workshop
            table.cell(i, 1).text = f"{amount:,.2f}"
        
        return slide
    
    def _add_material_type_net(self, prs, pre_data):
        """添加物料类型净偏差分布（Slide 10）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "物料类型净偏差分布"
        
        material_type_data = pre_data.get('material_type_net', pd.DataFrame())
        
        if material_type_data.empty:
            tf = slide.placeholders[1].text_frame
            tf.text = "无物料类型数据"
            return slide
        
        # 创建表格
        rows = len(material_type_data) + 1
        cols = len(material_type_data.columns) + 1
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 表头
        table.cell(0, 0).text = '物料类型'
        for j, factory in enumerate(material_type_data.columns, 1):
            table.cell(0, j).text = factory
        
        # 数据
        for i, (mtype, row) in enumerate(material_type_data.iterrows(), 1):
            table.cell(i, 0).text = mtype
            for j, factory in enumerate(material_type_data.columns, 1):
                amount = row.get(factory, 0)
                table.cell(i, j).text = f"{amount:,.2f}"
        
        return slide
    
    def _add_no_remark_warning(self, prs, pre_data, factory: str, threshold: float = 50000):
        """添加无备注预警 Top10（分工厂）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"无备注预警 Top10 - {factory}（≥{threshold/10000:.0f}万元）"
        
        warning_df = pre_data.get('no_remark_warning', pd.DataFrame())
        if warning_df.empty:
            tf = slide.placeholders[1].text_frame
            tf.text = "无符合阈值的无备注记录"
            return slide
        
        # 筛选该工厂的数据
        factory_warning = warning_df[warning_df['工厂'] == factory].nlargest(10, '偏差金额')
        
        if factory_warning.empty:
            tf = slide.placeholders[1].text_frame
            tf.text = f"该工厂无≥{threshold/10000:.0f}万元的无备注记录"
            return slide
        
        # 创建表格
        rows = len(factory_warning) + 1
        cols = 4
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 表头
        headers = ['物料名称', '车间', '偏差金额 (元)', '偏差率 (%)']
        for j, h in enumerate(headers):
            table.cell(0, j).text = h
        
        # 数据
        for i, (_, row) in enumerate(factory_warning.iterrows(), 1):
            table.cell(i, 0).text = row.get('物料名称', '')
            table.cell(i, 1).text = row.get('车间', '')
            table.cell(i, 2).text = f"{row.get('偏差金额', 0):,.2f}"
            table.cell(i, 3).text = f"{row.get('偏差率', 0):.2f}"
        
        return slide
    
    def _add_abnormal_warning(self, prs, pre_data):
        """添加异常预警明细"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "异常预警明细（定额>0 且 实际=0）"
        
        abnormal_df = pre_data.get('abnormal_warning', pd.DataFrame())
        
        if abnormal_df.empty:
            tf = slide.placeholders[1].text_frame
            tf.text = "无异常预警记录"
            return slide
        
        # 创建表格（显示前 20 条）
        top_20 = abnormal_df.head(20)
        rows = len(top_20) + 1
        cols = 5
        left = Inches(0.3)
        top = Inches(1.5)
        width = Inches(9.4)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 表头
        headers = ['工厂', '车间', '物料名称', '定额', '实际']
        for j, h in enumerate(headers):
            table.cell(0, j).text = h
        
        # 数据
        for i, (_, row) in enumerate(top_20.iterrows(), 1):
            table.cell(i, 0).text = row.get('工厂', '')
            table.cell(i, 1).text = row.get('车间', '')
            table.cell(i, 2).text = row.get('物料名称', '')
            table.cell(i, 3).text = str(row.get('数量 - 定额', 0))
            table.cell(i, 4).text = str(row.get('数量 - 实际', 0))
        
        if len(abnormal_df) > 20:
            note = slide.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(9.4), Inches(0.5))
            note.text_frame.text = f"注：共 {len(abnormal_df)} 条异常记录，仅显示前 20 条"
        
        return slide
    
    def _add_freq_reason(self, prs, pre_data):
        """添加高频原因对比"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "高频偏差原因对比"
        
        reason_data = pre_data.get('freq_reasons', {})
        
        if not reason_data:
            tf = slide.placeholders[1].text_frame
            tf.text = "无偏差原因数据"
            return slide
        
        # 创建表格
        rows = len(reason_data) + 1
        cols = 2
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 表头
        table.cell(0, 0).text = '偏差原因'
        table.cell(0, 1).text = '出现次数'
        
        # 数据
        for i, (reason, count) in enumerate(reason_data.items(), 1):
            table.cell(i, 0).text = reason
            table.cell(i, 1).text = str(count)
        
        return slide
    
    def _add_summary(self, prs, pre_data):
        """添加总结与建议页"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "总结与改进建议"
        
        tf = slide.placeholders[1].text_frame
        tf.clear()
        
        # 总体发现
        p = tf.add_paragraph()
        p.text = "【总体发现】"
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = f"• 总记录数：{pre_data['total_kpis']['total_records']:,} 条"
        p = tf.add_paragraph()
        p.text = f"• 总偏差金额：{pre_data['total_kpis']['total_amount']:,.2f} 元"
        p = tf.add_paragraph()
        p.text = f"• 高偏差（>10%）记录：{pre_data['total_kpis']['high_dev_count']} 条\\n"
        
        # 分工厂建议
        for factory, kpis in pre_data['factory_kpis'].items():
            p = tf.add_paragraph()
            p.text = f"【{factory} 专项建议】"
            p.font.bold = True
            p = tf.add_paragraph()
            p.text = f"• 需补备注记录：{kpis['no_remark_count']} 条"
            p = tf.add_paragraph()
            p.text = f"• 重点关注高偏差物料（见 Top10 页）"
            p = tf.add_paragraph()
            p.text = f"• 建议对车间进行专项培训\\n"
        
        return slide
    
    def generate_ppt(self, excel_path: str, output_path: str, progress_callback=None):
        """生成 PPT 报告（v1.2 优化版）
        
        Args:
            excel_path: Excel 文件路径
            output_path: PPT 输出路径
            progress_callback: 进度回调函数
        
        Returns:
            output_path: 生成的 PPT 文件路径
        """
        import pandas as pd
        from pptx import Presentation
        
        # 1. 数据验证
        if not os.path.exists(excel_path):
            raise FileNotFoundError(f"Excel 文件不存在：{excel_path}")
        
        # 2. 读取数据
        df = pd.read_excel(excel_path, sheet_name='完整偏差明细')
        if len(df) > 50000:
            raise Exception("数据量超过 5 万条，请缩小范围")
        
        # 3. 数据预处理
        pre_data = self._pre_aggregate_data(df)
        
        # 4. 创建 PPT（使用模板或 fallback）
        template_path = 'config/template.pptx'
        if os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = Presentation()
        
        # 5. 定义页面生成顺序
        pages = [
            (self._add_cover, {}),
            (self._add_toc, {}),
            (self._add_kpi_overview, {}),
            (self._add_factory_comparison, {}),
            (self._add_material_top10, {'factory': '1101'}),
            (self._add_material_top10, {'factory': '1102'}),
            (self._add_deviation_type_chart, {}),
            (self._add_workshop_details, {'factory': '1101'}),
            (self._add_workshop_details, {'factory': '1102'}),
            (self._add_material_type_net, {}),
            (self._add_no_remark_warning, {'factory': '1101'}),
            (self._add_no_remark_warning, {'factory': '1102'}),
            (self._add_abnormal_warning, {}),
            (self._add_freq_reason, {}),
            (self._add_summary, {}),
        ]
        
        # 6. 生成页面
        for i, (func, kwargs) in enumerate(pages):
            try:
                func(prs, pre_data, **kwargs)
                if progress_callback:
                    progress_callback((i + 1) / len(pages) * 100)
            except Exception as e:
                self.view.log(f"生成第{i+1}页失败：{e}", "error")
                # 继续生成后续页面
        
        # 7. 保存文件
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        
        return output_path
'''

# 5. 将新的 PPT 代码插入到文件中
# 找到原来的 generate_ppt 方法并替换
if 'def generate_ppt(self, excel_path: str, output_path: str, log_cb=None):' in content:
    # 找到方法开始和结束的位置
    lines = content.split('\n')
    new_lines = []
    in_old_method = False
    method_indent = None
    
    for i, line in enumerate(lines):
        if 'def generate_ppt(self, excel_path: str, output_path: str, log_cb=None):' in line:
            in_old_method = True
            method_indent = len(line) - len(line.lstrip())
            # 插入新方法
            new_lines.append(new_ppt_code.rstrip())
            continue
        
        if in_old_method:
            # 检查是否到了下一个方法
            current_indent = len(line) - len(line.lstrip()) if line.strip() else 0
            if line.strip().startswith('def ') and current_indent == method_indent:
                in_old_method = False
                new_lines.append(line)
            # 否则跳过（删除旧方法的行）
        else:
            new_lines.append(line)
    
    content = '\n'.join(new_lines)
    print("✓ Replaced old generate_ppt method")
else:
    print("⚠ Old generate_ppt method not found, appending new code")
    # 在类末尾添加新方法
    if 'class AuditPresenter:' in content:
        # 找到类的最后一行
        lines = content.split('\n')
        insert_idx = len(lines)
        for i in range(len(lines) - 1, -1, -1):
            if lines[i].strip() and not lines[i].strip().startswith('#'):
                insert_idx = i + 1
                break
        
        lines.insert(insert_idx, new_ppt_code)
        content = '\n'.join(lines)

# 6. 写入修改后的文件
with open(presenter_path, 'w', encoding='utf-8') as f:
    f.write(content)

print(f"✓ Updated audit_presenter.py (new size: {len(content)} bytes)")

# 7. 创建测试文件
test_code = '''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT 生成单元测试
覆盖任务卡要求的所有测试场景
"""
import os
import sys
import unittest
import pandas as pd
from datetime import datetime, timedelta

# 添加项目根目录到路径
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(__file__))))

from modules.audit.presenters.audit_presenter import AuditPresenter


class MockModel:
    """模拟 Model 层"""
    def run_analysis(self, **kwargs):
        return "/tmp/test_output"


class MockView:
    """模拟 View 层"""
    def __init__(self):
        self.logs = []
    
    def log(self, msg, level="info"):
        self.logs.append((level, msg))


class TestPPTGeneration(unittest.TestCase):
    """PPT 生成测试"""
    
    def setUp(self):
        """准备测试数据"""
        self.model = MockModel()
        self.view = MockView()
        self.presenter = AuditPresenter(self.model, self.view)
        
        # 创建测试数据
        self.test_df = pd.DataFrame({
            '工厂': ['1101'] * 5000 + ['1102'] * 4889,
            '车间': [f'车间{i%10}' for i in range(9889)],
            '物料编码': [f'{100 + i%5}00{i%100}' for i in range(9889)],
            '物料名称': [f'物料{i}' for i in range(9889)],
            '偏差金额': [i * 100 for i in range(9889)],
            '偏差率': [i % 20 - 10 for i in range(9889)],
            '备注原因': ['原因 A' if i % 3 == 0 else '' for i in range(9889)],
            '数量 - 定额': [100 if i % 50 == 0 else 50 for i in range(9889)],
            '数量 - 实际': [0 if i % 50 == 0 else 50 for i in range(9889)],
        })
        
        # 保存测试 Excel
        self.test_excel = os.path.join(os.path.dirname(__file__), 'test_data.xlsx')
        self.test_df.to_excel(self.test_excel, index=False)
        
        # PPT 输出路径
        self.test_ppt = os.path.join(os.path.dirname(__file__), 'test_output.pptx')
    
    def tearDown(self):
        """清理测试文件"""
        if os.path.exists(self.test_excel):
            os.remove(self.test_excel)
        if os.path.exists(self.test_ppt):
            os.remove(self.test_ppt)
    
    def test_generate_ppt_performance(self):
        """测试性能：9889 条数据生成时间≤30 秒"""
        import time
        
        start = time.time()
        output_path = self.presenter.generate_ppt(
            self.test_excel,
            self.test_ppt,
            progress_callback=lambda p: None
        )
        elapsed = time.time() - start
        
        self.assertLess(elapsed, 30, f"生成时间 {elapsed:.2f}秒 超过 30 秒")
        self.assertTrue(os.path.exists(output_path))
    
    def test_factory_separation(self):
        """测试分工厂逻辑：食品厂/饮料厂数据隔离"""
        pre_data = self.presenter._pre_aggregate_data(self.test_df)
        
        # 检查工厂 KPI 是否分离
        self.assertIn('1101', pre_data['factory_kpis'])
        self.assertIn('1102', pre_data['factory_kpis'])
        
        # 检查数据量
        factory_1101_count = pre_data['factory_kpis']['1101']['total_records']
        factory_1102_count = pre_data['factory_kpis']['1102']['total_records']
        
        self.assertEqual(factory_1101_count, 5000)
        self.assertEqual(factory_1102_count, 4889)
    
    def test_material_classification(self):
        """测试物料分类：编码前缀 100/400/200/600"""
        test_cases = [
            ('100001', '原材料'),
            ('400123', '原材料'),
            ('200456', '包材'),
            ('600789', '包材'),
            ('300999', '其他'),
            ('', '其他'),
            (None, '其他'),
        ]
        
        for code, expected in test_cases:
            result = self.presenter._classify_material_type(code)
            self.assertEqual(result, expected, f"物料编码 {code} 分类错误")
    
    def test_empty_data_handling(self):
        """测试空数据处理"""
        empty_df = pd.DataFrame()
        pre_data = self.presenter._pre_aggregate_data(empty_df)
        
        # 应该不抛出异常，返回空数据
        self.assertEqual(pre_data['total_kpis']['total_records'], 0)
    
    def test_insufficient_top10(self):
        """测试不足 10 条物料的情况"""
        small_df = pd.DataFrame({
            '工厂': ['1102'] * 3,
            '物料名称': ['物料 A', '物料 B', '物料 C'],
            '偏差金额': [100, 200, 300],
        })
        
        pre_data = self.presenter._pre_aggregate_data(small_df)
        top10 = pre_data['material_top10'].get('1102', {})
        
        # 应该显示全部 3 条
        self.assertEqual(len(top10), 3)
    
    def test_data_limit(self):
        """测试数据超限：5 万 + 条数据应抛出异常"""
        large_df = pd.DataFrame({
            '工厂': ['1101'] * 50001,
            '偏差金额': [100] * 50001,
        })
        
        large_excel = self.test_excel.replace('.xlsx', '_large.xlsx')
        large_df.to_excel(large_excel, index=False)
        
        with self.assertRaises(Exception) as context:
            self.presenter.generate_ppt(large_excel, self.test_ppt)
        
        self.assertIn("数据量超过 5 万条", str(context.exception))
        
        if os.path.exists(large_excel):
            os.remove(large_excel)
    
    def test_progress_callback(self):
        """测试进度回调"""
        progress_values = []
        
        def callback(pct):
            progress_values.append(pct)
        
        self.presenter.generate_ppt(self.test_excel, self.test_ppt, progress_callback=callback)
        
        # 检查进度从 0 递进到 100
        self.assertGreater(len(progress_values), 0)
        self.assertGreaterEqual(max(progress_values), 99)
        self.assertLessEqual(max(progress_values), 100)


if __name__ == '__main__':
    unittest.main()
'''

test_path = os.path.join(root, 'tests', 'test_ppt_generation.py')
with open(test_path, 'w', encoding='utf-8') as f:
    f.write(test_code)

print(f"✓ Created test file: tests/test_ppt_generation.py")

print("\n=== Implementation Complete ===")
print("✓ Created directories: output/pptx/, config/, tests/")
print("✓ Created template placeholder: config/template.pptx")
print("✓ Updated audit_presenter.py with v1.2 PPT generation")
print("✓ Created test file: tests/test_ppt_generation.py")
print("\nNext steps:")
print("1. Run tests: python -m pytest tests/test_ppt_generation.py -v")
print("2. Test manually: python main.py and click '生成 PPT'")

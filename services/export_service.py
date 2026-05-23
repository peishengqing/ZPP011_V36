#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Export Service - 导出服务层

职责：
- PPT 报告生成
- Excel 导出
- 文件输出管理

依赖：
- DataService（数据读取）
- FileService（文件保存）
"""

import os
import pandas as pd
from typing import Optional, Callable
from datetime import datetime


class ExportService:
    """导出服务类"""
    
    def __init__(self, data_service: 'DataService', file_service: 'FileService'):
        self.data_service = data_service
        self.file_service = file_service
    
    def generate_ppt(self, df: pd.DataFrame, output_path: str, 
                     progress_callback: Optional[Callable] = None) -> str:
        """生成 PPT 报告
        
        Args:
            df: DataFrame
            output_path: PPT 输出路径
            progress_callback: 进度回调函数
        
        Returns:
            生成的 PPT 文件路径
        """
        from pptx import Presentation
        
        # 数据验证
        if len(df) > 50000:
            raise Exception("数据量超过 5 万条，请缩小范围")
        
        # 创建 PPT
        prs = Presentation()
        
        # 添加封面
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "ZPP011 生产偏差分析报告"
        subtitle = slide.placeholders[1]
        subtitle.text = f"数据量：{len(df)} 条\n生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # 保存文件
        self.file_service.ensure_directory(os.path.dirname(output_path))
        prs.save(output_path)
        
        if progress_callback:
            progress_callback(100)
        
        return output_path
    
    def export_to_excel(self, df: pd.DataFrame, output_path: str, 
                        sheet_name: str = 'Data') -> str:
        """导出到 Excel
        
        Args:
            df: DataFrame
            output_path: Excel 输出路径
            sheet_name: 工作表名称
        
        Returns:
            生成的 Excel 文件路径
        """
        self.file_service.ensure_directory(os.path.dirname(output_path))
        df.to_excel(output_path, index=False, sheet_name=sheet_name)
        return output_path
    
    def export_to_csv(self, df: pd.DataFrame, output_path: str) -> str:
        """导出到 CSV
        
        Args:
            df: DataFrame
            output_path: CSV 输出路径
        
        Returns:
            生成的 CSV 文件路径
        """
        self.file_service.ensure_directory(os.path.dirname(output_path))
        df.to_csv(output_path, index=False, encoding='utf-8-sig')
        return output_path

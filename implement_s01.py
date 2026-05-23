#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
S01 模块整合实施脚本
创建 services/ 层架构，迁移业务逻辑
"""
import os
import sys
sys.stdout.reconfigure(encoding='utf-8')

root = r'E:\zpp011_dev\模块化脚本'

print("=== S01 Module Integration - Implementation ===\n")

# 1. 创建 services/ 目录
services_dir = os.path.join(root, 'services')
os.makedirs(services_dir, exist_ok=True)
print(f"✓ Created: services/")

# 2. 创建 __init__.py（服务入口）
init_code = '''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Services Layer - S01 Module Integration

服务层入口，导出所有 Service 类供外部使用。
遵循单一职责原则，无循环依赖。
"""

from services.audit_service import AuditService
from services.export_service import ExportService
from services.file_service import FileService
from services.data_service import DataService
from services.filter_service import FilterService
from services.config_service import ConfigService
from services.storage_service import StorageService

__all__ = [
    'AuditService',
    'ExportService',
    'FileService',
    'DataService',
    'FilterService',
    'ConfigService',
    'StorageService',
]
'''

with open(os.path.join(services_dir, '__init__.py'), 'w', encoding='utf-8') as f:
    f.write(init_code)
print("✓ Created: services/__init__.py")

# 3. 创建各个 Service 文件
services_to_create = {
    'file_service.py': '''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
File Service - 文件服务层

职责：
- 文件读取、保存、备份
- 路径处理、目录创建
- 文件验证（存在性、权限）

依赖：
- utils.helpers（纯工具函数）
"""

import os
import shutil
from datetime import datetime
from typing import Optional, List
from utils.helpers import standardize_remark


class FileService:
    """文件服务类"""
    
    def __init__(self):
        pass
    
    def ensure_directory(self, dir_path: str) -> str:
        """确保目录存在，不存在则创建
        
        Args:
            dir_path: 目录路径
        
        Returns:
            目录路径（绝对路径）
        """
        os.makedirs(dir_path, exist_ok=True)
        return dir_path
    
    def get_backup_path(self, file_path: str) -> str:
        """生成备份文件路径
        
        Args:
            file_path: 原文件路径
        
        Returns:
            备份文件路径（带时间戳）
        """
        dir_name = os.path.dirname(file_path)
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        ext = os.path.splitext(file_path)[1]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        backup_name = f"{base_name}_备份_{timestamp}{ext}"
        return os.path.join(dir_name, backup_name)
    
    def backup_file(self, file_path: str) -> Optional[str]:
        """备份文件
        
        Args:
            file_path: 原文件路径
        
        Returns:
            备份文件路径，失败返回 None
        """
        if not os.path.exists(file_path):
            return None
        
        try:
            backup_path = self.get_backup_path(file_path)
            shutil.copy2(file_path, backup_path)
            return backup_path
        except Exception as e:
            print(f"Backup failed: {e}")
            return None
    
    def find_latest_file(self, pattern: str, directory: str) -> Optional[str]:
        """查找匹配模式的最新文件
        
        Args:
            pattern: 文件名模式（如 "ZPP011 偏差分析最终版_*.xlsx"）
            directory: 搜索目录
        
        Returns:
            最新文件路径，未找到返回 None
        """
        import glob
        
        if not os.path.isdir(directory):
            return None
        
        search_pattern = os.path.join(directory, pattern)
        files = glob.glob(search_pattern)
        
        if not files:
            return None
        
        return max(files, key=os.path.getmtime)
    
    def validate_file(self, file_path: str, extensions: List[str] = None) -> bool:
        """验证文件
        
        Args:
            file_path: 文件路径
            extensions: 允许的扩展名列表（如 ['.xlsx', '.xls']）
        
        Returns:
            验证是否通过
        """
        if not os.path.exists(file_path):
            return False
        
        if extensions:
            ext = os.path.splitext(file_path)[1].lower()
            if ext not in extensions:
                return False
        
        return True
''',

    'data_service.py': '''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Data Service - 数据服务层

职责：
- DataFrame 预处理、清洗
- 数据聚合、KPI 计算
- 偏差分析逻辑

依赖：
- FileService（文件读取）
- utils.helpers（备注标准化）
"""

import pandas as pd
import numpy as np
from typing import Dict, Any, Optional
from utils.helpers import standardize_remark


class DataService:
    """数据服务类"""
    
    def __init__(self):
        self._df_cache: Optional[pd.DataFrame] = None
    
    def load_excel(self, file_path: str, sheet_name: str = 'Data') -> pd.DataFrame:
        """加载 Excel 文件
        
        Args:
            file_path: Excel 文件路径
            sheet_name: 工作表名称
        
        Returns:
            DataFrame
        """
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        self._df_cache = df
        return df
    
    def clean_column_names(self, df: pd.DataFrame) -> pd.DataFrame:
        """清洗列名（移除空格、统一格式）
        
        Args:
            df: 原始 DataFrame
        
        Returns:
            清洗后的 DataFrame
        """
        df_clean = df.copy()
        df_clean.columns = [col.strip().replace(' ', '') for col in df_clean.columns]
        
        # 统一偏差金额列名
        if '偏差金额 (含税)' in df_clean.columns:
            df_clean.rename(columns={'偏差金额 (含税)': '偏差金额'}, inplace=True)
        
        return df_clean
    
    def calculate_kpis(self, df: pd.DataFrame, factory: str = None) -> Dict[str, Any]:
        """计算 KPI 指标
        
        Args:
            df: DataFrame
            factory: 工厂筛选（可选）
        
        Returns:
            KPI 字典
        """
        if factory:
            df = df[df['工厂'] == factory].copy()
        
        total_records = len(df)
        total_amount = df['偏差金额'].sum() if '偏差金额' in df.columns else 0
        avg_dev_rate = df['偏差率'].mean() if '偏差率' in df.columns else 0
        high_dev_count = len(df[abs(df['偏差率']) > 10]) if '偏差率' in df.columns else 0
        no_remark_count = len(df[
            (df['备注原因'].isna()) | 
            (df['备注原因'] == '')
        ]) if '备注原因' in df.columns else 0
        
        return {
            'total_records': total_records,
            'total_amount': total_amount,
            'avg_dev_rate': avg_dev_rate,
            'high_dev_count': high_dev_count,
            'no_remark_count': no_remark_count,
        }
    
    def aggregate_by_factory(self, df: pd.DataFrame) -> Dict[str, Dict]:
        """按工厂维度聚合
        
        Args:
            df: DataFrame
        
        Returns:
            工厂维度 KPI 字典
        """
        factory_kpis = {}
        
        for factory in df['工厂'].unique():
            factory_df = df[df['工厂'] == factory]
            factory_kpis[factory] = self.calculate_kpis(factory_df)
        
        return factory_kpis
    
    def get_material_top10(self, df: pd.DataFrame, factory: str = None) -> Dict:
        """获取物料偏差金额 Top10
        
        Args:
            df: DataFrame
            factory: 工厂筛选（可选）
        
        Returns:
            Top10 字典 {物料名称：偏差金额}
        """
        if factory:
            df = df[df['工厂'] == factory].copy()
        
        df['偏差金额_abs'] = df['偏差金额'].abs()
        material_sum = df.groupby('物料名称')['偏差金额_abs'].sum()
        
        return material_sum.nlargest(10).to_dict()
    
    def classify_material_type(self, material_code: str) -> str:
        """物料类型分类（S01 专用）
        
        Args:
            material_code: 物料编码
        
        Returns:
            物料类型：'原材料' / '包材' / '其他'
        """
        if not material_code or not isinstance(material_code, str):
            return '其他'
        
        prefix = material_code.strip()[:3]
        
        if prefix in ('100', '400'):
            return '原材料'
        elif prefix in ('200', '600'):
            return '包材'
        else:
            return '其他'
    
    def get_workshop_stats(self, df: pd.DataFrame, factory: str = None) -> Dict:
        """获取车间维度统计
        
        Args:
            df: DataFrame
            factory: 工厂筛选（可选）
        
        Returns:
            车间统计字典 {车间：偏差金额}
        """
        if factory:
            df = df[df['工厂'] == factory].copy()
        
        workshop_sum = df.groupby('车间')['偏差金额'].sum()
        return workshop_sum.to_dict()
''',

    'audit_service.py': '''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Audit Service - 审核服务层

职责：
- AI 审核逻辑
- 自动结案
- 备注更新、状态管理

依赖：
- DataService（数据预处理）
- FilterService（筛选待审核数据）
- StorageService（审计记录存储）
"""

import pandas as pd
from typing import List, Dict, Any, Optional, Callable
from datetime import datetime


class AuditService:
    """审核服务类"""
    
    def __init__(self, data_service: 'DataService', filter_service: 'FilterService'):
        self.data_service = data_service
        self.filter_service = filter_service
        self._cancel_flag: bool = False
    
    def request_cancel(self):
        """请求取消当前操作"""
        self._cancel_flag = True
    
    def reset_cancel_flag(self):
        """重置取消标志"""
        self._cancel_flag = False
    
    def auto_close_cases(self, df: pd.DataFrame, threshold: float = 10.0) -> pd.DataFrame:
        """自动结案：为符合条件的行自动填写备注
        
        Args:
            df: DataFrame
            threshold: 偏差率阈值（默认 10%）
        
        Returns:
            更新后的 DataFrame
        """
        self.reset_cancel_flag()
        
        df_copy = df.copy()
        
        # 筛选需要自动结案的行
        mask = (
            (df_copy['偏差率'].abs() <= threshold) &
            ((df_copy['备注原因'].isna()) | (df_copy['备注原因'] == ''))
        )
        
        # 自动填写备注
        df_copy.loc[mask, '备注原因'] = f'自动结案（偏差率≤{threshold}%）'
        df_copy.loc[mask, '备注来源'] = '自动结案'
        
        return df_copy
    
    def ai_audit(self, df: pd.DataFrame, progress_callback: Optional[Callable] = None) -> pd.DataFrame:
        """AI 审核（简化版，实际应调用 AI 模型）
        
        Args:
            df: DataFrame
            progress_callback: 进度回调函数
        
        Returns:
            更新后的 DataFrame（含 AI 建议）
        """
        self.reset_cancel_flag()
        
        df_copy = df.copy()
        total = len(df_copy)
        
        # 初始化 AI 建议列
        if 'AI 建议' not in df_copy.columns:
            df_copy['AI 建议'] = ''
        
        # 逐行处理（模拟 AI 审核）
        for idx, row in df_copy.iterrows():
            if self._cancel_flag:
                break
            
            # 简单规则审核（实际应替换为 AI 模型调用）
            ai_suggestion = self._generate_ai_suggestion(row)
            df_copy.at[idx, 'AI 建议'] = ai_suggestion
            
            if progress_callback:
                progress_callback((idx + 1) / total * 100)
        
        return df_copy
    
    def _generate_ai_suggestion(self, row: pd.Series) -> str:
        """生成 AI 建议（简化版）
        
        Args:
            row: 数据行
        
        Returns:
            AI 建议文本
        """
        dev_rate = abs(row.get('偏差率', 0))
        has_note = pd.notna(row.get('备注原因')) and str(row.get('备注原因')).strip() != ''
        
        if not has_note:
            if dev_rate > 10:
                return "⚠️ 高偏差且无备注，请补充原因"
            else:
                return "ℹ️ 建议补充备注"
        else:
            return "✅ 备注完整"
    
    def update_remark(self, df: pd.DataFrame, row_index: int, new_remark: str, source: str = '人工填写') -> pd.DataFrame:
        """更新备注
        
        Args:
            df: DataFrame
            row_index: 行索引
            new_remark: 新备注
            source: 备注来源
        
        Returns:
            更新后的 DataFrame
        """
        df_copy = df.copy()
        df_copy.at[row_index, '备注原因'] = new_remark
        df_copy.at[row_index, '备注来源'] = source
        
        return df_copy
    
    def get_pending_audit(self, df: pd.DataFrame) -> pd.DataFrame:
        """获取待审核数据
        
        Args:
            df: DataFrame
        
        Returns:
            待审核数据子集
        """
        # 筛选无备注或备注为空的行
        mask = (
            (df['备注原因'].isna()) | 
            (df['备注原因'] == '')
        )
        return df[mask].copy()
''',

    'filter_service.py': '''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Filter Service - 筛选服务层

职责：
- 动态筛选条件管理
- 筛选历史记忆
- 多条件组合筛选

依赖：
- 无（纯业务逻辑）
"""

import pandas as pd
from typing import Dict, List, Optional, Any
import json
import os


class FilterService:
    """筛选服务类"""
    
    def __init__(self, config_path: str = None):
        self.config_path = config_path or os.path.expanduser("~/.zpp011_audit/filter_history.json")
        self._filter_history: Dict[str, Any] = {}
        self._load_history()
    
    def _load_history(self):
        """加载筛选历史"""
        if os.path.exists(self.config_path):
            try:
                with open(self.config_path, 'r', encoding='utf-8') as f:
                    self._filter_history = json.load(f)
            except Exception:
                self._filter_history = {}
    
    def _save_history(self):
        """保存筛选历史"""
        try:
            os.makedirs(os.path.dirname(self.config_path), exist_ok=True)
            with open(self.config_path, 'w', encoding='utf-8') as f:
                json.dump(self._filter_history, f, ensure_ascii=False, indent=2)
        except Exception:
            pass
    
    def apply_filters(self, df: pd.DataFrame, filters: Dict[str, Any]) -> pd.DataFrame:
        """应用筛选条件
        
        Args:
            df: DataFrame
            filters: 筛选条件字典 {列名：筛选值}
        
        Returns:
            筛选后的 DataFrame
        """
        df_copy = df.copy()
        
        for col, value in filters.items():
            if col not in df_copy.columns:
                continue
            
            if value is None or value == '' or value == '全部':
                continue
            
            # 处理不同类型的筛选
            if isinstance(value, list):
                # 多选
                df_copy = df_copy[df_copy[col].isin(value)]
            elif isinstance(value, str):
                # 文本筛选（支持模糊匹配）
                df_copy = df_copy[df_copy[col].astype(str).str.contains(value, na=False)]
            elif isinstance(value, (int, float)):
                # 数值筛选
                df_copy = df_copy[df_copy[col] == value]
        
        return df_copy
    
    def get_available_filters(self, df: pd.DataFrame) -> Dict[str, List]:
        """获取可用的筛选选项
        
        Args:
            df: DataFrame
        
        Returns:
            筛选选项字典 {列名：[选项列表]}
        """
        filters = {}
        
        for col in df.columns:
            unique_values = df[col].dropna().unique()
            
            # 只保留前 100 个唯一值，避免选项过多
            if len(unique_values) > 100:
                continue
            
            filters[col] = sorted([str(v) for v in unique_values])
        
        return filters
    
    def save_filter_state(self, filter_name: str, filters: Dict[str, Any]):
        """保存筛选状态
        
        Args:
            filter_name: 筛选方案名称
            filters: 筛选条件
        """
        self._filter_history[filter_name] = filters
        self._save_history()
    
    def load_filter_state(self, filter_name: str) -> Optional[Dict[str, Any]]:
        """加载筛选状态
        
        Args:
            filter_name: 筛选方案名称
        
        Returns:
            筛选条件字典，不存在返回 None
        """
        return self._filter_history.get(filter_name)
    
    def reset_filters(self) -> Dict[str, Any]:
        """重置筛选条件
        
        Returns:
            空筛选条件字典
        """
        return {}
''',

    'export_service.py': '''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Export Service - 导出服务层

职责：
- PPT 报告生成
- Excel 导出
- 文件输出管理

依赖：
- DataService（数据读取）
- FileService（文件保存）
"""

import os
import pandas as pd
from typing import Optional, Callable
from datetime import datetime


class ExportService:
    """导出服务类"""
    
    def __init__(self, data_service: 'DataService', file_service: 'FileService'):
        self.data_service = data_service
        self.file_service = file_service
    
    def generate_ppt(self, df: pd.DataFrame, output_path: str, 
                     progress_callback: Optional[Callable] = None) -> str:
        """生成 PPT 报告
        
        Args:
            df: DataFrame
            output_path: PPT 输出路径
            progress_callback: 进度回调函数
        
        Returns:
            生成的 PPT 文件路径
        """
        from pptx import Presentation
        
        # 数据验证
        if len(df) > 50000:
            raise Exception("数据量超过 5 万条，请缩小范围")
        
        # 创建 PPT
        prs = Presentation()
        
        # 添加封面
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "ZPP011 生产偏差分析报告"
        subtitle = slide.placeholders[1]
        subtitle.text = f"数据量：{len(df)} 条\\n生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # 保存文件
        self.file_service.ensure_directory(os.path.dirname(output_path))
        prs.save(output_path)
        
        if progress_callback:
            progress_callback(100)
        
        return output_path
    
    def export_to_excel(self, df: pd.DataFrame, output_path: str, 
                        sheet_name: str = 'Data') -> str:
        """导出到 Excel
        
        Args:
            df: DataFrame
            output_path: Excel 输出路径
            sheet_name: 工作表名称
        
        Returns:
            生成的 Excel 文件路径
        """
        self.file_service.ensure_directory(os.path.dirname(output_path))
        df.to_excel(output_path, index=False, sheet_name=sheet_name)
        return output_path
    
    def export_to_csv(self, df: pd.DataFrame, output_path: str) -> str:
        """导出到 CSV
        
        Args:
            df: DataFrame
            output_path: CSV 输出路径
        
        Returns:
            生成的 CSV 文件路径
        """
        self.file_service.ensure_directory(os.path.dirname(output_path))
        df.to_csv(output_path, index=False, encoding='utf-8-sig')
        return output_path
''',

    'config_service.py': '''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Config Service - 配置服务层

职责：
- 参数配置管理
- 阈值设置
- 颜色方案、路径配置

依赖：
- FileService（配置文件读写）
"""

import os
import json
from typing import Dict, Any, Optional
from datetime import datetime


class ConfigService:
    """配置服务类"""
    
    DEFAULT_CONFIG = {
        'thresholds': {
            'high_deviation_rate': 10.0,  # 高偏差阈值 (%)
            'no_remark_amount': 50000,    # 无备注预警阈值 (元)
        },
        'colors': {
            'positive': '#d4edda',  # 正偏差
            'negative': '#f8d7da',  # 负偏差
            'alt_material': '#fff3cd',  # 替代料
        },
        'paths': {
            'output_dir': '~/Desktop',
            'backup_dir': '~/Desktop/zpp011_backups',
        },
        'limits': {
            'max_rows': 50000,  # 最大数据行数
            'ppt_timeout': 30,  # PPT 生成超时 (秒)
        },
    }
    
    def __init__(self, config_path: str = None):
        self.config_path = config_path or os.path.expanduser("~/.zpp011_audit/config.json")
        self._config: Dict[str, Any] = {}
        self._load_config()
    
    def _load_config(self):
        """加载配置文件"""
        if os.path.exists(self.config_path):
            try:
                with open(self.config_path, 'r', encoding='utf-8') as f:
                    self._config = json.load(f)
            except Exception:
                self._config = self.DEFAULT_CONFIG.copy()
        else:
            self._config = self.DEFAULT_CONFIG.copy()
            self._save_config()
    
    def _save_config(self):
        """保存配置文件"""
        try:
            os.makedirs(os.path.dirname(self.config_path), exist_ok=True)
            with open(self.config_path, 'w', encoding='utf-8') as f:
                json.dump(self._config, f, ensure_ascii=False, indent=2)
        except Exception:
            pass
    
    def get(self, key_path: str, default: Any = None) -> Any:
        """获取配置值
        
        Args:
            key_path: 配置键路径（如 "thresholds.high_deviation_rate"）
            default: 默认值
        
        Returns:
            配置值
        """
        keys = key_path.split('.')
        value = self._config
        
        for key in keys:
            if isinstance(value, dict) and key in value:
                value = value[key]
            else:
                return default
        
        return value
    
    def set(self, key_path: str, value: Any):
        """设置配置值
        
        Args:
            key_path: 配置键路径
            value: 配置值
        """
        keys = key_path.split('.')
        config = self._config
        
        for key in keys[:-1]:
            if key not in config:
                config[key] = {}
            config = config[key]
        
        config[keys[-1]] = value
        self._save_config()
    
    def get_threshold(self, name: str) -> float:
        """获取阈值
        
        Args:
            name: 阈值名称
        
        Returns:
            阈值
        """
        return self.get(f'thresholds.{name}', 0.0)
    
    def get_color(self, name: str) -> str:
        """获取颜色
        
        Args:
            name: 颜色名称
        
        Returns:
            颜色值（十六进制）
        """
        return self.get(f'colors.{name}', '#000000')
    
    def get_path(self, name: str) -> str:
        """获取路径
        
        Args:
            name: 路径名称
        
        Returns:
            路径（展开 ~）
        """
        path = self.get(f'paths.{name}', '~/')
        return os.path.expanduser(path)
''',

    'storage_service.py': '''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Storage Service - 存储服务层

职责：
- SQLite 审计记录存储
- 历史数据恢复
- 数据持久化

依赖：
- FileService（数据库文件管理）
"""

import sqlite3
import pandas as pd
from typing import Optional, List, Dict, Any
from datetime import datetime
import os


class StorageService:
    """存储服务类"""
    
    def __init__(self, db_path: str = None):
        self.db_path = db_path or os.path.expanduser("~/.zpp011_audit/audit_history.db")
        self._conn: Optional[sqlite3.Connection] = None
        self._init_db()
    
    def _get_connection(self) -> sqlite3.Connection:
        """获取数据库连接"""
        if self._conn is None:
            os.makedirs(os.path.dirname(self.db_path), exist_ok=True)
            self._conn = sqlite3.connect(self.db_path)
        return self._conn
    
    def _init_db(self):
        """初始化数据库表"""
        conn = self._get_connection()
        cursor = conn.cursor()
        
        # 创建审计记录表
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS audit_records (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                order_date TEXT,
                order_no TEXT,
                material_code TEXT,
                material_name TEXT,
                factory TEXT,
                workshop TEXT,
                deviation_rate REAL,
                deviation_amount REAL,
                remark TEXT,
                remark_source TEXT,
                audit_result TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        
        # 创建索引
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_order_no ON audit_records(order_no)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_material ON audit_records(material_code)")
        
        conn.commit()
    
    def save_audit_records(self, df: pd.DataFrame) -> int:
        """保存审计记录
        
        Args:
            df: DataFrame（含审计结果）
        
        Returns:
            保存的记录数
        """
        conn = self._get_connection()
        
        # 映射列名
        column_mapping = {
            '订单日期': 'order_date',
            '流程订单': 'order_no',
            '组件物料号': 'material_code',
            '组件物料描述': 'material_name',
            '工厂': 'factory',
            '车间': 'workshop',
            '偏差率': 'deviation_rate',
            '偏差金额': 'deviation_amount',
            '备注原因': 'remark',
            '备注来源': 'remark_source',
            'audit_result': 'audit_result',
        }
        
        # 选择需要的列
        columns_to_save = [col for col in column_mapping.keys() if col in df.columns]
        df_save = df[columns_to_save].copy()
        df_save.columns = [column_mapping[col] for col in columns_to_save]
        
        # 写入数据库
        df_save.to_sql('audit_records', conn, if_exists='append', index=False)
        conn.commit()
        
        return len(df_save)
    
    def load_audit_records(self, days: int = 30) -> pd.DataFrame:
        """加载历史审计记录
        
        Args:
            days: 加载最近 N 天的记录
        
        Returns:
            DataFrame
        """
        conn = self._get_connection()
        
        query = """
            SELECT * FROM audit_records
            WHERE created_at >= datetime('now', '-' || ? || ' days')
            ORDER BY created_at DESC
        """
        
        df = pd.read_sql_query(query, conn, params=(days,))
        return df
    
    def restore_audit_from_db(self, df: pd.DataFrame, log_cb=None) -> int:
        """从数据库恢复审计记录到 DataFrame
        
        Args:
            df: 当前数据 DataFrame
            log_cb: 日志回调函数
        
        Returns:
            恢复的记录数
        """
        conn = self._get_connection()
        cursor = conn.cursor()
        
        # 查询历史审计记录
        cursor.execute("""
            SELECT order_no, material_code, remark, remark_source, audit_result
            FROM audit_records
            ORDER BY updated_at DESC
        """)
        
        history_records = cursor.fetchall()
        restored_count = 0
        
        # 匹配并恢复
        for idx, row in df.iterrows():
            order_no = row.get('流程订单', '')
            material_code = row.get('组件物料号', '')
            
            # 查找匹配的历史记录
            for hist in history_records:
                if hist[0] == order_no and hist[1] == material_code:
                    df.at[idx, '备注原因'] = hist[2]
                    df.at[idx, '备注来源'] = hist[3]
                    df.at[idx, 'audit_result'] = hist[4]
                    restored_count += 1
                    break
        
        conn.commit()
        
        if log_cb:
            log_cb(f"📌 恢复 {restored_count} 条历史审计记录", "info")
        
        return restored_count
    
    def clear_old_records(self, days: int = 90) -> int:
        """清理旧记录
        
        Args:
            days: 清理 N 天前的记录
        
        Returns:
            清理的记录数
        """
        conn = self._get_connection()
        cursor = conn.cursor()
        
        cursor.execute("""
            DELETE FROM audit_records
            WHERE created_at < datetime('now', '-' || ? || ' days')
        """, (days,))
        
        deleted_count = cursor.rowcount
        conn.commit()
        
        return deleted_count
''',
}

# 创建所有 Service 文件
for filename, content in services_to_create.items():
    filepath = os.path.join(services_dir, filename)
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"✓ Created: services/{filename}")

print("\n=== S01 Implementation Complete ===")
print("✓ services/ directory created")
print("✓ 7 Service classes implemented:")
print("  - AuditService (审核业务)")
print("  - ExportService (导出业务)")
print("  - FileService (文件服务)")
print("  - DataService (数据服务)")
print("  - FilterService (筛选服务)")
print("  - ConfigService (配置服务)")
print("  - StorageService (存储服务)")
print("\nNext: Update imports in events.py and other modules")

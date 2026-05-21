# modules/audit/presenters/audit_presenter.py
import os
import pandas as pd
import threading
from typing import Any, Dict, Optional, List


class AuditPresenter:
    """MVP Presenter 层：业务逻辑、协调 Model 和 View"""

    def __init__(self, model: 'AuditModel', view: 'AuditViewBridge'):
        self.model = model
        self.view = view  # view 是 events.py 中提供的一组回调接口（非完整 View 对象）
        self.running = False
        self.cancel_req = False
        # 其他状态可迁移至此
        self._ai_cancel_flag = None
        self.is_auditing = False

    # ---------- 业务逻辑方法（从 events.py 迁移）----------
    def start_analysis(self, input_file: str, alt_pairs: list,
                     start_date: str = None, end_date: str = None,
                     material_search: str = None,
                     progress_callback=None, cancel_check=None):
        """开始分析（同步，在 worker 线程中调用）
        
        Args:
            input_file: 输入文件路径
            alt_pairs: 替代配对列表
            start_date: 开始日期
            end_date: 结束日期
            material_search: 物料搜索关键词
            progress_callback: 进度回调函数
            cancel_check: 取消检查函数
        
        Returns:
            output_path: 分析输出目录
        
        Raises:
            Exception: 分析失败
        """
        self.view.log("开始分析...", "info")
        self.running = True
        self.cancel_req = False
        
        import tempfile
        import traceback
        import os
        
        _log = os.path.join(os.environ.get('TEMP', '.'), 'zpp011_debug.log')
        
        try:
            with open(_log, 'a', encoding='utf-8') as _f:
                _f.write(f"\n=== {__import__('datetime').datetime.now()} 开始分析 ===\n")
                _f.write(f"input_file={input_file}\n")
            
            # 创建临时目录
            temp_dir = tempfile.mkdtemp(prefix="zpp011_analysis_")
            self.view.log(f"临时目录: {temp_dir}", "debug")
            
            # 调用 Model 层进行分析
            output_path = self.model.run_analysis(
                input_file=input_file,
                output_dir=temp_dir,
                alt_pairs=alt_pairs,
                progress_callback=progress_callback,
                cancel_check=cancel_check,
                start_date=start_date,
                end_date=end_date,
                material_search=material_search,
            )
            
            with open(_log, 'a', encoding='utf-8') as _f:
                _f.write(f"{__import__('datetime').datetime.now()} 分析完成\n")
            
            self.view.log("分析完成", "success")
            return output_path
        
        except KeyboardInterrupt:
            with open(_log, 'a', encoding='utf-8') as _f:
                _f.write(f"{__import__('datetime').datetime.now()} KeyboardInterrupt\n")
            
            self.view.log("用户取消", "warn")
            self.cancel_req = True
            raise
        
        except Exception as e:
            with open(_log, 'a', encoding='utf-8') as _f:
                _f.write(f"{__import__('datetime').datetime.now()} Exception: {e}\n")
                _f.write(traceback.format_exc() + "\n")
            
            self.view.log(f"分析失败: {e}", "error")
            raise
        
        finally:
            self.running = False

    def request_cancel_analysis(self):
        """请求取消分析"""
        self.cancel_req = True
        self.view.log("取消请求已发送", "warn")

    def filter_audit_data(self, filters: Dict[str, Any]) -> 'pd.DataFrame':
        """筛选审核数据（纯业务逻辑，不操作UI）

        Args:
            filters: dict with keys: stat, color, search, order_date, columns,
                     factory, admin, name, status, dev_rate, is_alt, remark
        Returns:
            Filtered DataFrame
        """
        import pandas as pd
        audit_data = self.view.get_audit_data()
        if audit_data is None or len(audit_data) == 0:
            return audit_data

        df = audit_data.copy()

        # 1) Stat card filter
        stat = filters.get('stat')
        if stat == 'big_dev':
            rc = next((c for c in ['偏差率(%)', '偏差率'] if c in df.columns), None)
            if rc:
                df = df[pd.to_numeric(df[rc], errors='coerce').abs() > 10]
        elif stat == 'no_note':
            rc = next((c for c in ['备注原因', '备注'] if c in df.columns), '备注原因')
            df = df[df[rc].isna() | (df[rc].astype(str).str.strip() == '')]
        elif stat == 'approved':
            rc = next((c for c in ['备注原因', '备注'] if c in df.columns), '备注原因')
            df = df[df[rc].notna() & (df[rc].astype(str).str.strip() != '')]

        # 2) Color filter
        color = filters.get('color')
        if color and color != '全部' and '_priority_label' in df.columns:
            cmap = {'红': '红', '橙': '橙', '黄': '黄', '绿': '绿'}
            df = df[df['_priority_label'] == cmap.get(color, color)]

        # 3) Full-text search
        search = filters.get('search', '').strip()
        if search:
            mask = pd.Series(False, index=df.index)
            for col in df.columns:
                mask |= df[col].astype(str).str.contains(search, case=False, na=False)
            df = df[mask]

        # 4) Date range
        od = filters.get('order_date')
        if od and '订单日期' in df.columns:
            s, e = od
            if s:
                df = df[df['订单日期'].astype(str).str[:10] >= s]
            if e:
                df = df[df['订单日期'].astype(str).str[:10] <= e]

        # 5) Field filters
        col_map = filters.get('columns', {})
        for key, col_name in col_map.items():
            val = filters.get(key)
            if not val or val == '全部' or col_name is None:
                continue
            if key == 'is_alt':
                alt_pairs = self.view.get_alt_pairs()
                descs = set()
                for a, b in alt_pairs:
                    for item in (a, b):
                        if isinstance(item, (list, tuple)) and len(item) >= 3:
                            d = str(item[2]).strip()
                        elif isinstance(item, (list, tuple)):
                            d = str(item[1]).strip() if len(item) > 1 and item[1] else ''
                        else:
                            d = str(item).strip() if item else ''
                        if d and d != 'None':
                            descs.add(d)
                dc = next((c for c in ['组件物料描述', '物料名称'] if c in df.columns), None)
                if dc:
                    if val == '是':
                        df = df[df[dc].astype(str).str.strip().isin(descs)]
                    elif val == '否':
                        df = df[~df[dc].astype(str).str.strip().isin(descs)]
            elif key == 'dev_rate':
                rv = pd.to_numeric(df[col_name], errors='coerce')
                ops = {'>10%': rv > 10, '>20%': rv > 20, '>30%': rv > 30,
                       '绝对值>10%': rv.abs() > 10, '<-10%': rv < -10, '<-20%': rv < -20}
                if val in ops:
                    df = df[ops[val]]
            elif key == 'remark':
                # 备注筛选：动态定位列名，标准化空值后匹配
                remark_col = next(
                    (c for c in ['备注原因', '备注', 'remark'] if c in df.columns), None
                )
                if remark_col is not None:
                    temp = df[remark_col].fillna('').astype(str).str.strip().replace('nan', '')
                    if val == '为空':
                        df = df[temp == '']
                    elif val == '不为空':
                        df = df[temp != '']
                    else:
                        # 精确文本匹配
                        df = df[temp == val]
            else:
                if col_name in df.columns:
                    df = df[df[col_name].astype(str).str.contains(val, case=False, na=False)]


        # 6) AI审核结果筛选（_on_filter_changed 迁移）
        ai_result = filters.get('ai_result')
        if ai_result and ai_result != '全部':
            remark_src_col = next((c for c in ['审核来源', 'audit_source'] if c in df.columns), None)
            if remark_src_col:
                if ai_result == '合格':
                    df = df[df[remark_src_col] == '审核合格']
                elif ai_result == '需改进':
                    df = df[df[remark_src_col] == '审核待改进']
                elif ai_result == 'AI建议':
                    df = df[df[remark_src_col].str.startswith('AI建议', na=False)]
                elif ai_result == '未处理':
                    ai_sources = {'审核合格', '审核待改进', 'AI建议', 'AI建议（小偏差）'}
                    df = df[~(df[remark_src_col].isin(ai_sources) | df[remark_src_col].isna())]


        return df

    def run_ai_audit(self):
        """AI审核：启动异步审核流程（从 events.py 迁移）
        
        调用前由 View 层检查前置条件（is_auditing, audit_data 是否为空）
        返回 (audit_indices, df_to_audit) 供 View 层启动 task_manager
        """
        import pandas as pd
        
        audit_data = self.view.get_audit_data()
        alt_pairs = self.view.get_alt_pairs()
        
        # 确保必要列存在
        for col in ['audit_result', 'AI建议', '备注来源']:
            if col not in audit_data.columns:
                audit_data[col] = ''
        
        # 构建替代料名称集合
        alt_all_descs = set()
        for a, b in alt_pairs:
            for item in (a, b):
                if item:
                    if isinstance(item, (list, tuple)) and len(item) >= 3:
                        desc = str(item[2]).strip()
                    elif isinstance(item, (list, tuple)) and len(item) == 2:
                        desc = str(item[1]).strip() if item[1] else str(item[0]).strip()
                    else:
                        desc = str(item).strip()
                    if desc and desc != 'None':
                        alt_all_descs.add(desc)
        
        # 自动填充条件
        is_auto_fill = (
            audit_data['组件物料描述'].astype(str).str.contains('透明胶', na=False) |
            audit_data['组件物料号'].astype(str).str.startswith('600')
        )
        is_alt = audit_data['组件物料描述'].astype(str).str.strip().isin(alt_all_descs)
        exclude_sources = ['人工填写', '自动填充', '替代料', 'AI审核合格', 'AI审核待改进', 'AI生成']
        is_already_processed = audit_data['备注来源'].isin(exclude_sources)
        
        to_audit_mask = (
            (audit_data['audit_result'].isna() | (audit_data['audit_result'] == '')) &
            (~is_auto_fill) &
            (~is_alt) &
            (~is_already_processed)
        )
        
        audit_indices = audit_data[to_audit_mask].index.tolist()
        return audit_indices
    
    def ai_audit_worker(self, audit_indices, df_to_audit, ai_client, progress_callback=None):
        """AI审核工作线程函数（同步）
        
        Args:
            audit_indices: 需要审核的行索引列表
            df_to_audit: 审核数据的 DataFrame
            ai_client: AI客户端实例
            progress_callback: 进度回调
        
        Returns:
            list of dict: 审核结果列表
        """
        import pandas as pd
        
        total = len(audit_indices)
        popup_rows = []
        remark_col = next((c for c in ['备注原因', '备注'] if c in df_to_audit.columns), None)
        name_col = next((c for c in ['组件物料描述', '物料名称'] if c in df_to_audit.columns), None)
        rate_col = next((c for c in ['偏差率(%)', '偏差率'] if c in df_to_audit.columns), None)
        
        for seq, idx in enumerate(audit_indices):
            if self._ai_cancel_flag and self._ai_cancel_flag.is_set():
                raise InterruptedError("用户取消了AI审核")
            
            row = df_to_audit.loc[idx]
            remark = row[remark_col] if remark_col else ''
            dev_rate_raw = row[rate_col] if rate_col else 0
            try:
                dev_rate = float(dev_rate_raw or 0)
            except:
                dev_rate = 0.0
            material_desc = str(row[name_col]) if name_col else ''
            
            try:
                result = ai_client.audit(remark, dev_rate)
                audit_result = result.get('result', '需补备注')
                ai_suggestion = result.get('suggestion', '')
            except Exception as ex:
                audit_result = '审核失败'
                ai_suggestion = str(ex)
            
            popup_rows.append({
                'idx': idx,
                '物料': material_desc,
                '偏差率': f"{dev_rate:.1f}%",
                '原备注': str(remark) if not pd.isna(remark) else '',
                'AI建议': ai_suggestion,
                '审核结果': audit_result,
                '_audit_result': audit_result,
                '_ai_suggestion': ai_suggestion,
            })
            
            if progress_callback:
                progress_callback(int((seq + 1) / total * 100))
        
        return popup_rows
    
    def cancel_ai_audit(self):
        """取消当前 AI 审核"""
        if self._ai_cancel_flag is not None:
            self._ai_cancel_flag.set()

    def save_audit_back(self, auditor: str = None):
        """保存审核结果到 Excel 和数据库"""
        self.view.log("保存审核结果...", "info")
        
        # 从 View 获取当前审核数据
        save_df = self.view.get_audit_data().copy()
        
        # 确保必要列存在
        if '审核人' not in save_df.columns and '审核人员' in save_df.columns:
            save_df['审核人'] = save_df['审核人员']
        if '审核意见' not in save_df.columns:
            save_df['审核意见'] = save_df.get('审核意见', '')
        
        try:
            # 调用 Model 层保存到数据库
            self.model.save_audit_to_db(save_df, auditor=auditor, log_cb=self.view.log)
            self.view.log("审核结果已保存到数据库", "success")
            
            # 同时写回 Excel（如果需要）
            output_path = self.view.get_output_path()
            if output_path:
                self.model.write_audit_back_to_excel(save_df, output_path)
                self.view.log("审核结果已写回 Excel", "success")
        
        except Exception as e:
            self.view.log(f"保存失败: {e}", "error")
            raise
    def auto_close(self, df_to_audit, rule_engine_snapshot, progress_callback, cancel_flag):
        """自动结案核心逻辑（委托 AutoCloser.process）"""
        from core.auto_closer import AutoCloser
        return AutoCloser.process(
            df_to_audit, rule_engine_snapshot, progress_callback, cancel_flag
        )


    # ========== PPT 生成优化 v1.2 实现 ==========

    def _classify_material_type(self, material_code):
        """物料类型分类（Slide 10 专用）
        
        编码前缀 100/400 -> 原材料
        编码前缀 200/600 -> 包材
        其他 -> 其他
        """
        if not material_code or not isinstance(material_code, str):
            return '其他'
        prefix = str(material_code).strip()[:3]
        if prefix in ('100', '400'):
            return '原材料'
        elif prefix in ('200', '600'):
            return '包材'
        return '其他'

    def _pre_aggregate_data(self, df):
        """一次性计算所有聚合结果（避免重复查询 DataFrame）"""
        import pandas as pd
        import numpy as np

        pre = {}

        # ----- 标准化列名 -----
        COL_MAP = {}
        for target, candidates in [
            ('工厂', ['工厂', 'factory']),
            ('车间', ['车间', 'workshop']),
            ('偏差金额', ['偏差金额', '金额', 'dev_amount']),
            ('偏差率', ['偏差率(%)', '偏差率', 'dev_rate']),
            ('物料编码', ['物料编码', 'material_code']),
            ('物料名称', ['组件物料描述', '物料名称', 'material_name', '物料描述']),
            ('备注原因', ['备注原因', '备注', 'remark']),
            ('偏差类型', ['偏差类型', 'dev_type']),
            ('定额', ['数量 - 定额', '定额数量', 'quota']),
            ('实际', ['数量 - 实际', '实际数量', 'actual']),
        ]:
            for c in candidates:
                if c in df.columns:
                    COL_MAP[target] = c
                    break

        def get_col(name):
            return COL_MAP.get(name)

        def safe(df, col):
            if col and col in df.columns:
                return df[col]
            return pd.Series(dtype=float)

        df2 = df.copy()
        for col in ['工厂', '车间']:
            if get_col(col) and get_col(col) in df2.columns:
                df2[col] = df2[get_col(col)].fillna('未知')
            else:
                df2[col] = '未知'
        for col in ['偏差金额', '偏差率', '定额', '实际']:
            if get_col(col) and get_col(col) in df2.columns:
                df2[col] = pd.to_numeric(df2[get_col(col)], errors='coerce').fillna(0)
            else:
                df2[col] = 0.0

        # ----- 动态检测工厂 -----
        factory_col = get_col('工厂')
        if factory_col:
            factories = sorted(df2[factory_col].dropna().unique().tolist())
        else:
            factories = []
        pre['factories'] = factories
        pre['factory_col'] = factory_col

        # ----- 工厂维度 KPI -----
        factory_kpis = {}
        for fac in factories:
            mask = df2[df2[factory_col] == fac] if factory_col else df2
            sub = df2 if not factory_col else df2[df2[factory_col] == fac]
            dev_amt_col = get_col('偏差金额')
            dev_rate_col = get_col('偏差率')
            remark_col = get_col('备注原因')
            factory_kpis[fac] = {
                'total_records': len(sub),
                'total_amount': float(sub[dev_amt_col].sum()) if dev_amt_col else 0.0,
                'avg_dev_rate': float(sub[dev_rate_col].mean()) if dev_rate_col else 0.0,
                'high_dev_count': int((abs(sub[dev_rate_col]) > 10).sum()) if dev_rate_col else 0,
                'no_remark_count': int(sub[remark_col].isna().sum()) if remark_col else 0,
            }
        pre['factory_kpis'] = factory_kpis

        # ----- 总体 KPI -----
        dev_amt_col = get_col('偏差金额')
        dev_rate_col = get_col('偏差率')
        pre['total_kpis'] = {
            'total_records': len(df2),
            'total_amount': float(df2[dev_amt_col].sum()) if dev_amt_col else 0.0,
            'avg_dev_rate': float(df2[dev_rate_col].mean()) if dev_rate_col else 0.0,
            'high_dev_count': int((abs(df2[dev_rate_col]) > 10).sum()) if dev_rate_col else 0,
            'factories': len(factories),
        }

        # ----- 物料 Top10（分工厂） -----
        mat_name_col = get_col('物料名称')
        mat_top10 = {}
        for fac in factories:
            sub = df2[df2[factory_col] == fac] if factory_col else df2
            if mat_name_col:
                sub2 = sub.copy()
                sub2['_abs_amt'] = abs(sub2[dev_amt_col]) if dev_amt_col else 0
                top = sub2.groupby(mat_name_col)['_abs_amt'].sum().nlargest(10)
                mat_top10[fac] = top.to_dict()
            else:
                mat_top10[fac] = {}
        pre['material_top10'] = mat_top10

        # ----- 车间统计（分工厂） -----
        workshop_col = get_col('车间')
        workshop_stats = {}
        for fac in factories:
            sub = df2[df2[factory_col] == fac] if factory_col else df2
            if workshop_col:
                top5 = sub.groupby(workshop_col)[dev_amt_col].sum().nlargest(5)
                workshop_stats[fac] = top5.to_dict()
            else:
                workshop_stats[fac] = {}
        pre['workshop_stats'] = workshop_stats

        # ----- 偏差类型分布（分工厂，饼图用） -----
        dev_type_col = get_col('偏差类型')
        dev_type_dist = {}
        for fac in factories:
            sub = df2[df2[factory_col] == fac] if factory_col else df2
            if dev_type_col:
                dist = sub.groupby(dev_type_col)[dev_amt_col].sum()
                dev_type_dist[fac] = dist.to_dict()
            else:
                dev_type_dist[fac] = {}
        pre['dev_type_dist'] = dev_type_dist

        # ----- 物料类型净偏差（Slide 10，分工厂柱状图） -----
        mat_code_col = get_col('物料编码')
        if mat_code_col:
            df2['_mat_type'] = df2[mat_code_col].apply(self._classify_material_type)
            mat_type_col = '_mat_type'
        else:
            df2['_mat_type'] = '其他'
            mat_type_col = '_mat_type'

        mat_type_net = {}
        for fac in factories:
            sub = df2[df2[factory_col] == fac] if factory_col else df2
            net = sub.groupby(mat_type_col)[dev_amt_col].sum()
            mat_type_net[fac] = net.to_dict()
        pre['material_type_net'] = mat_type_net

        # ----- 无备注预警（分工厂，≥5万） -----
        remark_col = get_col('备注原因')
        no_remark = {}
        for fac in factories:
            sub = df2[df2[factory_col] == fac] if factory_col else df2
            if remark_col:
                mask = sub[remark_col].isna() | (sub[remark_col].astype(str).str.strip() == '')
                sub_warn = sub[mask & (abs(sub[dev_amt_col]) >= 50000)] if dev_amt_col else sub[mask]
                no_remark[fac] = sub_warn.nlargest(10, dev_amt_col) if dev_amt_col else sub_warn.head(10)
            else:
                no_remark[fac] = sub.nlargest(10, dev_amt_col) if dev_amt_col else sub.head(10)
        pre['no_remark_warning'] = no_remark

        # ----- 异常预警（定额>0 且 实际=0） -----
        quota_col = get_col('定额')
        actual_col = get_col('实际')
        if quota_col and actual_col:
            abnormal = df2[(df2[quota_col] > 0) & (df2[actual_col] == 0)]
            pre['abnormal_warning'] = abnormal
        else:
            pre['abnormal_warning'] = pd.DataFrame()

        # ----- 高频原因（分工厂 Top5） -----
        reason_col = get_col('备注原因')
        freq_reasons = {}
        for fac in factories:
            sub = df2[df2[factory_col] == fac] if factory_col else df2
            if reason_col:
                reasons = sub[reason_col].dropna()
                reasons = reasons[reasons.astype(str).str.strip() != '']
                freq = reasons.value_counts().head(5)
                freq_reasons[fac] = freq.to_dict()
            else:
                freq_reasons[fac] = {}
        pre['freq_reasons'] = freq_reasons

        return pre

    def _set_slide_title(self, slide, title_text, font_size=24):
        """设置幻灯片标题"""
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = title_text
        else:
            txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.7))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(font_size)
            p.font.bold = True

    def _add_page_number(self, slide, current, total):
        """在页面右下角添加页码"""
        txBox = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"第 {current} 页 / 共 {total} 页"
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.RIGHT

    def _add_cover(self, prs, pre_data, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        self._set_slide_title(slide, "ZPP011 生产偏差分析报告", 36)

        # 副标题/信息区
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
        tf = txBox.text_frame
        kpis = pre_data['total_kpis']
        lines = [
            f"总记录数：{kpis['total_records']:,} 条",
            f"工厂数量：{kpis['factories']} 个",
            f"总偏差金额：{kpis['total_amount']:,.2f} 元",
            f"平均偏差率：{kpis['avg_dev_rate']:.2f}%",
            f"高偏差记录：{kpis['high_dev_count']} 条",
            "",
            f"生成时间：{pd.Timestamp.now().strftime('%Y-%m-%d %H:%M')}",
        ]
        for i, line in enumerate(lines):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(16) if i < len(lines) - 1 else Pt(12)

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_toc(self, prs, pre_data, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题+内容
        self._set_slide_title(slide, "目  录")
        tf = slide.placeholders[1].text_frame
        tf.clear()

        factories = pre_data['factories']
        fac_names = ', '.join(str(f) for f in factories) if factories else ''

        toc = [
            "1. 核心指标总览",
            "2. 工厂维度对比",
        ]
        # 物料 Top10
        for fac in factories:
            toc.append(f"3. 物料偏差 Top10 — {fac}")
        # 偏差类型
        toc.append(f"{2 + len(factories) + 1}. 偏差类型分布")
        # 车间
        for fac in factories:
            toc.append(f"{3 + len(factories) + 1 + factories.index(fac)}. 车间详情 — {fac}")
        base = 3 + len(factories) * 2 + 1
        toc += [
            f"{base}. 物料类型净偏差分布",
        ]
        # 无备注预警
        for fac in factories:
            toc.append(f"{base + 1 + factories.index(fac)}. 无备注预警 — {fac}")
        base2 = base + 1 + len(factories)
        toc += [
            f"{base2 + 1}. 异常预警明细",
            f"{base2 + 2}. 高频偏差原因",
            f"{base2 + 3}. 总结与改进建议",
        ]

        for i, item in enumerate(toc):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(14)

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_kpi_overview(self, prs, pre_data, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        self._set_slide_title(slide, "核心指标总览")

        tf = slide.placeholders[1].text_frame
        tf.clear()

        kpis = pre_data['total_kpis']
        p = tf.paragraphs[0]
        p.text = "【总体指标】"
        p.font.bold = True
        p.font.size = Pt(16)

        summary_items = [
            f"  总记录数：{kpis['total_records']:,} 条",
            f"  总偏差金额：{kpis['total_amount']:,.2f} 元",
            f"  平均偏差率：{kpis['avg_dev_rate']:.2f}%",
            f"  高偏差（>10%）记录：{kpis['high_dev_count']} 条",
            "",
        ]
        for item in summary_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)

        for fac, fac_kpis in pre_data['factory_kpis'].items():
            p = tf.add_paragraph()
            p.text = f"【{fac}】"
            p.font.bold = True
            p.font.size = Pt(16)
            for item in [
                f"  记录数：{fac_kpis['total_records']:,}",
                f"  偏差金额：{fac_kpis['total_amount']:,.2f} 元",
                f"  高偏差记录：{fac_kpis['high_dev_count']}",
                f"  无备注记录：{fac_kpis['no_remark_count']}",
                "",
            ]:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(14)

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_factory_comparison(self, prs, pre_data, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_title(slide, "工厂维度对比")
        kpis = pre_data['factory_kpis']
        factories = pre_data['factories']

        rows = len(factories) + 1
        cols = 6
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.3), Inches(1.5),
            Inches(9.4), Inches(0.4 * rows)
        ).table

        headers = ['工厂', '记录数', '偏差金额(元)', '平均偏差率(%)', '高偏差数', '无备注数']
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            cell.text_frame.paragraphs[0].font.bold = True

        for i, fac in enumerate(factories, 1):
            k = kpis.get(fac, {})
            table.cell(i, 0).text = str(fac)
            table.cell(i, 1).text = f"{k.get('total_records', 0):,}"
            table.cell(i, 2).text = f"{k.get('total_amount', 0):,.2f}"
            table.cell(i, 3).text = f"{k.get('avg_dev_rate', 0):.2f}"
            table.cell(i, 4).text = str(k.get('high_dev_count', 0))
            table.cell(i, 5).text = str(k.get('no_remark_count', 0))

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_material_top10(self, prs, pre_data, factory, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_title(slide, f"物料偏差金额 Top10 — {factory}")

        top10 = pre_data['material_top10'].get(factory, {})

        if not top10:
            tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
            tf.text_frame.paragraphs[0].text = f"该工厂无有效物料记录"
            self._add_page_number(slide, page_num, total_pages)
            return slide

        rows = len(top10) + 1
        cols = 2
        table = slide.shapes.add_table(
            rows, cols,
            Inches(1.5), Inches(1.5),
            Inches(7), Inches(0.4 * rows)
        ).table

        table.cell(0, 0).text = '物料名称'
        table.cell(0, 1).text = '偏差金额(元)'
        table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(0, 1).text_frame.paragraphs[0].font.bold = True

        for i, (name, amt) in enumerate(top10.items(), 1):
            table.cell(i, 0).text = str(name)
            table.cell(i, 1).text = f"{amt:,.2f}"

        if len(top10) < 10:
            note = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(7), Inches(0.5))
            note.text_frame.paragraphs[0].text = f"注：该工厂仅有 {len(top10)} 条有效物料记录"
            note.text_frame.paragraphs[0].font.size = Pt(11)

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_deviation_type_chart(self, prs, pre_data, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_title(slide, "偏差类型分布")
        factories = pre_data['factories']
        dev_type_dist = pre_data['dev_type_dist']

        if len(factories) == 2:
            # 并排两个饼图
            positions = [Inches(0.3), Inches(5.0)]
        elif len(factories) == 1:
            positions = [Inches(3.0)]
        else:
            positions = [Inches(0.3)]

        for idx, fac in enumerate(factories):
            dist = dev_type_dist.get(fac, {})
            if not dist:
                # 无数据，显示占位
                tb = slide.shapes.add_textbox(positions[idx], Inches(1.5), Inches(4.5), Inches(0.5))
                tb.text_frame.paragraphs[0].text = f"{fac}：无数据"
                tb.text_frame.paragraphs[0].font.size = Pt(14)
                continue

            chart_data = CategoryChartData()
            chart_data.categories = list(dist.keys())
            chart_data.add_series('偏差金额', list(dist.values()))

            x, y, cx, cy = positions[idx], Inches(1.2), Inches(4.5), Inches(4.0)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

            # 添加数据标签
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.show_category_name = False

            # 标题
            lbl = slide.shapes.add_textbox(positions[idx], Inches(5.3), Inches(4.5), Inches(0.5))
            lbl.text_frame.paragraphs[0].text = f"{fac}"
            lbl.text_frame.paragraphs[0].font.size = Pt(14)
            lbl.text_frame.paragraphs[0].font.bold = True

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_workshop_details(self, prs, pre_data, factory, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_title(slide, f"车间偏差详情 — {factory}")

        ws_stats = pre_data['workshop_stats'].get(factory, {})
        dev_type_dist = pre_data['dev_type_dist'].get(factory, {})
        freq_reasons = pre_data['freq_reasons'].get(factory, {})

        if not ws_stats:
            tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
            tf.text_frame.paragraphs[0].text = f"该工厂无车间级数据"
            self._add_page_number(slide, page_num, total_pages)
            return slide

        rows = len(ws_stats) + 1
        cols = 2
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.5),
            Inches(4.5), Inches(0.4 * rows)
        ).table

        table.cell(0, 0).text = '车间'
        table.cell(0, 1).text = '偏差金额(元)'
        table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(0, 1).text_frame.paragraphs[0].font.bold = True

        for i, (ws, amt) in enumerate(ws_stats.items(), 1):
            table.cell(i, 0).text = str(ws)
            table.cell(i, 1).text = f"{amt:,.2f}"

        # 净偏差柱状图（Top5 车间）
        if ws_stats:
            chart_data = CategoryChartData()
            top5 = sorted(ws_stats.items(), key=lambda x: abs(x[1]), reverse=True)[:5]
            chart_data.categories = [str(k) for k, v in top5]
            chart_data.add_series('偏差金额', [abs(v) for k, v in top5])
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(5.2), Inches(1.5), Inches(4.2), Inches(3.5),
                chart_data
            ).chart
            chart.has_legend = False

        # 高频原因（右侧下方）
        if freq_reasons:
            tb = slide.shapes.add_textbox(Inches(5.2), Inches(5.2), Inches(4.2), Inches(2))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Top5 偏差原因："
            p.font.bold = True
            p.font.size = Pt(11)
            for reason, count in list(freq_reasons.items())[:5]:
                p = tf.add_paragraph()
                p.text = f"  {reason} ({count}次)"
                p.font.size = Pt(10)

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_material_type_net(self, prs, pre_data, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_title(slide, "物料类型净偏差分布（Slide 10）")

        mat_type_net = pre_data['material_type_net']
        factories = pre_data['factories']

        # 收集所有物料类型
        all_types = set()
        for fac_net in mat_type_net.values():
            all_types.update(fac_net.keys())
        all_types = sorted(all_types)

        if not all_types:
            tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
            tf.text_frame.paragraphs[0].text = "无物料类型数据"
            self._add_page_number(slide, page_num, total_pages)
            return slide

        # 柱状图
        chart_data = CategoryChartData()
        chart_data.categories = all_types
        for fac in factories:
            net = mat_type_net.get(fac, {})
            chart_data.add_series(str(fac), [net.get(t, 0) for t in all_types])

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(1.5), Inches(9), Inches(4.5),
            chart_data
        ).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        # 表格备用
        rows = len(all_types) + 1
        cols = len(factories) + 1
        if cols <= 6:
            tbl = slide.shapes.add_table(
                rows, cols,
                Inches(0.5), Inches(6.2),
                Inches(9), Inches(0.3 * rows)
            ).table
            tbl.cell(0, 0).text = '物料类型'
            tbl.cell(0, 0).text_frame.paragraphs[0].font.bold = True
            for j, fac in enumerate(factories, 1):
                tbl.cell(0, j).text = str(fac)
                tbl.cell(0, j).text_frame.paragraphs[0].font.bold = True
            for i, mtype in enumerate(all_types, 1):
                tbl.cell(i, 0).text = mtype
                for j, fac in enumerate(factories, 1):
                    val = mat_type_net.get(fac, {}).get(mtype, 0)
                    tbl.cell(i, j).text = f"{val:,.2f}"

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_no_remark_warning(self, prs, pre_data, factory, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_title(slide, f"无备注预警 Top10 — {factory}（≥5万元）")

        warn_df = pre_data['no_remark_warning'].get(factory, pd.DataFrame())

        if warn_df.empty:
            tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
            tf.text_frame.paragraphs[0].text = f"该工厂无≥5万元的无备注记录"
            self._add_page_number(slide, page_num, total_pages)
            return slide

        COL_MAP = {}
        for target, candidates in [
            ('物料名称', ['组件物料描述', '物料名称', '物料描述']),
            ('车间', ['车间', 'workshop']),
            ('偏差金额', ['偏差金额', '金额']),
            ('偏差率', ['偏差率(%)', '偏差率']),
        ]:
            for c in candidates:
                if c in warn_df.columns:
                    COL_MAP[target] = c
                    break

        mat_col = COL_MAP.get('物料名称')
        ws_col = COL_MAP.get('车间')
        amt_col = COL_MAP.get('偏差金额')
        rate_col = COL_MAP.get('偏差率')

        rows = min(len(warn_df), 10) + 1
        cols = 4
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.3), Inches(1.5),
            Inches(9.4), Inches(0.4 * rows)
        ).table

        headers = ['物料名称', '车间', '偏差金额(元)', '偏差率(%)']
        for j, h in enumerate(headers):
            table.cell(0, j).text = h
            table.cell(0, j).text_frame.paragraphs[0].font.bold = True

        for i, (_, row) in enumerate(warn_df.iterrows(), 1):
            table.cell(i, 0).text = str(row.get(mat_col, ''))[:50]
            table.cell(i, 1).text = str(row.get(ws_col, ''))
            table.cell(i, 2).text = f"{row.get(amt_col, 0):,.2f}"
            table.cell(i, 3).text = f"{row.get(rate_col, 0):.2f}"

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_abnormal_warning(self, prs, pre_data, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_title(slide, "异常预警明细（定额>0 且 实际=0）")

        abnormal_df = pre_data.get('abnormal_warning', pd.DataFrame())

        if abnormal_df.empty:
            tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
            tf.text_frame.paragraphs[0].text = "无异常预警记录"
            self._add_page_number(slide, page_num, total_pages)
            return slide

        display_df = abnormal_df.head(20)
        factories = pre_data['factories']
        factory_col = pre_data.get('factory_col')

        COL_MAP = {}
        for target, candidates in [
            ('物料名称', ['组件物料描述', '物料名称', '物料描述']),
            ('车间', ['车间', 'workshop']),
            ('定额', ['数量 - 定额', '定额数量', 'quota']),
            ('实际', ['数量 - 实际', '实际数量', 'actual']),
        ]:
            for c in candidates:
                if c in display_df.columns:
                    COL_MAP[target] = c
                    break

        mat_col = COL_MAP.get('物料名称')
        ws_col = COL_MAP.get('车间')
        quota_col = COL_MAP.get('定额')
        actual_col = COL_MAP.get('实际')

        cols = 5
        rows = len(display_df) + 1
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.3), Inches(1.5),
            Inches(9.4), Inches(0.35 * rows)
        ).table

        headers = ['工厂', '车间', '物料名称', '定额', '实际']
        for j, h in enumerate(headers):
            table.cell(0, j).text = h
            table.cell(0, j).text_frame.paragraphs[0].font.bold = True

        for i, (_, row) in enumerate(display_df.iterrows(), 1):
            factory_val = str(row.get(factory_col, '')) if factory_col else ''
            table.cell(i, 0).text = factory_val
            table.cell(i, 1).text = str(row.get(ws_col, ''))
            table.cell(i, 2).text = str(row.get(mat_col, ''))[:40]
            table.cell(i, 3).text = str(row.get(quota_col, 0))
            table.cell(i, 4).text = str(row.get(actual_col, 0))

        if len(abnormal_df) > 20:
            note = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.5))
            note.text_frame.paragraphs[0].text = f"注：共 {len(abnormal_df)} 条异常记录，仅显示前 20 条"
            note.text_frame.paragraphs[0].font.size = Pt(10)

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_freq_reason(self, prs, pre_data, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_title(slide, "高频偏差原因对比")
        factories = pre_data['factories']
        freq_reasons = pre_data['freq_reasons']

        if not freq_reasons or not any(freq_reasons.values()):
            tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
            tf.text_frame.paragraphs[0].text = "无偏差原因数据"
            self._add_page_number(slide, page_num, total_pages)
            return slide

        if len(factories) == 2:
            positions = [Inches(0.3), Inches(5.0)]
        elif len(factories) == 1:
            positions = [Inches(2)]
        else:
            positions = [Inches(0.3)]

        for idx, fac in enumerate(factories):
            reasons = freq_reasons.get(fac, {})
            if not reasons:
                tb = slide.shapes.add_textbox(positions[idx], Inches(1.5), Inches(4.5), Inches(0.5))
                tb.text_frame.paragraphs[0].text = f"{fac}：无数据"
                continue

            rows = len(reasons) + 1
            table = slide.shapes.add_table(
                rows, 2,
                positions[idx], Inches(1.5),
                Inches(4.5), Inches(0.4 * rows)
            ).table
            table.cell(0, 0).text = f'{fac} Top5 原因'
            table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
            table.cell(0, 1).text = '次数'
            table.cell(0, 1).text_frame.paragraphs[0].font.bold = True
            for i, (reason, count) in enumerate(reasons.items(), 1):
                table.cell(i, 0).text = str(reason)[:50]
                table.cell(i, 1).text = str(count)

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def _add_summary(self, prs, pre_data, page_num, total_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        self._set_slide_title(slide, "总结与改进建议")

        tf = slide.placeholders[1].text_frame
        tf.clear()

        kpis = pre_data['total_kpis']
        p = tf.paragraphs[0]
        p.text = "【总体发现】"
        p.font.bold = True
        p.font.size = Pt(16)

        for item in [
            f"  总记录数：{kpis['total_records']:,} 条，分布在 {kpis['factories']} 个工厂",
            f"  总偏差金额：{kpis['total_amount']:,.2f} 元",
            f"  高偏差（>10%）记录：{kpis['high_dev_count']} 条",
            "",
        ]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(13)

        for fac, fac_kpis in pre_data['factory_kpis'].items():
            p = tf.add_paragraph()
            p.text = f"【{fac} 专项建议】"
            p.font.bold = True
            p.font.size = Pt(16)
            for item in [
                f"  无备注记录：{fac_kpis['no_remark_count']} 条，需优先补充",
                f"  高偏差物料：见 Top10 页面，重点关注偏差率>10%的物料",
                "  建议：组织车间级物料管控培训，完善领料流程",
                "",
            ]:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(13)

        self._add_page_number(slide, page_num, total_pages)
        return slide

    def generate_ppt(self, output_path: str = None, progress_callback=None):
        """生成 PPT 报告（v1.2 优化版，直接读取当前审核数据）
        
        Args:
            output_path: PPT 输出路径（若为 None，从 View 获取 output_path）
            progress_callback: 进度回调函数 (0-100)
        
        Returns:
            output_path: 生成的 PPT 文件路径
        """
        import pandas as pd
        from pptx import Presentation

        # 1. 获取当前 DataFrame
        df = self.view.get_audit_data()
        if df is None or df.empty:
            raise ValueError("当前无审核数据，无法生成 PPT")

        if len(df) > 50000:
            raise Exception("数据量超过 5 万条，请缩小范围")

        # 2. 确定输出路径
        if output_path is None:
            import os
            out_dir = self.view.get_output_path()
            if out_dir:
                base = os.path.splitext(os.path.basename(out_dir))[0]
            else:
                base = "ZPP011偏差分析"
            out_dir = os.path.join(os.getcwd(), 'output', 'pptx')
            os.makedirs(out_dir, exist_ok=True)
            import datetime
            date_tag = datetime.datetime.now().strftime('%Y%m%d_%H%M')
            output_path = os.path.join(out_dir, f"ZPP011偏差分析_{date_tag}.pptx")

        # 3. 数据预处理
        pre_data = self._pre_aggregate_data(df)
        factories = pre_data['factories']

        # 4. 加载模板（fallback 机制）
        import os
        template_path = os.path.join(os.getcwd(), 'config', 'template.pptx')
        if os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = Presentation()

        # 5. 构建页面生成顺序（动态工厂数量）
        pages = []

        # P1 封面
        pages.append(('cover', None))
        # P2 目录
        pages.append(('toc', None))

        if not factories:
            # 无工厂数据时添加占位页
            pages.append(('text', ('无数据', '无有效工厂数据')))

        # P3 核心指标
        pages.append(('kpi', None))
        # P4 工厂对比
        pages.append(('factory_cmp', None))

        # 物料 Top10（每工厂一页）
        for fac in factories:
            pages.append(('mat_top10', fac))

        # 偏差类型饼图
        pages.append(('dev_type_chart', None))

        # 车间详情（每工厂一页）
        for fac in factories:
            pages.append(('workshop', fac))

        # Slide 10 物料类型净偏差
        pages.append(('mat_type_net', None))

        # 无备注预警（每工厂一页）
        for fac in factories:
            pages.append(('no_remark', fac))

        # 异常预警
        pages.append(('abnormal', None))

        # 高频原因
        pages.append(('freq_reason', None))

        # 总结
        pages.append(('summary', None))

        total_pages = len(pages)

        # 6. 按顺序生成页面
        for i, (page_type, factory) in enumerate(pages):
            try:
                if page_type == 'cover':
                    self._add_cover(prs, pre_data, i+1, total_pages)
                elif page_type == 'toc':
                    self._add_toc(prs, pre_data, i+1, total_pages)
                elif page_type == 'kpi':
                    self._add_kpi_overview(prs, pre_data, i+1, total_pages)
                elif page_type == 'factory_cmp':
                    self._add_factory_comparison(prs, pre_data, i+1, total_pages)
                elif page_type == 'mat_top10':
                    self._add_material_top10(prs, pre_data, factory, i+1, total_pages)
                elif page_type == 'dev_type_chart':
                    self._add_deviation_type_chart(prs, pre_data, i+1, total_pages)
                elif page_type == 'workshop':
                    self._add_workshop_details(prs, pre_data, factory, i+1, total_pages)
                elif page_type == 'mat_type_net':
                    self._add_material_type_net(prs, pre_data, i+1, total_pages)
                elif page_type == 'no_remark':
                    self._add_no_remark_warning(prs, pre_data, factory, i+1, total_pages)
                elif page_type == 'abnormal':
                    self._add_abnormal_warning(prs, pre_data, i+1, total_pages)
                elif page_type == 'freq_reason':
                    self._add_freq_reason(prs, pre_data, i+1, total_pages)
                elif page_type == 'summary':
                    self._add_summary(prs, pre_data, i+1, total_pages)
                elif page_type == 'text':
                    title, msg = factory if factory else ('无数据', '')
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    self._set_slide_title(slide, title)
                    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
                    tb.text_frame.paragraphs[0].text = msg
                    self._add_page_number(slide, i+1, total_pages)
            except Exception as e:
                self.view.log(f"[PPT] 生成第{i+1}页失败: {e}", "warn")

            if progress_callback:
                progress_callback(int((i + 1) / total_pages * 100))

        # 7. 保存
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
        prs.save(output_path)
        self.view.log(f"[PPT] 已生成: {output_path}", "success")
        return output_path


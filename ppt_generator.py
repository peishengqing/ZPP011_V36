#!/usr/bin/env python
# -*- coding: utf-8 -*-
""" PPT 生成模块（GUI 错误日志版） """

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
import os
import re
import traceback
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
import io
from matplotlib.ticker import FuncFormatter
import warnings
warnings.filterwarnings("ignore", category=UserWarning, module="matplotlib")

# 错误日志路径
_ERROR_LOG = os.path.join(os.path.dirname(__file__), 'ppt_error.log')

def _log_error(msg):
    try:
        with open(_ERROR_LOG, 'a', encoding='utf-8') as f:
            f.write(f"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')} - {msg}\n")
    except:
        pass

def _log(msg):
    try:
        print(msg)
    except:
        pass

plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

def _calc_note_rate(summary_df):
    try:
        rate_col = summary_df['备注覆盖率']
        rates = rate_col.astype(str).str.replace('%', '').astype(float) / 100.0
        weights = summary_df['总条数'].astype(float)
        if weights.sum() > 0:
            return (rates * weights).sum() / weights.sum()
        return rates.mean()
    except:
        try:
            return float(rate_col.iloc[0].replace('%', '')) / 100.0
        except:
            return 0.0

def _generate_smart_summary(total_rows, pos_count, neg_count,
                            pos_qty, neg_qty, note_rate,
                            no_note_count, anomaly_count, repeat_count,
                            high_freq_count, factory_stats,
                            date_start, date_end):
    lines = []
    lines.append(f"【{date_start} ~ {date_end} 偏差分析报告】")
    lines.append("")
    lines.append(f"本期共涉及 {total_rows:,} 条偏差记录，其中正偏差（多耗）{pos_count:,} 条、负偏差（少耗）{neg_count:,} 条。")
    if pos_qty > 0 or neg_qty > 0:
        parts = []
        if pos_qty > 0:
            parts.append(f"多耗金额 ¥{pos_qty:,.0f}")
        if neg_qty > 0:
            parts.append(f"少耗金额 ¥{neg_qty:,.0f}")
        lines.append(f"偏差金额方面：{'；'.join(parts)}。")
    lines.append(f"整体备注覆盖率为 {note_rate:.1%}，")
    if note_rate < 0.2:
        lines.append("覆盖率偏低，建议各车间加强备注填写。")
    elif note_rate < 0.4:
        lines.append("覆盖率有待提升，部分车间需重点关注。")
    else:
        lines.append("整体表现良好。")
    alerts = []
    if no_note_count > 0:
        alerts.append(f"{no_note_count} 条高偏差无备注记录需跟进")
    if anomaly_count > 0:
        alerts.append(f"{anomaly_count} 条异常记录待核查")
    if repeat_count > 0:
        alerts.append(f"{repeat_count} 条屡犯记录需重点管控")
    if high_freq_count > 0:
        alerts.append(f"{high_freq_count} 种高频偏差原因")
    if alerts:
        lines.append("")
        lines.append("⚠️ 重点预警：" + "；".join(alerts) + "。")
    if factory_stats:
        lines.append("")
        lines.append("📊 各工厂偏差情况：")
        for fs in factory_stats[:5]:
            lines.append(f"  • {fs['工厂']}：{fs['记录数']}条记录，多耗{fs['多耗']:,.0f}，少耗{fs['少耗']:,.0f}")
    lines.append("")
    lines.append("建议结合各车间 Top5 原因分析，针对性制定改进措施。")
    return "\n".join(lines)

def _read_excel_data(excel_path, log_cb=None):
    _log("[PPT] 读取 Excel 数据...")
    from openpyxl import load_workbook
    wb = load_workbook(excel_path, data_only=True, read_only=True)
    data = {"file_name": os.path.basename(excel_path), "workshop_summary": [], "no_note_list": [], "total_rows": 0}
    tgt = next((n for n in wb.sheetnames if "汇总" in n), None)
    if tgt:
        ws = wb[tgt]
        rows = list(ws.iter_rows(min_row=2, values_only=True))
        for row in rows[1:]:
            if not row or not row[0]:
                continue
            data["workshop_summary"].append({
                "workshop": str(row[0]).strip() if row[0] else "",
                "mat_type": str(row[1]).strip() if len(row) > 1 and row[1] else "",
                "total_cnt": row[2] if len(row) > 2 else None,
                "pos_cnt": row[3] if len(row) > 3 else None,
                "neg_cnt": row[4] if len(row) > 4 else None,
                "total_dev": row[5] if len(row) > 5 else None,
                "note_cover": row[6] if len(row) > 6 else None,
            })
            try:
                data["total_rows"] += int(row[2]) if row[2] else 0
            except BaseException:
                pass
        _log(f"  [PPT] 汇总: {len(data['workshop_summary'])} 行")
    tgt2 = next((n for n in wb.sheetnames if "预警" in n or "无备注" in n), None)
    if tgt2:
        ws2 = wb[tgt2]
        rows2 = list(ws2.iter_rows(min_row=3, values_only=True))
        for i, row in enumerate(rows2[:10]):
            if not row or not row[0]:
                continue
            data["no_note_list"].append({
                "rank": i + 1,
                "workshop": str(row[0]).strip() if row[0] else "",
                "mat_type": str(row[1]).strip() if len(row) > 1 and row[1] else "",
                "material": str(row[2]).strip() if len(row) > 2 and row[2] else "",
                "dev_type": str(row[3]).strip() if len(row) > 3 and row[3] else "",
                "dev_rate": row[4] if len(row) > 4 else None,
            })
        _log(f"  [PPT] 预警: {len(data['no_note_list'])} 条")
    wb.close()
    _log("[PPT] Excel 读取完成")
    return data

def _get_blank_layout(prs):
    for lo in prs.slide_layouts:
        if 'blank' in lo.name.lower():
            return lo
    return prs.slide_layouts[-1]

def _add_slide_title(prs, text, color, log_cb=None):
    blank = _get_blank_layout(prs)
    s = prs.slides.add_slide(blank)
    s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.4)).fill.solid()
    s.shapes[0].fill.fore_color.rgb = RGBColor(*color)
    s.shapes[0].line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.4), Inches(0.2), prs.slide_width - Inches(0.8), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    _log(f"  ▶ {text[:10]} 页完成")

def _add_content_box(slide, lines, font_size=15, color=(0x2A, 0x2A, 0x2A), prs=None):
    if prs is None:
        raise ValueError("prs 参数不能为 None")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), prs.slide_width - Inches(1), prs.slide_height - Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.color.rgb = RGBColor(*color)

def _create_bar_chart_image(df_summary, factory_name, colors):
    _factory_col = '工厂' if '工厂' in df_summary.columns else ('工厂名称' if '工厂名称' in df_summary.columns else None)
    if _factory_col is None or _factory_col not in df_summary.columns:
        return None
    df_f = df_summary[df_summary[_factory_col].astype(str).str.strip() == str(factory_name).strip()].copy()
    if df_f.empty:
        return None
    _dev_col = '总偏差金额(含税)' if '总偏差金额(含税)' in df_summary.columns else '总偏差金额'
    df_f['abs_dev'] = df_f[_dev_col].abs()
    df_f = df_f.sort_values('abs_dev', ascending=True)
    workshops = df_f['车间'].tolist()
    dev_vals = df_f[_dev_col].tolist()
    fig, ax = plt.subplots(figsize=(8, 4))
    bar_colors = [colors['pos'] if v > 0 else colors['neg'] for v in dev_vals]
    ax.barh(workshops, dev_vals, color=bar_colors)
    ax.axvline(0, color='black', linewidth=0.8)
    ax.set_xlabel('偏差金额（元）')
    ax.set_title(f'{factory_name} 各车间偏差金额分布')
    ax.xaxis.set_major_formatter(FuncFormatter(lambda x, p: f'{x/1000:.0f}k' if abs(x) >= 1000 else f'{x:.0f}'))
    plt.tight_layout()
    img_buf = io.BytesIO()
    plt.savefig(img_buf, format='png', dpi=100, bbox_inches='tight')
    plt.close()
    img_buf.seek(0)
    return img_buf

def _create_trend_chart_image(trend_df, top_n=5):
    if trend_df.empty:
        return None
    df = trend_df.copy()
    for col in ['早期偏差率', '中期偏差率', '近期偏差率']:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col].astype(str).str.replace('%', ''), errors='coerce')
    df = df.dropna(subset=['早期偏差率', '中期偏差率', '近期偏差率'])
    if df.empty:
        return None
    df['change'] = (df['近期偏差率'] - df['早期偏差率']).abs()
    df = df.sort_values('change', ascending=False).head(top_n)
    if df.empty:
        return None
    fig, ax = plt.subplots(figsize=(8, 4))
    x = ['早期', '中期', '近期']
    for _, row in df.iterrows():
        y = [row['早期偏差率'], row['中期偏差率'], row['近期偏差率']]
        ax.plot(x, y, marker='o', label=row['物料名称'][:12])
    ax.axhline(0, color='gray', linestyle='--', linewidth=0.5)
    ax.set_ylabel('偏差率 (%)')
    ax.set_title('关键物料偏差率趋势')
    ax.legend(loc='best', fontsize=8)
    ax.grid(True, linestyle=':', alpha=0.6)
    plt.tight_layout()
    img_buf = io.BytesIO()
    plt.savefig(img_buf, format='png', dpi=100, bbox_inches='tight')
    plt.close()
    img_buf.seek(0)
    return img_buf


def _add_doughnut_chart(slide, over_amount, under_amount):
    """添加多耗/少耗环形饼图"""
    if over_amount == 0 and under_amount == 0:
        return  # 无数据，跳过

    chart_data = ChartData()
    chart_data.categories = ['多耗', '少耗']
    chart_data.add_series('偏差金额', [over_amount, under_amount])

    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    ).chart

    # 设置标题
    chart.has_title = True
    chart.chart_title.text_frame.text = "多耗与少耗占比"

    # 设置数据点颜色
    series = chart.series[0]
    points = series.points
    if len(points) >= 2:
        # 多耗 -> 红色
        points[0].format.fill.solid()
        points[0].format.fill.fore_color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
        # 少耗 -> 绿色
        points[1].format.fill.solid()
        points[1].format.fill.fore_color.rgb = RGBColor(0x2E, 0xCC, 0x71)


def run_ppt_generation(excel_path, output_path, log_cb=None):
    _log_error(f"开始生成 PPT，Excel: {excel_path}, 输出: {output_path}")
    try:
        xl = pd.ExcelFile(excel_path)
        sheets = xl.sheet_names
        _log(f"[PPT] Excel 完整路径: {excel_path}")
        _log(f"[PPT] Excel 工作表列表: {sheets}")

        def safe_read(name, default_cols=None):
            if name in sheets:
                return pd.read_excel(excel_path, sheet_name=name)
            for s in sheets:
                if name in s:
                    return pd.read_excel(excel_path, sheet_name=s)
            _log(f"  [PPT] 未找到 Sheet「{name}」，将跳过相关内容")
            return pd.DataFrame(columns=default_cols or [])

        summary_df = safe_read('汇总统计')
        _log(f"  [PPT] summary_df: {summary_df.shape[0]} rows, columns: {summary_df.columns.tolist()[:5]}")
        no_note_df = safe_read('无备注预警')
        abnormal_df = safe_read('异常预警')
        freq_loss_df = safe_read('原因汇总')
        try:
            cause_top5_df = safe_read('偏差原因分析')
        except:
            cause_top5_df = pd.DataFrame()
        dev_detail_df = safe_read('完整偏差明细')
        alt_df = safe_read('替代料明细')
        try:
            cause_analysis_df = safe_read('偏差原因分析')
        except:
            cause_analysis_df = pd.DataFrame()
        info_df = safe_read('分析说明')
        trend_df = safe_read('趋势分析')

        if not summary_df.empty:
            total_rows = int(summary_df['总条数'].sum())
            _log(f"  [PPT] total_rows={total_rows}")
            total_pos_count = int(summary_df['正偏差条数'].sum())
            total_neg_count = int(summary_df['负偏差条数'].sum())
            total_pos_qty = summary_df['正偏差数量'].sum()
            total_neg_qty = summary_df['负偏差数量'].sum()
            total_pos_amount = summary_df['正偏差金额(含税)'].sum()
            total_neg_amount = abs(summary_df['负偏差金额(含税)'].sum())
            note_rate_avg = _calc_note_rate(summary_df)
        else:
            total_rows = total_pos_count = total_neg_count = 0
            total_pos_qty = total_neg_qty = 0
            total_pos_amount = total_neg_amount = 0
            note_rate_avg = 0

        no_note_count = len(no_note_df)
        anomaly2_count = len(abnormal_df) if not abnormal_df.empty else 0

        high_freq_count = len(freq_loss_df)
        high_freq_total = freq_loss_df['净偏差'].abs().sum() if not freq_loss_df.empty and '净偏差' in freq_loss_df.columns else 0

        repeat_count = 0
        if not dev_detail_df.empty and '是否屡犯' in dev_detail_df.columns:
            repeat_count = dev_detail_df[dev_detail_df['是否屡犯'].str.startswith('屡犯')].shape[0]

        # 日期提取
        date_str_start = 'N/A'
        date_str_end = 'N/A'
        date_extracted = False
        if not info_df.empty and info_df.shape[0] > 1:
            for col in info_df.columns:
                val = str(info_df.iloc[1][col]) if info_df.shape[1] > 0 else ''
                if '～' in val or '~' in val:
                    date_match = re.search(r'(\d{4}[-/]\d{2}[-/]\d{2})\s*[～~]\s*(\d{4}[-/]\d{2}[-/]\d{2})', val)
                    if date_match:
                        date_str_start = date_match.group(1).replace('/', '-')
                        date_str_end = date_match.group(2).replace('/', '-')
                        date_extracted = True
                        _log(f"  [PPT] 从分析说明提取日期: {date_str_start} ~ {date_str_end}")
                        break
        if not date_extracted and '订单日期' in dev_detail_df.columns and not dev_detail_df.empty:
            dates = pd.to_datetime(dev_detail_df['订单日期'], errors='coerce')
            if not dates.isna().all():
                date_str_start = dates.min().strftime('%Y-%m-%d')
                date_str_end = dates.max().strftime('%Y-%m-%d')
                date_extracted = True
                _log(f"  [PPT] 从订单日期提取日期: {date_str_start} ~ {date_str_end}")
        if not date_extracted:
            file_date_match = re.search(r'(\d{8})-(\d{8})', os.path.basename(excel_path))
            if file_date_match:
                d1, d2 = file_date_match.group(1), file_date_match.group(2)
                date_str_start = f"{d1[:4]}-{d1[4:6]}-{d1[6:8]}"
                date_str_end = f"{d2[:4]}-{d2[4:6]}-{d2[6:8]}"
                date_extracted = True
                _log(f"  [PPT] 从文件名提取日期: {date_str_start} ~ {date_str_end}")
        if not date_extracted:
            today_str = datetime.now().strftime('%Y-%m-%d')
            date_str_start = date_str_end = today_str
            _log(f"  [PPT] 使用当前日期: {date_str_start}")

        # 工厂统计
        factory_stats = []
        if not summary_df.empty:
            _grp_factory_col = '工厂' if '工厂' in summary_df.columns else ('工厂名称' if '工厂名称' in summary_df.columns else '工厂')
            for factory, grp in summary_df.groupby(_grp_factory_col):
                factory_stats.append({
                    '工厂': factory,
                    '记录数': grp['总条数'].sum(),
                    '多耗': grp['正偏差数量'].sum(),
                    '少耗': abs(grp['负偏差数量'].sum()),
                    '净偏差': grp['总数量'].sum(),
                    '备注率': f"{grp['总条数'].sum()/total_rows:.0%}" if total_rows else "N/A"
                })

        workshop_top5 = []
        if not cause_analysis_df.empty and '车间' in cause_analysis_df.columns:
            circled = ['①', '②', '③', '④', '⑤', '⑥', '⑦', '⑧', '⑨', '⑩']
            for (factory, workshop), grp in cause_analysis_df.groupby([_grp_factory_col, '车间']):
                agg_grp = grp.groupby('备注原因').agg({
                    '涉及物料数': 'sum',
                    '净偏差': 'sum'
                }).reset_index()
                top5 = agg_grp.nlargest(5, '涉及物料数')
                top5_lines = []
                for idx, (_, row) in enumerate(top5.iterrows(), 1):
                    prefix = circled[idx-1] if idx <= len(circled) else f"{idx}."
                    top5_lines.append(
                        f"{prefix} {row['备注原因']}（{int(row['涉及物料数'])}次，净偏差 {row['净偏差']:,.0f}）"
                    )
                workshop_top5.append({
                    '工厂': factory,
                    '车间': workshop,
                    '多耗': pd.to_numeric(grp['多耗'], errors='coerce').fillna(0).sum(),
                    '少耗': abs(pd.to_numeric(grp['少耗'], errors='coerce').fillna(0).sum()),
                    '原因数': len(grp),
                    'top5': top5_lines
                })

        # 智能解读
        summary_text = _generate_smart_summary(
            total_rows, total_pos_count, total_neg_count,
            total_pos_qty, total_neg_qty,
            note_rate_avg, no_note_count, anomaly2_count,
            repeat_count, high_freq_count, factory_stats,
            date_str_start, date_str_end
        )
        _log(f"  [PPT] 智能解读生成完成 ({len(summary_text)} 字符)")

        # 创建 PPT
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        C_PRIMARY = RGBColor(30, 60, 114)
        C_ACCENT = RGBColor(0, 120, 215)
        C_RED = RGBColor(220, 53, 69)
        C_GREEN = RGBColor(40, 167, 69)
        C_ORANGE = RGBColor(255, 152, 0)
        C_WHITE = RGBColor(255, 255, 255)
        C_TEXT = RGBColor(51, 51, 51)
        C_BG = RGBColor(248, 249, 250)
        C_CARD_BG = RGBColor(255, 255, 255)

        def set_slide_bg(slide, color):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = color

        def add_title_bar(slide, title, subtitle=''):
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
            bar.fill.solid()
            bar.fill.fore_color.rgb = C_PRIMARY
            bar.line.fill.background()
            txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.7))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            if subtitle:
                txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12), Inches(0.4))
                tf2 = txBox2.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = subtitle
                p2.font.size = Pt(14)
                p2.font.color.rgb = RGBColor(180, 200, 230)

        def add_card(slide, left, top, width, height, bg_color=C_CARD_BG):
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            card.fill.solid()
            card.fill.fore_color.rgb = bg_color
            card.line.color.rgb = RGBColor(230, 230, 230)
            card.line.width = Pt(1)
            return card

        def add_kpi_card(slide, left, top, value, label, color=C_ACCENT):
            add_card(slide, left, top, Inches(2.8), Inches(1.6))
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.1), Inches(1.6))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = color
            stripe.line.fill.background()
            txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), Inches(2.3), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = color
            txBox2 = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.0), Inches(2.3), Inches(0.4))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = label
            p2.font.size = Pt(13)
            p2.font.color.rgb = RGBColor(120, 120, 120)

        def add_table(slide, left, top, width, headers, rows, col_widths=None):
            n_rows = len(rows) + 1
            n_cols = len(headers)
            row_h = Inches(0.45)
            tbl_h = row_h * n_rows
            table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, tbl_h)
            table = table_shape.table
            if col_widths:
                for i, w in enumerate(col_widths):
                    table.columns[i].width = w
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(h)
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_PRIMARY
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = C_WHITE
                    p.alignment = PP_ALIGN.CENTER
            for r_idx, row in enumerate(rows, 1):
                for c_idx, val in enumerate(row):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = str(val)
                    if r_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 245, 255)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(11)
                        p.font.color.rgb = C_TEXT
                        p.alignment = PP_ALIGN.CENTER
                    if c_idx in [2, 3, 4] and isinstance(val, (int, float)) and val < 0:
                        for p in cell.text_frame.paragraphs:
                            p.font.color.rgb = C_GREEN

        # 新增：多耗/少耗环形饼图
        over_total = 0.0
        under_total = 0.0
        if not dev_detail_df.empty and '偏差金额' in dev_detail_df.columns:
            over_total = dev_detail_df[dev_detail_df['偏差金额'] > 0]['偏差金额'].sum()
            under_total = dev_detail_df[dev_detail_df['偏差金额'] < 0]['偏差金额'].abs().sum()

        if over_total > 0 or under_total > 0:
            doughnut_slide = prs.slides.add_slide(_get_blank_layout(prs))
            set_slide_bg(doughnut_slide, C_BG)
            add_title_bar(doughnut_slide, '多耗与少耗占比', '偏差金额分布')
            _add_doughnut_chart(doughnut_slide, over_total, under_total)
            _log("  [PPT] 环形饼图完成")

        # 页面生成
        slide1 = prs.slides.add_slide(_get_blank_layout(prs))
        set_slide_bg(slide1, C_PRIMARY)
        deco = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
        deco.fill.solid(); deco.fill.fore_color.rgb = C_ACCENT; deco.line.fill.background()
        txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10), Inches(1.5))
        tf = txBox.text_frame; p = tf.paragraphs[0]
        p.text = 'ZPP011 生产偏差分析报告'; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = C_WHITE
        txBox2 = slide1.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(10), Inches(0.6))
        tf2 = txBox2.text_frame; p2 = tf2.paragraphs[0]
        p2.text = f'{date_str_start} ~ {date_str_end}'; p2.font.size = Pt(24); p2.font.color.rgb = RGBColor(180,200,230)
        txBox3 = slide1.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(10), Inches(1))
        tf3 = txBox3.text_frame; p3 = tf3.paragraphs[0]
        p3.text = f'数据范围：{total_rows:,} 条记录 | 多个工厂与车间'; p3.font.size = Pt(16); p3.font.color.rgb = RGBColor(150,170,200)
        _log("  [PPT] 封面完成")

        slide_info = prs.slides.add_slide(_get_blank_layout(prs))
        set_slide_bg(slide_info, C_BG)
        add_title_bar(slide_info, '分析说明')
        th_method = "固定阈值（公司规定）：±10%"
        th_val = "±10.0%"
        if not info_df.empty and info_df.shape[0] > 1:
            if info_df.shape[1] > 0:
                th_method = str(info_df.iloc[1, 0]) if pd.notna(info_df.iloc[1, 0]) else th_method
            if info_df.shape[1] > 1:
                th_val = str(info_df.iloc[1, 1]) if pd.notna(info_df.iloc[1, 1]) else th_val
        sheets_desc = [
            ("汇总统计", "按车间×物料分类统计偏差条数、数量、金额、备注覆盖率"),
            ("替代料明细", "识别到的替代料配对及净偏差"),
            ("无备注预警", "偏差率超过阈值但未填备注的记录"),
            ("偏差金额分析", "按物料汇总正/负偏差金额"),
            ("偏差原因分析", "按标准原因类别汇总分析"),
            ("趋势分析", "按10天周期分析偏差率趋势"),
        ]
        info_text = f"分析日期范围：{date_str_start} ～ {date_str_end}\n动态阈值方法：{th_method}\n动态阈值数值：{th_val}\n\n各Sheet功能说明：\n" + "\n".join([f"  • {s[0]}：{s[1]}" for s in sheets_desc])
        txBox = slide_info.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for line in info_text.split('\n'):
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
            p.font.color.rgb = C_TEXT
        _log("  [PPT] 分析说明页完成")

        slide2 = prs.slides.add_slide(_get_blank_layout(prs))
        set_slide_bg(slide2, C_BG)
        add_title_bar(slide2, '总体概况 & 智能解读')
        add_kpi_card(slide2, Inches(0.5), Inches(1.6), f'{total_rows:,}', '总记录数', C_PRIMARY)
        add_kpi_card(slide2, Inches(3.6), Inches(1.6), f'{total_pos_count:,}', '正偏差（多耗）', C_RED)
        add_kpi_card(slide2, Inches(6.7), Inches(1.6), f'{total_neg_count:,}', '负偏差（少耗）', C_GREEN)
        add_kpi_card(slide2, Inches(9.8), Inches(1.6), f'{note_rate_avg:.1%}' if total_rows else 'N/A', '备注覆盖率', C_ACCENT)
        add_kpi_card(slide2, Inches(0.5), Inches(3.4), f'¥{total_pos_qty:,.0f}', '正偏差金额', C_RED)
        add_kpi_card(slide2, Inches(3.6), Inches(3.4), f'¥{total_neg_qty:,.0f}', '负偏差金额', C_GREEN)
        if summary_text:
            txBox = slide2.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(12), Inches(2))
            tf = txBox.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = summary_text; p.font.size = Pt(12); p.font.color.rgb = C_TEXT
        _log("  [PPT] 概况页完成")

        slide_chart = prs.slides.add_slide(_get_blank_layout(prs))
        set_slide_bg(slide_chart, C_BG)
        add_title_bar(slide_chart, '各车间偏差金额分布')
        factory_col = '工厂' if '工厂' in summary_df.columns else ('工厂名称' if '工厂名称' in summary_df.columns else None)
        if factory_col is None:
            _log(" [PPT] 警告：找不到工厂列，无法生成车间偏差金额分布图")
        else:
            factories = summary_df[factory_col].astype(str).str.strip().unique()
            _log(f" [PPT] 发现工厂: {list(factories)}")
            preferred = ['云南达利-食品厂', '云南达利-饮料厂']
            images_added = 0
            for idx, fac in enumerate(preferred):
                if fac in factories:
                    img = _create_bar_chart_image(summary_df, fac, {'pos': '#dc3545', 'neg': '#28a745'})
                    if img:
                        slide_chart.shapes.add_picture(img, Inches(0.5 + (idx * 5.8)), Inches(1.6), width=Inches(5.5))
                        images_added += 1
            if images_added == 0 and len(factories) >= 1:
                for idx, fac in enumerate(factories[:2]):
                    img = _create_bar_chart_image(summary_df, fac, {'pos': '#dc3545', 'neg': '#28a745'})
                    if img:
                        slide_chart.shapes.add_picture(img, Inches(0.5 + (idx * 5.8)), Inches(1.6), width=Inches(5.5))
                        images_added += 1
            if images_added == 0:
                txBox = slide_chart.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = "未找到有效的车间偏差数据，请检查 Excel 中是否存在相关数据。"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(128, 128, 128)
        _log("  [PPT] 条形图页完成")

        slide3 = prs.slides.add_slide(_get_blank_layout(prs))
        set_slide_bg(slide3, C_BG)
        add_title_bar(slide3, '工厂维度统计')
        if factory_stats:
            headers_fac = ['工厂', '记录数', '多耗', '少耗', '净偏差', '备注覆盖率']
            rows_fac = [[s['工厂'], s['记录数'], s['多耗'], s['少耗'], s['净偏差'], s['备注率']] for s in factory_stats]
            add_table(slide3, Inches(0.8), Inches(1.8), Inches(9.7), headers_fac, rows_fac,
                      [Inches(2.5), Inches(1.2), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)])
        _log("  [PPT] 工厂维度完成")

        slide4 = prs.slides.add_slide(_get_blank_layout(prs))
        set_slide_bg(slide4, C_BG)
        add_title_bar(slide4, '物料分类 & 高频低额预警')
        if not summary_df.empty:
            mat_cat_stats = summary_df.groupby('物料分类').agg(
                total_qty=('总数量', 'sum'),
                pos_qty=('正偏差数量', 'sum'),
                neg_qty=('负偏差数量', 'sum')
            ).reset_index()
            headers_mat = ['物料分类', '总偏差数量', '多耗', '少耗']
            rows_mat = [[r['物料分类'], r['total_qty'], r['pos_qty'], r['neg_qty']] for _, r in mat_cat_stats.iterrows()]
            add_table(slide4, Inches(0.8), Inches(1.8), Inches(6), headers_mat, rows_mat,
                      [Inches(2), Inches(2), Inches(1.5), Inches(1.5)])
        add_kpi_card(slide4, Inches(7.5), Inches(1.8), f'{high_freq_count:,}', '高频低额物料数', C_ORANGE)
        txBox = slide4.shapes.add_textbox(Inches(7.8), Inches(3.6), Inches(5), Inches(1))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = f'累计损失金额：{high_freq_total:,.0f} 元'; p.font.size = Pt(14); p.font.color.rgb = C_TEXT
        _log("  [PPT] 物料分类完成")

        if not dev_detail_df.empty and '物料名称' in dev_detail_df.columns and '偏差金额' in dev_detail_df.columns:
            slide_mat = prs.slides.add_slide(_get_blank_layout(prs))
            set_slide_bg(slide_mat, C_BG)
            add_title_bar(slide_mat, '物料偏差金额分析', 'Top5 正偏差 / Top5 负偏差')
            dev_mat = dev_detail_df.groupby('物料名称').agg({'偏差金额': 'sum'}).reset_index()
            dev_mat = dev_mat.sort_values('偏差金额', ascending=False)
            top5_pos = dev_mat.head(5)
            top5_neg = dev_mat.tail(5).sort_values('偏差金额')
            headers = ['物料名称', '偏差金额(元)']
            rows_pos = [[r['物料名称'][:20], f"{r['偏差金额']:,.2f}"] for _, r in top5_pos.iterrows()]
            rows_neg = [[r['物料名称'][:20], f"{r['偏差金额']:,.2f}"] for _, r in top5_neg.iterrows()]
            add_table(slide_mat, Inches(0.6), Inches(1.8), Inches(5.5), headers, rows_pos, [Inches(3.5), Inches(2.0)])
            add_table(slide_mat, Inches(6.8), Inches(1.8), Inches(5.5), headers, rows_neg, [Inches(3.5), Inches(2.0)])
            _log("  [PPT] 物料金额Top5完成")

        if not trend_df.empty:
            slide_trend = prs.slides.add_slide(_get_blank_layout(prs))
            set_slide_bg(slide_trend, C_BG)
            add_title_bar(slide_trend, '偏差趋势分析')
            trend_img = _create_trend_chart_image(trend_df)
            if trend_img:
                slide_trend.shapes.add_picture(trend_img, Inches(0.8), Inches(1.8), width=Inches(11))
            classification = ""
            if '趋势' in trend_df.columns:
                trend_counts = trend_df['趋势'].value_counts().head(5).items()
                classification = "\n".join([f"{k}: {v}" for k, v in trend_counts])
            if classification:
                txBox = slide_trend.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11), Inches(1.2))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = "趋势分类：\n" + classification
                p.font.size = Pt(10)
                p.font.color.rgb = C_TEXT
            _log("  [PPT] 趋势分析页完成")

        slide5 = prs.slides.add_slide(_get_blank_layout(prs))
        set_slide_bg(slide5, C_BG)
        add_title_bar(slide5, '预警分析', '无备注预警 · 异常预警 · 重复偏差')
        add_kpi_card(slide5, Inches(0.5), Inches(1.6), f'{no_note_count:,}', '无备注预警', C_RED)
        txBox = slide5.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(3.5), Inches(3))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = '偏差率超阈值且未备注'; p.font.size = Pt(14); p.font.color.rgb = C_TEXT
        add_kpi_card(slide5, Inches(4.5), Inches(1.6), f'{anomaly2_count:,}', '异常预警', C_ORANGE)
        txBox2 = slide5.shapes.add_textbox(Inches(4.8), Inches(3.4), Inches(3.5), Inches(3))
        tf2 = txBox2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; p2.text = '定额>0但未投料'; p2.font.size = Pt(14); p2.font.color.rgb = C_TEXT
        add_kpi_card(slide5, Inches(8.5), Inches(1.6), f'{repeat_count:,}', '重复偏差物料', C_PRIMARY)
        txBox3 = slide5.shapes.add_textbox(Inches(8.8), Inches(3.4), Inches(3.5), Inches(3))
        tf3 = txBox3.text_frame; tf3.word_wrap = True
        p3 = tf3.paragraphs[0]; p3.text = '多次出现偏差的物料'; p3.font.size = Pt(14); p3.font.color.rgb = C_TEXT
        _log("  [PPT] 预警分析完成")

        for ws in workshop_top5:
            slide = prs.slides.add_slide(_get_blank_layout(prs))
            set_slide_bg(slide, C_BG)
            add_title_bar(slide, f'{ws["工厂"]} · {ws["车间"]} 偏差原因 Top5')
            add_kpi_card(slide, Inches(0.5), Inches(1.6), f'{ws["多耗"]:,.1f}', '多耗', C_RED)
            add_kpi_card(slide, Inches(3.6), Inches(1.6), f'{ws["少耗"]:,.1f}', '少耗', C_GREEN)
            add_kpi_card(slide, Inches(6.7), Inches(1.6), f'{ws["多耗"]+ws["少耗"]:,.1f}', '净偏差', C_ACCENT)
            add_kpi_card(slide, Inches(9.8), Inches(1.6), f'{ws["原因数"]}', '原因数', C_PRIMARY)
            content_top = Inches(3.6)
            if ws['top5']:
                txBox_t = slide.shapes.add_textbox(Inches(0.6), content_top, Inches(12), Inches(0.5))
                tf_t = txBox_t.text_frame; p_t = tf_t.paragraphs[0]
                p_t.text = '[PPT] 主要偏差原因'; p_t.font.size = Pt(18); p_t.font.bold = True; p_t.font.color.rgb = C_PRIMARY
                content_top += Inches(0.6)
                for line in ws['top5']:
                    txBox_l = slide.shapes.add_textbox(Inches(1.0), content_top, Inches(11.5), Inches(0.55))
                    tf_l = txBox_l.text_frame; tf_l.word_wrap = True
                    p_l = tf_l.paragraphs[0]; p_l.text = line; p_l.font.size = Pt(16); p_l.font.color.rgb = C_TEXT
                    content_top += Inches(0.55)
            _log(f"  [PPT] 车间页：{ws['车间']}")

        slide_end = prs.slides.add_slide(_get_blank_layout(prs))
        set_slide_bg(slide_end, C_PRIMARY)
        txBox = slide_end.shapes.add_textbox(Inches(1.2), Inches(2), Inches(10), Inches(1.5))
        tf = txBox.text_frame; p = tf.paragraphs[0]
        p.text = '总结与建议'; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = C_WHITE
        txBox2 = slide_end.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(10), Inches(3))
        tf2 = txBox2.text_frame; tf2.word_wrap = True
        conclusions = [
            f'共 {total_rows:,} 条记录，备注覆盖率待提升',
            f'{no_note_count:,} 条无备注预警需跟进',
            f'高频低额损失 {high_freq_total:,.0f} 元，需关注系统性原因',
            f'重复偏差物料 {repeat_count} 个，建议重点排查',
        ]
        for i, c in enumerate(conclusions):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = f'▸ {c}'; p.font.size = Pt(16); p.font.color.rgb = RGBColor(200,215,240); p.space_after = Pt(14)

        prs.save(output_path)
        _log(f"[PPT] 详细汇报 PPT 已保存: {output_path} (共{len(prs.slides)}页)")
        _log_error("PPT 生成成功")
        return True
    except Exception as e:
        error_msg = f"PPT生成失败: {str(e)}\n{traceback.format_exc()}"
        _log_error(error_msg)
        if log_cb:
            log_cb(error_msg)
        return False

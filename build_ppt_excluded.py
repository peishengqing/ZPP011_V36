# -*- coding: utf-8 -*-
"""
ZPP011 偏差分析 PPT 生成脚本（剔除替代料口径）
==============================================
功能：
    读取 ZPP011 偏差分析 Excel（完整偏差明细），按 [是否替代料=是] 整条剔除 92 行后，
    重新计算 总/分厂/分物料/分车间/物料类型/高偏差率/预警 等指标，
    生成 11 页改进版分析报告 PPT（深青+琥珀配色，python-pptx 原生绘制）。

用法：
    python build_ppt_excluded.py
    可在下方 CONFIG 区修改输入/输出路径。

依赖：
    pandas, openpyxl, python-pptx
"""

import os
import datetime
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ===================== CONFIG =====================
INPUT_XLSX = r"E:\zpp011_v2\ZPP011偏差分析最终版_20260724_090611.xlsx"
OUTPUT_PPTX = r"E:\zpp011_v2\ZPP011偏差分析_改进版_剔除替代料_20260724.pptx"
FONT = "Microsoft YaHei"

# ===================== 配色 =====================
TEAL = RGBColor(0x0E, 0x7C, 0x86)
TEAL_D = RGBColor(0x0A, 0x5A, 0x62)
AMBER = RGBColor(0xE8, 0x83, 0x3A)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
INK = RGBColor(0x22, 0x2B, 0x2E)
GREY = RGBColor(0x6B, 0x77, 0x7A)
LIGHT = RGBColor(0xF2, 0xF6, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xEC, 0xF3, 0xF3)

PW = Inches(13.333)
PH = Inches(7.5)


# ===================== 工具函数 =====================
def set_run(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    # 设置东亚字体，避免中文回退为默认字体
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT})
    rPr.append(ea)


def txt(slide, x, y, w, h, text, size=14, color=INK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        set_run(r, size, color, bold)
    return tb


def rect(slide, x, y, w, h, fill, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def bg(slide, color):
    rect(slide, 0, 0, PW, PH, color)


def header(slide, idx, title, subtitle=None):
    """顶部色条 + 编号徽章 + 标题"""
    rect(slide, 0, 0, PW, Inches(1.15), TEAL)
    rect(slide, 0, Inches(1.15), PW, Inches(0.06), AMBER)
    # 编号徽章
    rect(slide, Inches(0.45), Inches(0.28), Inches(0.6), Inches(0.6), AMBER)
    txt(slide, Inches(0.45), Inches(0.28), Inches(0.6), Inches(0.6),
        f"{idx:02d}", 22, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    txt(slide, Inches(1.25), Inches(0.22), Inches(11), Inches(0.55),
        title, 26, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    if subtitle:
        txt(slide, Inches(1.27), Inches(0.74), Inches(11), Inches(0.35),
            subtitle, 13, RGBColor(0xCF, 0xE8, 0xE8), False)
    # 页脚
    txt(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
        "ZPP011 生产偏差分析 · 剔除替代料口径", 9, GREY)
    txt(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.35),
        f"{idx}/11", 9, GREY, False, PP_ALIGN.RIGHT)


def rate_num(x):
    if pd.isna(x):
        return np.nan
    s = str(x).replace("%", "").strip()
    return float(s) if s not in ("", "nan") else np.nan


def w_money(v):
    """万元，带正负号，保留一位小数"""
    return f"{v/10000:+,.1f}万"


# ===================== 数据计算 =====================
def load_metrics(xlsx):
    d = pd.read_excel(xlsx, "完整偏差明细")
    d["rd"] = d["偏差率"].apply(rate_num)
    d["rdn"] = d["净偏差率"].apply(rate_num)

    alt_mask = d["是否替代料"] == "是"
    alt = d[alt_mask].copy()
    keep = d[~alt_mask].copy()

    M = {}
    M["total_rows"] = len(d)
    M["alt_rows"] = int(alt_mask.sum())
    M["keep_rows"] = len(keep)

    M["tot_b"] = d["偏差金额"].sum()
    M["tot_n"] = d["净偏差金额"].sum()
    M["keep_b"] = keep["偏差金额"].sum()
    M["keep_n"] = keep["净偏差金额"].sum()
    M["alt_b"] = alt["偏差金额"].sum()
    M["alt_n"] = alt["净偏差金额"].sum()

    # 分厂（剔除后）
    g = keep.groupby("工厂").agg(行数=("工厂", "size"),
                                 偏差金额=("偏差金额", "sum"),
                                 净偏差金额=("净偏差金额", "sum"))
    M["by_factory"] = g

    # 分厂（全表，用于对比）
    gf = d.groupby("工厂").agg(偏差金额=("偏差金额", "sum"),
                               净偏差金额=("净偏差金额", "sum"))
    M["by_factory_all"] = gf

    # 物料 Top10（剔除后，净偏差金额）
    mt = keep.groupby("物料名称").agg(净偏差金额=("净偏差金额", "sum"),
                                      行数=("物料名称", "size")).sort_values(
        "净偏差金额", ascending=False).head(10)
    M["mat_top10"] = mt

    # 车间（剔除后）
    ws = keep.groupby(["工厂", "车间"]).agg(净偏差金额=("净偏差金额", "sum"),
                                           行数=("车间", "size"))
    M["by_workshop"] = ws

    # 物料类型（剔除后）
    mt_type = keep.groupby("物料类型").agg(行数=("物料类型", "size"),
                                          净偏差金额=("净偏差金额", "sum")).sort_values(
        "净偏差金额", ascending=False)
    M["by_mat_type"] = mt_type

    # 高偏差率（>=5%）
    hi = keep[keep["rd"] >= 5]
    M["hi_cnt"] = len(hi)
    M["hi_all_cnt"] = len(d[d["rd"] >= 5])
    # 主数据型（系统无额定/无定额）
    sys_mask = hi["备注"].fillna("").str.contains("系统无额定|系统无定额")
    M["hi_sys_cnt"] = int(sys_mask.sum())

    # 预警 sheets
    try:
        wm = pd.read_excel(xlsx, "无备注预警")
        M["warn_note"] = len(wm)
    except Exception:
        M["warn_note"] = 0
    try:
        ab = pd.read_excel(xlsx, "异常预警")
        M["warn_abn"] = len(ab)
    except Exception:
        M["warn_abn"] = 0
    try:
        mid = pd.read_excel(xlsx, "中间地带明细")
        M["warn_mid"] = len(mid)
    except Exception:
        M["warn_mid"] = 0

    return M, d, keep


# ===================== 各页构建 =====================
def build_cover(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, TEAL_D)
    rect(s, 0, 0, PW, Inches(0.18), AMBER)
    txt(s, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.6),
        "ZPP011 生产偏差分析报告", 40, WHITE, True)
    txt(s, Inches(0.82), Inches(1.75), Inches(11.5), Inches(0.5),
        "剔除替代料口径 · 改进版", 20, AMBER, True)
    txt(s, Inches(0.82), Inches(2.35), Inches(11.5), Inches(0.4),
        "数据来源：ZPP011偏差分析最终版_20260724_090611.xlsx", 13,
        RGBColor(0xCF, 0xE8, 0xE8))

    cards = [
        (f"{M['total_rows']:,}", "偏差记录总数", TEAL),
        (f"{M['alt_rows']}", "替代料记录（已剔除）", AMBER),
        (f"{M['keep_rows']:,}", "有效分析记录", RGBColor(0x3A, 0xA8, 0xB0)),
        (w_money(M['keep_n']), "剔除后总净偏差", RED if M['keep_n'] > 0 else GREEN),
    ]
    cw = Inches(2.85)
    gap = Inches(0.3)
    x0 = Inches(0.8)
    y0 = Inches(3.4)
    for i, (big, lab, col) in enumerate(cards):
        x = x0 + i * (cw + gap)
        rect(s, x, y0, cw, Inches(2.0), WHITE)
        rect(s, x, y0, cw, Inches(0.12), col)
        txt(s, x, y0 + Inches(0.35), cw, Inches(0.9), big, 30, col, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x, y0 + Inches(1.35), cw, Inches(0.5), lab, 13, GREY, False,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
        f"生成时间：{datetime.datetime.now():%Y-%m-%d %H:%M}  ·  云南达利生产基地",
        12, RGBColor(0xCF, 0xE8, 0xE8))


def build_toc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 2, "目录", "CONTENTS")
    items = [
        "核心指标总览", "替代料影响说明", "工厂维度对比",
        "物料偏差 Top10", "车间偏差详情", "物料类型分布",
        "偏差原因分析", "预警核查结果", "总结与改进建议",
    ]
    col_w = Inches(6.0)
    x1 = Inches(0.8)
    x2 = Inches(7.0)
    y0 = Inches(1.6)
    rh = Inches(0.55)
    for i, it in enumerate(items):
        col = 0 if i < 5 else 1
        row = i if i < 5 else i - 5
        x = x1 if col == 0 else x2
        y = y0 + row * rh
        rect(s, x, y + Inches(0.05), Inches(0.45), Inches(0.45), TEAL)
        txt(s, x, y + Inches(0.05), Inches(0.45), Inches(0.45),
            f"{i+1:02d}", 14, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.6), y, Inches(5.2), Inches(0.55), it, 16, INK, False,
            PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def build_metrics(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 3, "核心指标总览", "剔除替代料后两家工厂均为净省料")
    bf = M["by_factory"]
    cards = []
    cards.append(("总体", M["keep_n"], M["keep_rows"]))
    for f, row in bf.iterrows():
        cards.append((f, row["净偏差金额"], int(row["行数"])))
    cw = Inches(3.9)
    gap = Inches(0.35)
    x0 = Inches(0.8)
    y0 = Inches(1.7)
    for i, (name, net, rows) in enumerate(cards):
        x = x0 + i * (cw + gap)
        col = TEAL if i == 0 else (RED if net > 0 else GREEN)
        rect(s, x, y0, cw, Inches(2.4), CARD_BG)
        rect(s, x, y0, cw, Inches(0.14), col)
        txt(s, x, y0 + Inches(0.3), cw, Inches(0.5), name, 20, INK, True,
            PP_ALIGN.CENTER)
        txt(s, x, y0 + Inches(0.95), cw, Inches(0.9), w_money(net), 30, col, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x, y0 + Inches(1.9), cw, Inches(0.4), f"有效记录 {rows:,} 条",
            13, GREY, False, PP_ALIGN.CENTER)

    # 洞察条
    rect(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.9), LIGHT)
    rect(s, Inches(0.8), Inches(4.6), Inches(0.14), Inches(1.9), AMBER)
    txt(s, Inches(1.1), Inches(4.75), Inches(11.2), Inches(0.4),
        "关键洞察", 16, AMBER, True)
    txt(s, Inches(1.1), Inches(5.2), Inches(11.2), Inches(1.2),
        "• 全表剔除 92 行替代料后，总净偏差由 +89.8万 转为 −14.3万（净省料）。\n"
        "• 饮料厂原 +93.5万「超耗」几乎全部来自替代料对冲噪音，剔除后转为 −5.4万（省料）。\n"
        "• 食品厂剔除后 −8.9万，局部存在正向超耗物料，需重点核查（见物料 Top10）。",
        13, INK)


def build_alt_explain(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 4, "替代料影响说明", "为什么说剔除替代料前后结论完全相反")
    # 左：全表
    rect(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(2.6), CARD_BG)
    rect(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(0.14), GREY)
    txt(s, Inches(1.0), Inches(1.9), Inches(5.2), Inches(0.4),
        "全表（含替代料）", 16, INK, True)
    txt(s, Inches(1.0), Inches(2.45), Inches(5.2), Inches(0.5),
        f"总偏差金额(毛)  {w_money(M['tot_b'])}", 14, INK, True)
    txt(s, Inches(1.0), Inches(3.0), Inches(5.2), Inches(0.5),
        f"总净偏差金额      {w_money(M['tot_n'])}", 14, RED, True)
    txt(s, Inches(1.0), Inches(3.6), Inches(5.2), Inches(0.5),
        f"替代料 {M['alt_rows']} 行：正向 +{M['alt_b_pos']/10000:,.1f}万 / 负向 −{abs(M['alt_b_neg'])/10000:,.1f}万",
        12, GREY)

    # 右：剔除后
    rect(s, Inches(6.9), Inches(1.7), Inches(5.6), Inches(2.6), CARD_BG)
    rect(s, Inches(6.9), Inches(1.7), Inches(5.6), Inches(0.14), TEAL)
    txt(s, Inches(7.1), Inches(1.9), Inches(5.2), Inches(0.4),
        "剔除替代料后", 16, INK, True)
    txt(s, Inches(7.1), Inches(2.45), Inches(5.2), Inches(0.5),
        f"保留 {M['keep_rows']:,} 行", 14, INK, True)
    txt(s, Inches(7.1), Inches(3.0), Inches(5.2), Inches(0.5),
        f"总净偏差金额  {w_money(M['keep_n'])}", 14, GREEN, True)
    txt(s, Inches(7.1), Inches(3.6), Inches(5.2), Inches(0.5),
        f"替代料净残差仅 {w_money(M['alt_n'])}（正负对冲后）", 12, GREY)

    rect(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.8), LIGHT)
    rect(s, Inches(0.8), Inches(4.7), Inches(0.14), Inches(1.8), AMBER)
    txt(s, Inches(1.1), Inches(4.85), Inches(11.2), Inches(0.4),
        "结论", 16, AMBER, True)
    txt(s, Inches(1.1), Inches(5.3), Inches(11.2), Inches(1.1),
        "替代料本质是授权替换、正负对冲（正向 +%.1f万、负向 −%.1f万），毛偏差制造了巨大数字噪音，\n"
        "但真实残差极小。因此原「饮料厂 +93.5万超耗」是替代料对冲撑出的假象，剔除后两家均净省料。"
        % (M['alt_b_pos']/10000, abs(M['alt_b_neg'])/10000),
        13, INK)


def build_factory_cmp(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 5, "工厂维度对比", "剔除前后净偏差金额变化")
    bf_all = M["by_factory_all"]
    bf = M["by_factory"]
    rows = []
    for f in bf.index:
        fb = bf_all.loc[f, "净偏差金额"]
        v = bf.loc[f, "净偏差金额"]
        rows.append((f, fb, v))
    # 画双柱
    base_y = Inches(5.4)
    maxv = max(abs(r[1]) for r in rows)
    maxv = max(maxv, 1)
    bw = Inches(1.6)
    for i, (f, fb, v) in enumerate(rows):
        x = Inches(1.6) + i * Inches(3.2)
        # 全表柱
        h1 = Inches(abs(fb) / maxv * 2.8)
        col1 = RED if fb > 0 else GREEN
        if fb >= 0:
            rect(s, x, base_y - h1, bw, h1, col1)
        else:
            rect(s, x, base_y, bw, h1, col1)
        txt(s, x, base_y - h1 - Inches(0.35), bw, Inches(0.3),
            w_money(fb), 12, col1, True, PP_ALIGN.CENTER)
        # 剔除后柱
        x2 = x + bw + Inches(0.3)
        h2 = Inches(abs(v) / maxv * 2.8)
        col2 = RED if v > 0 else GREEN
        if v >= 0:
            rect(s, x2, base_y - h2, bw, h2, col2)
        else:
            rect(s, x2, base_y, bw, h2, col2)
        txt(s, x2, base_y - h2 - Inches(0.35), bw, Inches(0.3),
            w_money(v), 12, col2, True, PP_ALIGN.CENTER)
        txt(s, x, base_y + Inches(0.1), bw * 2 + Inches(0.3), Inches(0.4),
            f, 13, INK, True, PP_ALIGN.CENTER)

    txt(s, Inches(0.8), Inches(1.7), Inches(6), Inches(0.4),
        "（深/浅两柱：全表 vs 剔除替代料后）", 12, GREY)
    rect(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7), LIGHT)
    txt(s, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.6),
        "饮料厂：+93.5万 → −5.4万（由超耗翻为省料）；食品厂：−3.8万 → −8.9万（更省）。",
        13, INK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def build_mat_top10(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 6, "物料偏差 Top10", "按净偏差金额排序（剔除替代料后）")
    mt = M["mat_top10"]
    maxv = mt["净偏差金额"].abs().max()
    x0 = Inches(0.8)
    y0 = Inches(1.6)
    rh = Inches(0.48)
    bw_max = Inches(7.5)
    for i, (name, row) in enumerate(mt.iterrows()):
        y = y0 + i * rh
        v = row["净偏差金额"]
        col = RED if v > 0 else GREEN
        txt(s, x0, y, Inches(3.6), rh, f"{i+1:02d}. {name}", 11, INK,
            False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        bw = Inches(abs(v) / maxv * bw_max.inches) if maxv else Inches(0)
        bx = x0 + Inches(3.7)
        if v >= 0:
            rect(s, bx, y + Inches(0.08), bw, Inches(0.3), col)
        else:
            rect(s, bx + (bw_max - bw), y + Inches(0.08), bw, Inches(0.3), col)
        txt(s, bx + bw_max + Inches(0.1), y, Inches(1.6), rh,
            f"{v/10000:+,.1f}万", 11, col, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    # 洞察卡
    rect(s, Inches(11.6), Inches(1.6), Inches(1.0), Inches(0.0), WHITE)  # 占位
    rect(s, Inches(11.5), Inches(1.6), Inches(1.3), Inches(4.8), CARD_BG)
    txt(s, Inches(11.6), Inches(1.7), Inches(1.2), Inches(4.6),
        "红色=超耗\n绿色=省料\n\n第一名\n30.2g\n金黄胚1810\n超耗最突出\n（食品厂）",
        11, INK, False, PP_ALIGN.CENTER, MSO_ANCHOR.TOP)


def build_workshop(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 7, "车间偏差详情", "各工厂车间净偏差金额 Top")
    ws = M["by_workshop"]
    cols = [("食品厂", Inches(0.8)), ("饮料厂", Inches(7.0))]
    for f, x0 in cols:
        sub = ws.xs(f, level=0) if f in ws.index.get_level_values(0) else None
        if sub is None or len(sub) == 0:
            continue
        sub = sub.sort_values("净偏差金额", ascending=False).head(6)
        maxv = sub["净偏差金额"].abs().max()
        txt(s, x0, Inches(1.55), Inches(5.5), Inches(0.4), f"{f}", 16, TEAL, True)
        y0 = Inches(2.1)
        rh = Inches(0.62)
        bw_max = Inches(3.6)
        for i, (name, row) in enumerate(sub.iterrows()):
            y = y0 + i * rh
            v = row["净偏差金额"]
            col = RED if v > 0 else GREEN
            txt(s, x0, y, Inches(2.4), rh, name, 11, INK, False,
                PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
            bw = Inches(abs(v) / maxv * bw_max.inches) if maxv else Inches(0)
            bx = x0 + Inches(2.5)
            if v >= 0:
                rect(s, bx, y + Inches(0.18), bw, Inches(0.26), col)
            else:
                rect(s, bx + (bw_max - bw), y + Inches(0.18), bw, Inches(0.26), col)
            txt(s, bx + bw_max + Inches(0.05), y, Inches(1.4), rh,
                f"{v/10000:+,.1f}万", 11, col, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7), LIGHT)
    txt(s, Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.6),
        "食品厂超耗集中在棕榈油/麦芽糖浆类物料车间；饮料厂各车间剔除替代料后均为省料。",
        13, INK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def build_mat_type(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 8, "物料类型分布", "各物料类型净偏差金额（剔除替代料后）")
    mt = M["by_mat_type"]
    maxv = mt["净偏差金额"].abs().max()
    x0 = Inches(1.0)
    y0 = Inches(1.8)
    rh = Inches(0.7)
    bw_max = Inches(8.0)
    for i, (name, row) in enumerate(mt.iterrows()):
        y = y0 + i * rh
        v = row["净偏差金额"]
        col = RED if v > 0 else GREEN
        txt(s, x0, y, Inches(2.6), rh, name, 13, INK, False,
            PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        bw = Inches(abs(v) / maxv * bw_max.inches) if maxv else Inches(0)
        bx = x0 + Inches(2.7)
        if v >= 0:
            rect(s, bx, y + Inches(0.22), bw, Inches(0.3), col)
        else:
            rect(s, bx + (bw_max - bw), y + Inches(0.22), bw, Inches(0.3), col)
        txt(s, bx + bw_max + Inches(0.1), y, Inches(2.0), rh,
            f"{v/10000:+,.1f}万", 12, col, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0), LIGHT)
    rect(s, Inches(0.8), Inches(5.9), Inches(0.14), Inches(1.0), AMBER)
    txt(s, Inches(1.1), Inches(6.0), Inches(11.2), Inches(0.9),
        "所有物料类型均为净省料（绿色），无整体超耗类型。其余类型（如「其他」）占比极小，\n"
        "提示主数据维护（BOM/定额）缺失才是高偏差率的主因，而非真实物料浪费。",
        13, INK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def build_cause(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 9, "偏差原因分析", "高偏差率（≥5%）记录归因")
    total_hi = M["hi_cnt"]
    sys_cnt = M["hi_sys_cnt"]
    other = total_hi - sys_cnt
    pct_sys = sys_cnt / total_hi * 100 if total_hi else 0
    pct_other = 100 - pct_sys
    # 左栏 主数据型
    rect(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.2), CARD_BG)
    rect(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(0.5), TEAL)
    txt(s, Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.5),
        f"主数据型  {pct_sys:.0f}%", 18, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.0), Inches(2.4), Inches(5.2), Inches(3.3),
        f"• 系统无额定 / 无定额 导致 {sys_cnt} 条\n"
        f"• 占高偏差率记录的 {pct_sys:.1f}%\n"
        f"• 本质：BOM/定额主数据未维护\n"
        f"• 并非真实超耗，而是计量无基准\n"
        f"• 建议：补全定额主数据即可大幅消除",
        14, INK)
    # 右栏 其他
    rect(s, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.2), RGBColor(0xFB,0xEF,0xE3))
    rect(s, Inches(6.9), Inches(1.7), Inches(5.6), Inches(0.5), AMBER)
    txt(s, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.5),
        f"耗用/订单型  {pct_other:.0f}%", 18, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.1), Inches(2.4), Inches(5.2), Inches(3.3),
        f"• 真实耗用偏差 / 订单差异等 {other} 条\n"
        f"• 占高偏差率记录的 {pct_other:.1f}%\n"
        f"• 含工艺损耗、计量误差、替代料残留\n"
        f"• 需结合车间现场核查\n"
        f"• 建议：区分真实超耗与数据缺失",
        14, INK)


def build_warn(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 10, "预警核查结果", "无备注 / 异常 / 中间地带 三类预警")
    cards = [
        ("无备注预警", M["warn_note"], "该工厂无≥5万元无备注记录" if M["warn_note"] == 0 else f"共 {M['warn_note']} 条需补备注", TEAL),
        ("异常预警", M["warn_abn"], f"共 {M['warn_abn']} 条偏差异常记录" if M["warn_abn"] else "无异常预警记录", AMBER),
        ("中间地带明细", M["warn_mid"], f"共 {M['warn_mid']} 条待复核" if M["warn_mid"] else "无中间地带记录", RGBColor(0x3A,0xA8,0xB0)),
    ]
    cw = Inches(3.7)
    gap = Inches(0.3)
    x0 = Inches(0.8)
    y0 = Inches(1.8)
    for i, (name, cnt, desc, col) in enumerate(cards):
        x = x0 + i * (cw + gap)
        rect(s, x, y0, cw, Inches(2.6), CARD_BG)
        rect(s, x, y0, cw, Inches(0.14), col)
        txt(s, x, y0 + Inches(0.35), cw, Inches(0.5), name, 16, INK, True,
            PP_ALIGN.CENTER)
        txt(s, x, y0 + Inches(1.0), cw, Inches(0.8), f"{cnt}", 34, col, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x, y0 + Inches(1.95), cw, Inches(0.5), desc, 12, GREY, False,
            PP_ALIGN.CENTER)
    rect(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.6), LIGHT)
    rect(s, Inches(0.8), Inches(4.9), Inches(0.14), Inches(1.6), AMBER)
    txt(s, Inches(1.1), Inches(5.05), Inches(11.2), Inches(0.4),
        "核查结论", 16, AMBER, True)
    txt(s, Inches(1.1), Inches(5.5), Inches(11.2), Inches(0.9),
        f"无备注预警 {M['warn_note']} 条、异常预警 {M['warn_abn']} 条、中间地带 {M['warn_mid']} 条，"
        "均存在真实需跟进记录（非「无记录」）。建议优先补全主数据缺失项以消除大部分异常。",
        13, INK)


def build_summary(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    bf = M["by_factory"]
    food_net = bf[bf.index.str.contains("食品")].iloc[0]["净偏差金额"] if bf.index.str.contains("食品").any() else bf.iloc[0]["净偏差金额"]
    bev_net = bf[bf.index.str.contains("饮料")].iloc[0]["净偏差金额"] if bf.index.str.contains("饮料").any() else bf.iloc[-1]["净偏差金额"]
    header(s, 11, "总结与改进建议", "分工厂针对性行动项")
    # 食品厂
    rect(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(4.8), CARD_BG)
    rect(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(0.5), TEAL)
    txt(s, Inches(1.0), Inches(1.7), Inches(5.3), Inches(0.5),
        f"食品厂（剔除后 {w_money(food_net)}）", 15, WHITE, True,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.0), Inches(2.4), Inches(5.3), Inches(4.0),
        "• 局部正向超耗物料：二次精炼24°棕榈油、\n"
        "  二次精炼10°棕榈油、麦芽糖浆80° 等\n"
        "• 高偏差率主因：系统无额定/无定额（主数据）\n"
        "• 建议：\n"
        "  1. 优先补全 BOM 与定额主数据\n"
        "  2. 核查棕榈油/糖浆类物料真实耗用\n"
        "  3. 规范替代料登记，避免噪音混入",
        13, INK)
    # 饮料厂
    rect(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.8), RGBColor(0xFB,0xEF,0xE3))
    rect(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(0.5), AMBER)
    txt(s, Inches(7.0), Inches(1.7), Inches(5.3), Inches(0.5),
        f"饮料厂（剔除后 {w_money(bev_net)}）", 15, WHITE, True,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.0), Inches(2.4), Inches(5.3), Inches(4.0),
        "• 原 +93.5万「超耗」实为替代料对冲噪音\n"
        "• 剔除后 −5.4万，整体净省料\n"
        "• 建议：\n"
        "  1. 保留替代料成对登记，单独核算\n"
        "  2. 不再将替代料计入偏差总额\n"
        "  3. 关注少数真实超耗产线计量准确性",
        13, INK)


# ===================== 主流程 =====================
def main():
    M, d, keep = load_metrics(INPUT_XLSX)
    # 补充替代料正负拆分（用于说明页）
    alt = d[d["是否替代料"] == "是"]
    M["alt_b_pos"] = alt.loc[alt["偏差金额"] > 0, "偏差金额"].sum()
    M["alt_b_neg"] = alt.loc[alt["偏差金额"] < 0, "偏差金额"].sum()

    prs = Presentation()
    prs.slide_width = PW
    prs.slide_height = PH

    build_cover(prs, M)
    build_toc(prs)
    build_metrics(prs, M)
    build_alt_explain(prs, M)
    build_factory_cmp(prs, M)
    build_mat_top10(prs, M)
    build_workshop(prs, M)
    build_mat_type(prs, M)
    build_cause(prs, M)
    build_warn(prs, M)
    build_summary(prs, M)

    prs.save(OUTPUT_PPTX)
    print(f"已生成：{OUTPUT_PPTX}")
    print(f"  总记录 {M['total_rows']:,} / 替代料 {M['alt_rows']} / 有效 {M['keep_rows']:,}")
    print(f"  剔除后总净偏差：{M['keep_n']/10000:+,}万")


if __name__ == "__main__":
    main()

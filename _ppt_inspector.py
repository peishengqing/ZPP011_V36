import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.oxml.ns import qn

prs = Presentation(r'C:\Users\Administrator\Desktop\ZPP011偏差分析_专业版.pptx')

print(f'Total slides: {len(prs.slides)}')
print(f'Slide size: {prs.slide_width} x {prs.slide_height}')
bg_color = None
for slide in prs.slides:
    bg_ph = slide.background.fill
    if bg_ph.type is not None:
        try:
            bg_color = str(bg_ph.fore_color.rgb)
        except:
            bg_color = str(bg_ph.type)
    break
print(f'Background type: {bg_color}')
print()

for i, slide in enumerate(prs.slides):
    print(f'=== Slide {i+1} (layout: {slide.slide_layout.name}) ===')
    for shape in slide.shapes:
        print(f'  [{shape.shape_type.name if hasattr(shape.shape_type, "name") else shape.shape_type}] name={shape.name}')
        print(f'    pos=({shape.left},{shape.top}) size=({shape.width},{shape.height})')
        
        # Check fill/background color
        try:
            fill = shape.fill
            if fill.type is not None:
                print(f'    fill_type={fill.type}, color={fill.fore_color.rgb if fill.fore_color and fill.fore_color.type else "?"}')
        except:
            pass
            
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if not txt:
                    continue
                fi = ''
                if para.runs:
                    r = para.runs[0]
                    rgb = '?'
                    if r.font.color and r.font.color.type:
                        rgb = str(r.font.color.rgb)
                    fi = f'font={r.font.name}, size={r.font.size}, bold={r.font.bold}, color={rgb}, align={para.alignment}'
                print(f'    Text: {txt[:60]} | {fi}')
        
        if shape.has_table:
            tbl = shape.table
            print(f'    Table: {len(tbl.rows)}x{len(tbl.columns)}')
            for ri in range(min(2, len(tbl.rows))):
                cells = [tbl.cell(ri, ci).text.strip()[:25] for ci in range(len(tbl.columns))]
                print(f'      Row{ri}: {cells}')
        
        if shape.shape_type.name == 'PICTURE':
            print(f'    [IMAGE]')
    print()

# Extract theme colors from slide layouts / slide masters
print('=== Theme / Master ===')
master = prs.slide_masters[0] if prs.slide_masters else None
if master:
    for layout in master.slide_layouts:
        pass

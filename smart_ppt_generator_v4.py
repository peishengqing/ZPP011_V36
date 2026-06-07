# -*- coding: utf-8 -*-
"""
ZPP011 智能PPT生成器 v4 (10维高阶思维 + 不限页数)
10维分析框架：
  1. 执行摘要    2. 数据质量评估  3. 整体偏差概览
  4. 食品/饮料对比  5. 车间维度分析  6. 物料分类分析
  7. 偏差原因根因  8. 趋势预测      9. 异常预警
 10. 改进建议 + 动态附录（按数据量自动扩展页数）
"""
import os, re, warnings, math
from datetime import datetime
from io import BytesIO
from typing import List, Dict, Tuple, Optional, Any

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData

warnings.filterwarnings("ignore", category=UserWarning)
plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

# ============================================================
# 全局配置
# ============================================================
COMPANY = "云南达利生产基地"
TITLE   = "ZPP011 生产偏差深度分析报告"
SUBTITLE = "数据驱动决策 · 精准管控成本"

COLOR = {
    'primary': RGBColor(41, 128, 185),
    'secondary': RGBColor(142, 68, 173),
    'accent':   RGBColor(230, 126, 34),
    'success':  RGBColor(39, 174, 96),
    'danger':   RGBColor(231, 76, 60),
    'warning':  RGBColor(241, 196, 15),
    'dark':     RGBColor(44, 62, 80),
    'light':    RGBColor(236, 240, 241),
    'white':    RGBColor(255, 255, 255),
    'text':     RGBColor(52, 73, 94),
    'text_l':   RGBColor(149, 165, 166),
    'food':     RGBColor(231, 76, 60),
    'bev':      RGBColor(41, 128, 185),
}
MPL = {'pos': '#e74c3c', 'neg': '#27ae60', 'pri': '#2980b9', 'acc': '#e67e22'}

ROWS_PER_SLIDE = 18
ITEMS_PER_SLIDE = 8

# ============================================================
# 工具函数
# ============================================================
def _log(msg, log_cb=None):
    if log_cb:
        log_cb(msg)
    else:
        print(msg)

def _find_col(df, names):
    for n in names:
        if n in df.columns:
            return n
    return None

def _fmt(n, p=0):
    try:
        f = float(n)
        if abs(f) >= 1e8:
            return '%.*f亿' % (p, f/1e8)
        if abs(f) >= 1e4:
            return '%.*f万' % (p, f/1e4)
        return ('%,.*f' % (p, f)).replace(' ', '')
    except:
        return str(n)

def _safe(v, d=None):
    try:
        return v if pd.notna(v) else d
    except:
        return d

# ============================================================
# PPT 设计系统
# ============================================================
class DS:
    """Design System - 统一管理所有样式"""
    def __init__(s, prs):
        s.prs = prs
        s.W   = prs.slide_width
        s.H   = prs.slide_height

    def grad(s, slide, c1, c2, ang=45):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, s.W, s.H)
        f  = sh.fill
        f.gradient()
        f.gradient_angle = ang
        f.gradient_stops[0].color.rgb = c1
        f.gradient_stops[1].color.rgb = c2
        sh.line.fill.background()
        return sh

    def card(s, slide, L, T, W, H, bg=COLOR['white']):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, L, T, W, H)
        sh.fill.solid()
        sh.fill.fore_color.rgb = bg
        sh.line.color.rgb = COLOR['light']
        sh.line.width = Pt(1)
        try:
            sh.adjustments[0] = 0.08
        except: pass
        return sh

    def title_bar(s, slide, text, sub='', color=COLOR['primary']):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, s.W, Pt(68))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Pt(8), Pt(68))
        acc.fill.solid()
        acc.fill.fore_color.rgb = COLOR['accent']
        acc.line.fill.background()
        tx  = slide.shapes.add_textbox(Pt(30), Pt(10), s.W-Pt(60), Pt(38))
        tf  = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR['white']
        p.font.name = 'Microsoft YaHei'
        if sub:
            tx2 = slide.shapes.add_textbox(Pt(30), Pt(46), s.W-Pt(60), Pt(18))
            tf2 = tx2.text_frame
            p2  = tf2.paragraphs[0]
            p2.text = sub
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(200,220,240)
            p2.font.name = 'Microsoft YaHei'

    def kpi(s, slide, L, T, val, lbl, color=COLOR['primary']):
        txv = slide.shapes.add_textbox(L, T, Pt(160), Pt(55))
        tfv = txv.text_frame
        tfv.vertical_anchor = MSO_ANCHOR.MIDDLE
        p   = tfv.paragraphs[0]
        p.text = str(val)
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
        txl = slide.shapes.add_textbox(L, T+Pt(58), Pt(160), Pt(20))
        tfl = txl.text_frame
        pl  = tfl.paragraphs[0]
        pl.text = lbl
        pl.font.size = Pt(11)
        pl.font.color.rgb = COLOR['text_l']
        pl.font.name = 'Microsoft YaHei'
        pl.alignment = PP_ALIGN.CENTER

# ============================================================
# 图表生成器
# ============================================================
class CG:
    @staticmethod
    def bar_h(labels, vals, title, colors=None):
        fig, ax = plt.subplots(figsize=(9, max(4, len(labels)*0.45)))
        if colors is None:
            colors = [MPL['pos'] if v>=0 else MPL['neg'] for v in vals]
        bars = ax.barh(range(len(labels)), vals, color=colors, alpha=0.8)
        ax.set_yticks(range(len(labels)))
        ax.set_yticklabels(labels, fontsize=9)
        ax.axvline(0, color='k', linewidth=0.8)
        ax.set_title(title, fontsize=13, fontweight='bold')
        ax.grid(axis='x', linestyle='--', alpha=0.3)
        for bar, v in zip(bars, vals):
            ax.text(v, bar.get_y()+bar.get_height()/2,
                     '%.1f万'% (v/1e4) if abs(v)>1e4 else '%.1f'%v,
                     va='center', fontsize=8)
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        plt.close()
        buf.seek(0)
        return buf

    @staticmethod
    def pie(sizes, labels, title):
        fig, ax = plt.subplots(figsize=(6,5))
        ax.pie(sizes, labels=labels, autopct='%1.1f%%',
               startangle=90, colors=['#e74c3c','#27ae60','#2980b9','#f1c40f','#9b59b6'])
        ax.set_title(title, fontsize=13, fontweight='bold')
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        plt.close()
        buf.seek(0)
        return buf

    @staticmethod
    def bar_v(cats, vals1, vals2, labels, title):
        fig, ax = plt.subplots(figsize=(10,5))
        x = range(len(cats))
        w = 0.35
        ax.bar([i-w/2 for i in x], vals1, w, label=labels[0], color=MPL['pri'], alpha=0.8)
        ax.bar([i+w/2 for i in x], vals2, w, label=labels[1], color=MPL['acc'], alpha=0.8)
        ax.set_xticks(x)
        ax.set_xticklabels(cats, rotation=15, fontsize=8)
        ax.set_title(title, fontsize=13, fontweight='bold')
        ax.legend()
        ax.grid(axis='y', linestyle='--', alpha=0.3)
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        plt.close()
        buf.seek(0)
        return buf

# ============================================================
# 数据分析引擎（10维）
# ============================================================
class AX:
    """Analyzer v4 - 10维分析引擎"""
    def __init__(s, excel_path, log_cb=None):
        s.path   = excel_path
        s.cb     = log_cb
        s.d      = {}      # {key: DataFrame}
        s.has_fb = False
        s.food   = {}
        s.bev    = {}
        s.ins    = []      # 洞察列表
        s.alerts = []      # 预警列表
        s.top_pos = []      # 多耗TOP
        s.top_neg = []      # 少耗TOP
        s.workshops = []    # 车间列表
        s.materials = []     # 物料分类列表

    def load(s):
        _log("[AX] 加载数据...", s.cb)
        xl = pd.ExcelFile(s.path)
        sm = {
            'sum':  ['汇总统计','汇总','统计'],
            'det':  ['完整偏差明细','偏差明细','明细','Sheet1'],
            'alt':  ['替代料明细','替代料'],
            'cause':['偏差原因分析','原因分析','原因'],
            'trend':['趋势分析','Trend','Sheet10'],
            'abn':  ['异常预警','预警'],
            'note': ['无备注预警','无备注'],
            'mid':  ['中等偏差明细','中等偏差'],
            'full': ['完整数据','Sheet5'],
        }
        for key, names in sm.items():
            for n in names:
                if n in xl.sheet_names:
                    s.d[key] = pd.read_excel(s.path, sheet_name=n)
                    _log("[AX] 加载: %s" % n, s.cb)
                    break
        _log("[AX] 完成，共%d个数据集" % len(s.d), s.cb)
        return s

    def go(s):
        _log("[AX] 开始10维分析...", s.cb)
        s._s1_basic()
        s._s2_quality()
        s._s3_overview()
        s._s4_food_bev()
        s._s5_workshop()
        s._s6_material()
        s._s7_root_cause()
        s._s8_trend()
        s._s9_anomaly()
        s._s10_recommend()
        _log("[AX] 完成，洞察%d条" % len(s.ins), s.cb)
        return s
# ========== AX 分析方法实现 ==========
    def _s1_basic(s):
        """维度1：基础统计"""
        df = s.d.get('sum')
        if df is None: return
        c = {}
        for k,ns in [('total',['总条数','记录数']),('pos',['正偏差条数','正偏差']),
                     ('neg',['负偏差条数','负偏差']),('pos_amt',['正偏差金额(含税)','正偏差金额']),
                     ('neg_amt',['负偏差金额(含税)','负偏差金额'])]:
            col = _find_col(df, ns)
            c[k] = df[col].sum() if col else 0
        rc = _find_col(df, ['备注覆盖率','覆盖率'])
        c['note_rate'] = 0
        if rc:
            try:
                rates = df[rc].astype(str).str.replace('%','').astype(float)
                weights = df[_find_col(df,['总条数','记录数'])]
                c['note_rate'] = (rates*weights).sum()/weights.sum()/100
            except: pass
        s.stats = c
        _log('[AX] 维度1完成: %s'%str(c)[:80], s.cb)

    def _s2_quality(s):
        """维度2：数据质量评估"""
        df = s.d.get('det')
        if df is None: return
        total = len(df)
        # 空值检查
        null_cols = []
        for col in ['偏差率(%)','偏差金额(含税)','备注原因']:
            if col in df.columns:
                null_rate = df[col].isna().sum()/total*100
                if null_rate > 5:
                    null_cols.append('%s空值%.1f%%'%(col,null_rate))
        # 异常值检查
        if '偏差率(%)' in df.columns:
            extreme = (df['偏差率(%)'].astype(float)>100).sum() if df['偏差率(%)'].dtype!=object else 0
            if extreme>0:
                s.alerts.append({'type':'数据质量','desc':'偏差率>100%%的记录有%d条'%extreme,'sev':'medium'})
        s.quality = {'total':total,'null_cols':null_cols,'extreme':int(extreme) if 'extreme' in dir() else 0}
        _log('[AX] 维度2完成: 质量评分%.0f'%(100-max(len(null_cols)*10,0)), s.cb)

    def _s3_overview(s):
        """维度3：整体偏差概览"""
        df = s.d.get('det')
        if df is None: return
        # 正负偏差分布
        if '偏差率(%)' in df.columns:
            rates = pd.to_numeric(df['偏差率(%)'],errors='coerce')
            s.overview = {
                'pos_rate': (rates>0).sum(),
                'neg_rate': (rates<0).sum(),
                'avg_pos': rates[rates>0].mean() if (rates>0).any() else 0,
                'avg_neg': rates[rates<0].mean() if (rates<0).any() else 0,
            }
        # TOP物料
        if '物料编码' in df.columns:
            top = df.groupby('物料编码').size().sort_values(ascending=False).head(10)
            s.top_mat = [{'code':k,'count':v} for k,v in top.items()]
        _log('[AX] 维度3完成', s.cb)

    def _s4_food_bev(s):
        """维度4：食品/饮料对比"""
        df = s.d.get('sum')
        if df is None: return
        fc = _find_col(df, ['工厂名称','工厂'])
        if fc is None: return
        all_f = df[fc].dropna().astype(str)
        s.has_fb = any('食品'in x for x in all_f) and any('饮料'in x for x in all_f)
        if not s.has_fb: return
        cols = ['总条数','正偏差条数','负偏差条数','正偏差金额(含税)','负偏差金额(含税)']
        for label, kw in [('food','食品'),('bev','饮料')]:
            m = df[df[fc].astype(str).str.contains(kw,na=False)]
            st = {}
            for k,ns in [('total',['总条数','记录数']),('pos',['正偏差条数','正偏差']),
                         ('neg',['负偏差条数','负偏差']),('pos_amt',['正偏差金额(含税)','正偏差金额']),
                         ('neg_amt',['负偏差金额(含税)','负偏差金额'])]:
                c = _find_col(m, ns)
                st[k] = m[c].sum() if c else 0
            setattr(s, label, st)
        _log('[AX] 维度4完成: 食品vs饮料', s.cb)

    def _s5_workshop(s):
        """维度5：车间维度分析"""
        df = s.d.get('det')
        if df is None: return
        wc = _find_col(df, ['车间','工厂名称','工厂'])
        if wc is None: return
        grp = df.groupby(df[wc]).size().sort_values(ascending=False)
        s.workshops = [{'name':k,'count':int(v)} for k,v in grp.items()]
        _log('[AX] 维度5完成: %d个车间'%len(s.workshops), s.cb)

    def _s6_material(s):
        """维度6：物料分类分析"""
        df = s.d.get('det')
        if df is None: return
        mc = _find_col(df, ['物料分类','物料类型','分类'])
        if mc is None: return
        grp = df.groupby(df[mc]).size().sort_values(ascending=False)
        s.materials = [{'name':k,'count':int(v)} for k,v in grp.items()]
        _log('[AX] 维度6完成: %d个分类'%len(s.materials), s.cb)
    def _s7_root_cause(s):
        """维度7：根因分析"""
        df = s.d.get('cause')
        if df is None: return
        rc = _find_col(df, ['备注原因','原因','偏差原因'])
        if rc is None: return
        counts = df[rc].value_counts().head(10)
        s.root = [{'cause':str(k),'count':int(v),'pct':v/len(df)*100}
                  for k,v in counts.items()]
        _log('[AX] 维度7完成: %d个原因'%len(s.root), s.cb)

    def _s8_trend(s):
        """维度8：趋势预测"""
        df = s.d.get('trend')
        s.trend = None
        if df is None or df.empty: return
        s.trend = {'has':True,
                    'summary':'趋势数据可用，建议重点关注偏差率上升的物料'}
        _log('[AX] 维度8完成', s.cb)

    def _s9_anomaly(s):
        """维度9：异常预警"""
        s.alerts = []
        # 无备注预警
        nd = s.d.get('note')
        if nd is not None and not nd.empty:
            s.alerts.append({'type':'无备注预警','count':len(nd),'sev':'high',
                               'desc':'%d条高偏差无备注'%len(nd)})
        # 异常预警
        ab = s.d.get('abn')
        if ab is not None and not ab.empty:
            s.alerts.append({'type':'异常预警','count':len(ab),'sev':'medium',
                               'desc':'%d条异常记录'%len(ab)})
        _log('[AX] 维度9完成: %d个预警'%len(s.alerts), s.cb)

    def _s10_recommend(s):
        """维度10：改进建议（生成洞察）"""
        s.ins = []
        st = getattr(s, 'stats', {})
        if st.get('note_rate',0) < 0.5:
            s.ins.append('备注覆盖率仅%.1f%%，急需提升'%(st.get('note_rate',0)*100))
        if len(s.alerts) > 0:
            s.ins.append('发现%d类异常预警，需优先处理'%len(s.alerts))
        if hasattr(s, 'overview') and s.overview.get('pos_rate',0) > s.overview.get('neg_rate',0)*2:
            s.ins.append('正偏差（多耗）数量远超负偏差，关注浪费')
        _log('[AX] 维度10完成: %d条建议'%len(s.ins), s.cb)

# ============================================================
# PPT 生成器 v4（不限页数）
# ============================================================
class PG:
    """PPT Generator - 10维动态生成"""
    def __init__(s, log_cb=None):
        s.cb   = log_cb
        s.prs  = Presentation()
        s.prs.slide_width  = Inches(13.333)
        s.prs.slide_height = Inches(7.5)
        s.ds  = DS(s.prs)
        s.cg  = CG()
        s.ax  = None

    def gen(s, excel_path, output_path):
        _log("[v4] 开始: %s" % excel_path, s.cb)
        s.ax = AX(excel_path, s.cb).load().go()

        s._cov()
        s._toc()
        s._s1()
        s._s2()
        s._s3()
        if s.ax.has_fb:
            s._s4()
        s._s5()
        s._s6()
        s._s7()
        s._s8()
        s._s9()
        s._s10()
        s._appendix()

        s.prs.save(output_path)
        _log("[v4] 完成: %s (%d页)" % (output_path, len(s.prs.slides)), s.cb)
        return True

    # ========== 封面 ==========
    def _cov(s):
        sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
        s.ds.grad(sl, COLOR['primary'], COLOR['secondary'], 135)
        for x,y,sz in [(Pt(50),Pt(100),Pt(200)),
                         (Pt(800),Pt(50),Pt(150)),
                         (Pt(900),Pt(400),Pt(100))]:
            c = sl.shapes.add_shape(MSO_SHAPE.OVAL, x,y,sz,sz)
            c.fill.solid()
            c.fill.fore_color.rgb = COLOR['white']
            c.fill.fore_color.theme_color = 1
            c.line.fill.background()
        tx = sl.shapes.add_textbox(Pt(100),Pt(150),Pt(800),Pt(100))
        tf = tx.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p  = tf.paragraphs[0]; p.text = TITLE
        p.font.size = Pt(54); p.font.bold = True
        p.font.color.rgb = COLOR['white']; p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        tx2 = sl.shapes.add_textbox(Pt(100),Pt(270),Pt(800),Pt(50))
        tf2 = tx2.text_frame; p2 = tf2.paragraphs[0]
        p2.text = SUBTITLE; p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(200,220,240); p2.font.name = 'Microsoft YaHei'
        p2.alignment = PP_ALIGN.CENTER
        tx3 = sl.shapes.add_textbox(Pt(100),Pt(350),Pt(800),Pt(40))
        tf3 = tx3.text_frame; p3 = tf3.paragraphs[0]
        p3.text = COMPANY; p3.font.size = Pt(18)
        p3.font.color.rgb = COLOR['accent']; p3.font.name = 'Microsoft YaHei'
        p3.alignment = PP_ALIGN.CENTER
        d = datetime.now().strftime('%Y年%m月%d日')
        tx4 = sl.shapes.add_textbox(Pt(100),Pt(400),Pt(800),Pt(30))
        tf4 = tx4.text_frame; p4 = tf4.paragraphs[0]
        p4.text = d; p4.font.size = Pt(14)
        p4.font.color.rgb = RGBColor(180,200,220); p4.alignment = PP_ALIGN.CENTER
        _log("[v4] 封面完成", s.cb)
    # ============================================================
    # 目录页
    # ============================================================
    def _toc(s):
        sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
        bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, s.prs.slide_width, s.prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = COLOR['light']; bg.line.fill.background()
        s.ds.title_bar(sl, "目录", "Report Contents")
        items = [("01","执行摘要"),("02","数据质量评估"),("03","整体偏差概览")]
        if s.ax.has_fb: items.append(("04","食品/饮料对比"))
        items += [("05","车间维度分析"),("06","物料分类分析"),("07","根因诊断"),("08","趋势预测"),("09","异常预警"),("10","改进建议")]
        for i,(num,title) in enumerate(items):
            y = Pt(90) + Pt(55)*i
            t = sl.shapes.add_textbox(Pt(50), y, Pt(60), Pt(40))
            tf = t.text_frame; p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = COLOR['accent']
            t2 = sl.shapes.add_textbox(Pt(130), y, Pt(400), Pt(40))
            tf2 = t2.text_frame; p2 = tf2.paragraphs[0]; p2.text = title; p2.font.size = Pt(18); p2.font.bold = True; p2.font.color.rgb = COLOR['dark']
        _log("[v4] 目录完成", s.cb)

    # ============================================================
    # 维度1：执行摘要
    # ============================================================
    def _s1(s):
        sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
        s.ds.title_bar(sl, "01 执行摘要", "Executive Summary")
        st = getattr(s.ax, 'stats', {})
        insights = [
            "本期共分析 %s 条偏差记录" % _fmt(st.get('total',0)),
            "正偏差（多耗）%s 条，金额 ¥%s" % (st.get('pos',0), _fmt(st.get('pos_amt',0))),
            "负偏差（少耗）%s 条，金额 ¥%s" % (st.get('neg',0), _fmt(st.get('neg_amt',0))),
            "备注覆盖率 %.1f%%" % (st.get('note_rate',0)*100),
        ]
        for a in getattr(s.ax, 'alerts', []):
            insights.append("⚠ %s" % a.get('desc',''))
        y = Pt(90)
        for txt in insights:
            tx = sl.shapes.add_textbox(Pt(30), y, Pt(600), Pt(28))
            tf = tx.text_frame; p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(13); p.font.color.rgb = COLOR['text']; p.space_after = Pt(6)
            y += Pt(35)
        recs = ["立即补录无备注高偏差记录","重点核查异常预警物料","建立偏差根因分类标准","推动系统化防错机制"]
        y2 = Pt(90)
        for txt in recs:
            tx = sl.shapes.add_textbox(Pt(650), y2, Pt(540), Pt(28))
            tf = tx.text_frame; p = tf.paragraphs[0]; p.text = "▶ "+txt; p.font.size = Pt(12); p.font.color.rgb = COLOR['secondary']; p.space_after = Pt(6)
            y2 += Pt(35)
        _log("[v4] 维度1完成", s.cb)
    # ============================================================
    # 维度2：数据质量评估
    # ============================================================
    def _s2(s):
        sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
        s.ds.title_bar(sl, "02 数据质量评估", "Data Quality")
        q  = getattr(s.ax, 'quality', {})
        total = q.get('total', 0)
        nulls = q.get('null_cols', [])
        items = [
            "数据总行数：%d" % total,
            "空值列数：%d" % len(nulls),
        ] + ["⚠ " + x for x in nulls] + [
            "评分：%d/100" % max(100 - len(nulls)*10, 0),
            "建议：及时补全缺失的备注原因和金额数据",
        ]
        y = Pt(90)
        for txt in items:
            tx = sl.shapes.add_textbox(Pt(30), y, Pt(600), Pt(30))
            tf = tx.text_frame; p = tf.paragraphs[0]
            p.text = txt; p.font.size = Pt(13); p.font.color.rgb = COLOR['dark']; p.space_after = Pt(5)
            y += Pt(36)
        _log("[v4] 维度2完成", s.cb)

    # ============================================================
    # 维度3：整体偏差概览
    # ============================================================
    def _s3(s):
        sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
        s.ds.title_bar(sl, "03 整体偏差概览", "Overall Deviation Overview")
        ov = getattr(s.ax, 'overview', {})
        kpis = [
            ("总偏差数", _fmt(ov.get('pos_rate', 0) + ov.get('neg_rate', 0))),
            ("正偏差均", "%.1f%%" % ov.get('avg_pos', 0)),
            ("负偏差均", "%.1f%%" % ov.get('avg_neg', 0)),
        ]
        x = Pt(30)
        for lbl, val in kpis:
            s.ds.kpi(sl, x, Pt(100), val, lbl)
            x += Pt(180)
        # 饼图
        try:
            buf = CG.pie(
                [ov.get('pos_rate', 1), ov.get('neg_rate', 1)],
                ['正偏差','负偏差'],
                "正负偏差分布"
            )
            pic = sl.shapes.add_picture(buf, Pt(500), Pt(100), width=Pt(380))
        except: pass
        _log("[v4] 维度3完成", s.cb)
    # ============================================================
    # 维度5：车间维度分析（动态页数）
    # ============================================================
    def _s5(s):
        ws = getattr(s.ax, 'workshops', [])
        if not ws: return
        chunk = 6
        pages = math.ceil(len(ws)/chunk)
        for p in range(pages):
            sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
            s.ds.title_bar(sl, "05 车间维度分析（%d/%d）" % (p+1, pages), "Workshop Analysis")
            for i in range(chunk):
                idx = p*chunk + i
                if idx >= len(ws): break
                w = ws[idx]
                L = Pt(30) + (i % 3) * Pt(420)
                T = Pt(90) + (i // 3) * Pt(220)
                s.ds.card(sl, L, T, Pt(400), Pt(200))
                tx = sl.shapes.add_textbox(L+Pt(8), T+Pt(8), Pt(380), Pt(30))
                tf = tx.text_frame; tf.paragraphs[0].text = w['name']; tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.bold = True
                tx2 = sl.shapes.add_textbox(L+Pt(8), T+Pt(40), Pt(380), Pt(150))
                tf2 = tx2.text_frame; tf2.paragraphs[0].text = "偏差条数：%d\n点击查看明细" % w['count']; tf2.paragraphs[0].font.size = Pt(11)
            _log("[v4] 维度5页%d完成" % (p+1), s.cb)

    # ============================================================
    # 维度6：物料分类分析（动态页数）
    # ============================================================
    def _s6(s):
        ms = getattr(s.ax, 'materials', [])
        if not ms: return
        chunk = 6
        pages = math.ceil(len(ms)/chunk)
        for p in range(pages):
            sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
            s.ds.title_bar(sl, "06 物料分类分析（%d/%d）" % (p+1, pages), "Material Category Analysis")
            for i in range(chunk):
                idx = p*chunk + i
                if idx >= len(ms): break
                m = ms[idx]
                L = Pt(30) + (i % 3) * Pt(420)
                T = Pt(90) + (i // 3) * Pt(220)
                s.ds.card(sl, L, T, Pt(400), Pt(200))
                tx = sl.shapes.add_textbox(L+Pt(8), T+Pt(8), Pt(380), Pt(30))
                tf = tx.text_frame; tf.paragraphs[0].text = m['name']; tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.bold = True
                tx2 = sl.shapes.add_textbox(L+Pt(8), T+Pt(40), Pt(380), Pt(150))
                tf2 = tx2.text_frame; tf2.paragraphs[0].text = "偏差条数：%d" % m['count']; tf2.paragraphs[0].font.size = Pt(11)
            _log("[v4] 维度6页%d完成" % (p+1), s.cb)
    # ===========================================================
    # 维度7：根因诊断（动态页数）
    # ===========================================================
    def _s7(s):
        root = getattr(s.ax, 'root', [])
        if not root: return
        chunk = 8
        pages = math.ceil(len(root)/chunk)
        for p in range(pages):
            sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
            s.ds.title_bar(sl, "07 根因诊断（%d/%d）" % (p+1, pages), "Root Cause Analysis")
            for i in range(chunk):
                idx = p*chunk + i
                if idx >= len(root): break
                r = root[idx]
                L = Pt(30) + (i % 2) * Pt(620)
                T = Pt(90) + (i // 2) * Pt(100)
                s.ds.card(sl, L, T, Pt(580), Pt(85))
                tx = sl.shapes.add_textbox(L+Pt(8), T+Pt(6), Pt(560), Pt(28))
                tf = tx.text_frame; tf.paragraphs[0].text = "%d. %s" % (idx+1, r['cause'])
                tf.paragraphs[0].font.size = Pt(12); tf.paragraphs[0].font.bold = True
                tx2 = sl.shapes.add_textbox(L+Pt(8), T+Pt(34), Pt(560), Pt(22))
                tf2 = tx2.text_frame; tf2.paragraphs[0].text = "数量：%d  占比：%.1f%%" % (r['count'], r['pct'])
                tf2.paragraphs[0].font.size = Pt(10); tf2.paragraphs[0].font.color.rgb = COLOR['text_l']
            _log("[v4] 维度7页%d完成" % (p+1), s.cb)

    # ===========================================================
    # 维度8：趋势预测
    # ===========================================================
    def _s8(s):
        trend = getattr(s.ax, 'trend', None)
        if not trend: return
        sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
        s.ds.title_bar(sl, "08 趋势预测", "Trend Analysis")
        tx = sl.shapes.add_textbox(Pt(30), Pt(90), Pt(900), Pt(400))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = trend.get('summary', '趋势数据已加载，建议结合历史偏差率分析未来走势。')
        p.font.size = Pt(16); p.font.color.rgb = COLOR['dark']
        _log("[v4] 维度8完成", s.cb)
    # ===========================================================
    # 维度9：异常预警（动态页数）
    # ===========================================================
    def _s9(s):
        alerts = getattr(s.ax, 'alerts', [])
        if not alerts: return
        chunk = 5
        pages = math.ceil(len(alerts)/chunk)
        for p in range(pages):
            sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
            s.ds.title_bar(sl, "09 异常预警（%d/%d）" % (p+1, pages), "Anomaly Warning")
            for i in range(chunk):
                idx = p*chunk + i
                if idx >= len(alerts): break
                a = alerts[idx]
                tx = sl.shapes.add_textbox(Pt(30), Pt(90)+Pt(50)*i, Pt(900), Pt(40))
                tf = tx.text_frame
                p = tf.paragraphs[0]
                p.text = u"⚠ %s：%s" % (a.get('type',''), a.get('desc',''))
                p.font.size = Pt(12)
                p.font.color.rgb = COLOR['danger'] if a.get('sev')=='high' else COLOR['warning']

    # ===========================================================
    # 维度10：改进建议
    # ===========================================================
    def _s10(s):
        sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
        s.ds.title_bar(sl, "10 改进建议", "Recommendations")
        ins = getattr(s.ax, 'ins', [])
        recs = [
            "立即补录无备注高偏差记录",
            "重点核查异常预警物料",
            "建立偏差根因分类标准",
            "推动系统化防错机制",
        ]
        all_items = ins + recs[:max(0,4-len(ins))]
        for i, txt in enumerate(all_items):
            tx = sl.shapes.add_textbox(Pt(30), Pt(90)+Pt(45)*i, Pt(900), Pt(38))
            tf = tx.text_frame
            p = tf.paragraphs[0]
            p.text = u"▶ %s" % txt
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR['dark']
            p.space_after = Pt(6)
        _log("[v4] 维度10完成", s.cb)
    # ==========================================================
    # 维度4：食品/饮料对比
    # ==========================================================
    def _s4(s):
        if not getattr(s.ax, 'has_fb', False):
            return
        sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
        s.ds.title_bar(sl, "04 食品/饮料对比", "Food vs Beverage")
        food = getattr(s.ax, 'food', {})
        bev  = getattr(s.ax, 'bev', {})
        # 左侧：食品 KPI
        xL = Pt(30); y0 = Pt(90)
        for i, (lbl, val) in enumerate([
            ("食品·总条数", _fmt(food.get('total', 0))),
            ("食品·正偏差金额", "¥" + _fmt(food.get('pos_amt', 0))),
            ("食品·负偏差金额", "¥" + _fmt(food.get('neg_amt', 0))),
            ("食品·备注覆盖率", "%.1f%%" % (food.get('note_rate', 0)*100)),
        ]):
            y = y0 + Pt(95)*i
            s.ds.card(sl, xL, y, Pt(300), Pt(80))
            tx = sl.shapes.add_textbox(xL+Pt(8), y+Pt(6), Pt(280), Pt(25))
            tx.text_frame.paragraphs[0].text = lbl
            tx.text_frame.paragraphs[0].font.size = Pt(10)
            tx.text_frame.paragraphs[0].font.color.rgb = COLOR['text_l']
            tx2 = sl.shapes.add_textbox(xL+Pt(8), y+Pt(35), Pt(280), Pt(35))
            tx2.text_frame.paragraphs[0].text = str(val)
            tx2.text_frame.paragraphs[0].font.size = Pt(18)
            tx2.text_frame.paragraphs[0].font.bold = True
            tx2.text_frame.paragraphs[0].font.color.rgb = COLOR['food']
        # 右侧：饮料 KPI
        xR = Pt(350)
        for i, (lbl, val) in enumerate([
            ("饮料·总条数", _fmt(bev.get('total', 0))),
            ("饮料·正偏差金额", "¥" + _fmt(bev.get('pos_amt', 0))),
            ("饮料·负偏差金额", "¥" + _fmt(bev.get('neg_amt', 0))),
            ("饮料·备注覆盖率", "%.1f%%" % (bev.get('note_rate', 0)*100)),
        ]):
            y = y0 + Pt(95)*i
            s.ds.card(sl, xR, y, Pt(300), Pt(80))
            tx = sl.shapes.add_textbox(xR+Pt(8), y+Pt(6), Pt(280), Pt(25))
            tx.text_frame.paragraphs[0].text = lbl
            tx.text_frame.paragraphs[0].font.size = Pt(10)
            tx.text_frame.paragraphs[0].font.color.rgb = COLOR['text_l']
            tx2 = sl.shapes.add_textbox(xR+Pt(8), y+Pt(35), Pt(280), Pt(35))
            tx2.text_frame.paragraphs[0].text = str(val)
            tx2.text_frame.paragraphs[0].font.size = Pt(18)
            tx2.text_frame.paragraphs[0].font.bold = True
            tx2.text_frame.paragraphs[0].font.color.rgb = COLOR['bev']
        _log("[v4] 维度4完成", s.cb)

    # ===========================================================
    # 附录：偏差明细（动态页数，不限页）
    # ===========================================================
    def _appendix(s):
        df = s.ax.d.get('det')
        if df is None or df.empty:
            _log("[v4] 附录：无明细数据，跳过", s.cb)
            return
        total = len(df)
        chunk = ROWS_PER_SLIDE
        pages = math.ceil(total / chunk)
        _log("[v4] 附录：共%d行，生成%d页" % (total, pages), s.cb)
        for p in range(pages):
            sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
            s.ds.title_bar(sl, "附录%d：偏差明细（%d/%d）" % (p+1, p+1, pages), "Appendix：Deviation Details")
            start = p * chunk
            end = min(start + chunk, total)
            sub = df.iloc[start:end]
            # 只取前6列写入文本（表格太复杂，用文本框代替）
            cols = list(sub.columns)[:6]
            tx = sl.shapes.add_textbox(Pt(20), Pt(90), Pt(900), Pt(380))
            tf = tx.text_frame
            tf.word_wrap = True
            for idx, row in sub.iterrows():
                txt = " | ".join([str(row[c])[:15] for c in cols])
                p = tf.add_paragraph()
                p.text = txt
                p.font.size = Pt(8)
                p.font.color.rgb = COLOR['text']
                p.space_after = Pt(2)
        _log("[v4] 附录完成，共%d页" % pages, s.cb)

    # ===========================================================
    # 结尾页
    # ===========================================================
    def _closing(s):
        sl = s.prs.slides.add_slide(s.prs.slide_layouts[6])
        s.ds.grad(sl, COLOR['secondary'], COLOR['primary'], 135)
        tx = sl.shapes.add_textbox(Pt(100), Pt(200), Pt(800), Pt(100)))
        tf = tx.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "谢谢！"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = COLOR['white']
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        tx2 = sl.shapes.add_textbox(Pt(100), Pt(320), Pt(800), Pt(40)))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "数据驱动决策 · 精准管控成本"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 220, 240)
        p2.font.name = 'Microsoft YaHei'
        p2.alignment = PP_ALIGN.CENTER
        _log("[v4] 结尾页完成", s.cb)

# ===========================================================
# 入口函数
# ===========================================================
def generate_smart_ppt_v4(excel_path, output_path, log_cb=None):
    """v4 入口：10维不限页数"""
    try:
        pg = PG(log_cb)
        return pg.gen(excel_path, output_path)
    except Exception as e:
        _log("[v4] 失败: %s" % str(e), log_cb)
        import traceback
        traceback.print_exc()
        return False

if __name__ == '__main__':
    import sys
    if len(sys.argv) < 3:
        print("用法: python smart_ppt_generator_v4.py 输入.xlsx 输出.pptx")
    else:
        generate_smart_ppt_v4(sys.argv[1], sys.argv[2])

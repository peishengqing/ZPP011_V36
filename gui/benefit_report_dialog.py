# -*- coding: utf-8 -*-
"""
效益报告生成对话框（完整版）
生成8页PPT：封面、执行摘要、车间排行、物料排行、趋势分析、成本换算、改进建议、附录
"""

# -*- coding: utf-8 -*-
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import os
from datetime import datetime
import threading
import pandas as pd

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class BenefitReportDialog:
    def __init__(self, parent, audit_data, output_dir=None):
        self.parent = parent
        self.audit_data = audit_data
        self.output_dir = output_dir or os.path.expanduser("~/Desktop")
        self.window = None
        self._build_ui()

    def _build_ui(self):
        self.window = tk.Toplevel(self.parent.root)
        self.window.title("生成效益报告")
        self.window.geometry("400x220")
        self.window.transient(self.parent.root)
        self.window.grab_set()
        self.window.resizable(False, False)

        self.window.update_idletasks()
        x = self.parent.root.winfo_x() + (self.parent.root.winfo_width() - 400) // 2
        y = self.parent.root.winfo_y() + (self.parent.root.winfo_height() - 220) // 2
        self.window.geometry(f"+{x}+{y}")

        tk.Label(
            self.window, text="📊 生成效益报告（完整版）", font=("微软雅黑", 12, "bold")
        ).pack(pady=10)
        tk.Label(self.window, text="数据范围：当前分析结果", fg="gray").pack()
        tk.Label(self.window, text="输出格式：PowerPoint 8页报告", fg="gray").pack(pady=5)

        path_frame = tk.Frame(self.window)
        path_frame.pack(fill=tk.X, padx=20, pady=10)
        tk.Label(path_frame, text="保存到：").pack(side=tk.LEFT)
        self.save_path_var = tk.StringVar(value=self._default_path())
        entry = tk.Entry(path_frame, textvariable=self.save_path_var, width=30)
        entry.pack(side=tk.LEFT, padx=5)
        tk.Button(path_frame, text="浏览", command=self._browse_path).pack(side=tk.LEFT)

        btn_frame = tk.Frame(self.window)
        btn_frame.pack(pady=15)
        tk.Button(
            btn_frame,
            text="生成",
            command=self._generate,
            bg="#28a745",
            fg="white",
            width=10,
        ).pack(side=tk.LEFT, padx=5)
        tk.Button(btn_frame, text="取消", command=self.window.destroy, width=10).pack(
            side=tk.LEFT, padx=5
        )

    def _default_path(self):
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        return os.path.join(self.output_dir, f"效益报告_{ts}.pptx")

    def _browse_path(self):
        file_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            initialfile=os.path.basename(self.save_path_var.get()),
            filetypes=[("PowerPoint 文件", "*.pptx")],
            title="保存效益报告",
        )
        if file_path:
            self.save_path_var.set(file_path)

    def _generate(self):
        if not PPTX_AVAILABLE:
            messagebox.showerror(
                "错误", "缺少 python-pptx 库，请执行：pip install python-pptx"
            )
            return

        output_path = self.save_path_var.get()
        if not output_path.endswith(".pptx"):
            output_path += ".pptx"

        self.window.grab_release()
        self.window.destroy()

        def task():
            try:
                self.parent.log("📊 开始生成效益报告（8页完整版）...", "info")
                self._generate_ppt(output_path)
                self.parent.root.after(0, lambda: self._on_success(output_path))
            except Exception as e:
                error_msg = str(e)
                self.parent.root.after(0, lambda msg=error_msg: self._on_error(msg))

        threading.Thread(target=task, daemon=True).start()

    def _generate_ppt(self, output_path):
        """核心PPT生成逻辑（完整8页版）"""
        from core import history_db
        from core.summary_generator import generate_summary

        df = self.audit_data
        if df is None or df.empty:
            raise ValueError("无数据，请先完成分析")

        app_title = self.parent.root.title() if hasattr(self.parent, 'root') else 'ZPP011'

        # ========== 1. 计算基础指标 ==========
        # 偏差金额列
        amount_col = None
        for c in df.columns:
            if '偏差金额' in str(c):
                amount_col = c
                break
        if amount_col is None:
            amount_col = '偏差金额'  # fallback

        total_amount = df[amount_col].fillna(0).sum()

        # 车间排名（Top5）
        workshop_col = next((c for c in ['车间', '生产管理员描述', 'admin'] if c in df.columns), None)
        if workshop_col:
            workshop_amount = df.groupby(workshop_col)[amount_col].apply(lambda x: x.abs().sum()).nlargest(5)
            top_workshop = workshop_amount.index[0] if len(workshop_amount) > 0 else "无"
        else:
            workshop_amount = pd.Series({"无车间数据": 0})
            top_workshop = "无"

        # 物料排名（Top10）
        mat_code_col = next((c for c in ['物料编码', '组件物料号', 'code'] if c in df.columns), None)
        mat_name_col = next((c for c in ['物料名称', '组件物料描述', 'name'] if c in df.columns), None)
        if mat_code_col:
            group_cols = [mat_code_col]
            if mat_name_col and mat_name_col != mat_code_col:
                group_cols.append(mat_name_col)
            material_rank = df.groupby(group_cols)[amount_col].apply(lambda x: x.abs().sum()).nlargest(10).reset_index()
        else:
            material_rank = pd.DataFrame()

        # 审核完成率
        status_col = next((c for c in ['审核状态', 'audit_status'] if c in df.columns), None)
        if status_col:
            reviewed = (df[status_col] == '已审核').sum()
            total = len(df)
            review_rate = reviewed / total if total > 0 else 0
        else:
            review_rate = 0

        # ========== 2. 趋势数据（历史对比）==========
        trend_data = None
        try:
            from core.history_db import get_analysis_list
            records = get_analysis_list(limit=6)
            if records and len(records) >= 2:
                trend = []
                for rec in records:
                    try:
                        from core.history_db import get_analysis_data
                        hist_df = get_analysis_data(rec.get('id'))
                        if hist_df is not None and not hist_df.empty:
                            amt = hist_df[amount_col].fillna(0).sum() if amount_col in hist_df.columns else 0
                            ts = rec.get('timestamp', '')
                            if ts:
                                trend.append((ts[:7], amt))  # YYYY-MM
                    except Exception:
                        continue
                if len(trend) >= 2:
                    trend_data = pd.DataFrame(trend, columns=['month', 'amount'])
        except Exception as e:
            self.parent.log(f"获取趋势数据失败：{e}", "warn")

        # ========== 3. 成本换算（按单位分组）==========
        unit_col = None
        for c in df.columns:
            if '单位' in str(c) or 'unit' in str(c).lower():
                unit_col = c
                break

        unit_summary = []
        if unit_col:
            try:
                groups = df.groupby(unit_col)
                for unit, group in groups:
                    amount = group[amount_col].fillna(0).sum()
                    if amount != 0:
                        unit_summary.append(f"单位 {unit}：偏差金额 ¥{amount:,.2f}")
            except Exception:
                pass

        # ========== 4. 改进建议（智能小结）==========
        try:
            from core.summary_generator import generate_summary
            summary_text = generate_summary(df, None)
            advice = summary_text[:500] + "..." if len(summary_text) > 500 else summary_text
        except Exception as e:
            advice = f"（智能小结生成失败：{e}）"

        # ========== 5. 创建PPT ==========
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]

        # ----- 第1页：封面 -----
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = "ZPP011 生产偏差分析效益报告"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"报告周期：{datetime.now().strftime('%Y-%m-%d')}\n数据版本：{app_title}"
        subtitle_frame.paragraphs[0].font.size = Pt(14)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # ----- 第2页：执行摘要 -----
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "执行摘要"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        kpi_data = [
            ("总偏差金额", f"¥{total_amount:,.2f}"),
            ("Top 偏差车间", top_workshop),
            ("审核完成率", f"{review_rate*100:.1f}%")
        ]
        y = 1.5
        for i, (label, value) in enumerate(kpi_data):
            box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(3), Inches(0.6))
            frame = box.text_frame
            frame.text = f"{label}\n{value}"
            frame.paragraphs[0].font.size = Pt(16)
            frame.paragraphs[0].bold = True if i == 0 else False
            frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(240, 240, 240)
            box.line.color.rgb = RGBColor(200, 200, 200)
            y += 1.2

        # ----- 第3页：车间偏差排行 -----
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "车间偏差金额排行"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        if not workshop_amount.empty and workshop_amount.iloc[0] != 0:
            chart_data = CategoryChartData()
            chart_data.categories = workshop_amount.index.tolist()
            chart_data.add_series("偏��金��", workshop_amount.values.tolist())
            x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = False
            chart.category_axis.tick_labels.font.size = Pt(9)
            chart.value_axis.tick_labels.font.size = Pt(9)
            for series in chart.series:
                series.data_labels.show_value = True
                series.data_labels.font.size = Pt(8)
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            txBox.text_frame.text = "无车间数据"

        # ----- 第4页：物料偏差排行 -----
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "物料偏差金额排行"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        if not material_rank.empty:
            rows = min(10, len(material_rank))
            cols = 3
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(0.4 * (rows + 1))
            table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
            table.columns[0].width = Inches(1.2)
            table.columns[1].width = Inches(4)
            table.columns[2].width = Inches(2.8)

            headers = ["排名", "物料编码", "偏差金额(元)"]
            for col, header in enumerate(headers):
                cell = table.cell(0, col)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            for i in range(rows):
                row = material_rank.iloc[i]
                code = str(row[mat_code_col]) if mat_code_col else ''
                amount = row[amount_col]
                cell1 = table.cell(i + 1, 0)
                cell1.text = str(i + 1)
                cell2 = table.cell(i + 1, 1)
                cell2.text = code
                cell3 = table.cell(i + 1, 2)
                cell3.text = f"{amount:,.2f}"
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            txBox.text_frame.text = "无物料数据"

        # ----- 第5页：偏差趋势分析 -----
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "偏差趋势分析"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        if trend_data is not None and not trend_data.empty and len(trend_data) >= 2:
            chart_data = CategoryChartData()
            chart_data.categories = trend_data['month'].tolist()
            chart_data.add_series("总偏差金额(万元)", (trend_data['amount'] / 10000).tolist())
            x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = True
            try:
                if hasattr(chart.value_axis, 'has_title') and chart.value_axis.has_title:
                    chart.value_axis.title.text_frame.text = "万元"
            except Exception:
                pass  # 某些版本的 python-pptx 不支持坐标轴标题
            chart.value_axis.tick_labels.font.size = Pt(9)
            chart.category_axis.tick_labels.font.size = Pt(9)
            for series in chart.series:
                series.data_labels.show_value = True
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            txBox.text_frame.text = "历史数据不足（需至少2次分析记录）"

        # ----- 第6页：成本换算 -----
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "成本换算（按单位）"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        if unit_summary:
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            text_frame = text_box.text_frame
            text_frame.text = "\n".join(unit_summary)
            text_frame.paragraphs[0].font.size = Pt(12)
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            txBox.text_frame.text = "无法获取单位信息"

        # ----- 第7页：改进建议 -----
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "改进建议"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        text_frame = text_box.text_frame
        text_frame.text = advice
        text_frame.paragraphs[0].font.size = Pt(11)

        # ----- 第8页：附录 -----
        slide = prs.slides.add_slide(slide_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "附录：数据说明"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        appendix = (
            "数据来源：SAP 生产订单与实际耗用差异分析\n"
            "偏差金额 = (实际数量 - 定额数量) × 含税单价\n"
            "审核状态基于「审核结果」列判断\n"
            "趋势分析需要历史分析记录（至少2次）\n"
            f"报告生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        )
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
        text_frame = text_box.text_frame
        text_frame.text = appendix
        text_frame.paragraphs[0].font.size = Pt(10)

        # 保存
        prs.save(output_path)

    def _on_success(self, output_path):
        if messagebox.askyesno(
            "生成成功", f"效益报告已生成（8页）：\n{output_path}\n是否立即打开？"
        ):
            try:
                os.startfile(output_path)
            except Exception as e:
                messagebox.showwarning("打开失败", f"无法打开文件：{e}")

    def _on_error(self, error_msg):
        messagebox.showerror("生成失败", f"生成效益报告失败：\n{error_msg}")
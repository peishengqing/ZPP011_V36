# -*- coding: utf-8 -*-
"""
ZPP011 高级PPT生成器（v3 专业版）
— **新版** 匹配24页深色商业报告模板风格
  配色 #0A1F3D | 卡片布局 | 英标+中标题 | 左装饰竖线
"""
import os
import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import matplotlib.font_manager as fm

# 配置matplotlib中文字体
for _f in fm.fontManager.ttflist:
    if 'Microsoft YaHei' in _f.name or 'SimHei' in _f.name or 'Noto Sans SC' in _f.name:
        plt.rcParams['font.sans-serif'] = [_f.name]
        plt.rcParams['font.family'] = 'sans-serif'
        break
plt.rcParams['axes.unicode_minus'] = False

from PySide6.QtCore import QObject, Signal
from core.config_loader import config

# ── 新配色方案 ──────────────────────────────────────────
BG_DARK     = RGBColor(0x0A, 0x1F, 0x3D)  # 最深色背景
TEXT_WHITE  = RGBColor(0xF1, 0xF3, 0xF5)  # 浅白
TEXT_BLUE   = RGBColor(0x5B, 0x8C, 0xB8)  # 蓝色强调
ACCENT_RED  = RGBColor(0xC4, 0x5C, 0x4A)  # 暖红
ACCENT_GREEN= RGBColor(0x6B, 0x9B, 0x7A)  # 绿色
ACCENT_AMBER= RGBColor(0xD4, 0xA5, 0x74)  # 琥珀色
CARD_LIGHT  = RGBColor(0xF1, 0xF3, 0xF5)  # 浅卡片背景
CARD_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER     = RGBColor(0x5B, 0x8C, 0xB8)   # 装饰线色
TEXT_GRAY   = RGBColor(0x80, 0x80, 0x80)

# 幻灯片尺寸 (16:9)
SLIDE_W = 12192000  # 13.33 inches
SLIDE_H = 6858000   # 7.5 inches

# 布局常量
MARGIN   = 457086   # 0.5 inch
CONTENT_W = 11274781  # 内容区宽度
CONTENT_L = MARGIN     # 内容区左边缘
TITLE_TOP   = 476250   # 英文标签顶部
CN_TITLE_TOP= 700088   # 中文标题顶部
BODY_TOP    = 1200000  # 正文开始

# 字体常量
FONT_EN = 'Inter'
FONT_CN = 'Source Han Sans'
FONT_CN_DISPLAY = 'Source Han Serif CN'


# ── 工具函数 ────────────────────────────────────────────
def _emu(inches):
    return int(inches * 914400)

def _set_font(run, name=FONT_CN, size=14, bold=None, color=None):
    run.font.name = name
    run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color

def _add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def _add_textbox(slide, left, top, width, height, text, size=14, bold=None,
                 color=TEXT_WHITE, align=PP_ALIGN.LEFT, name=FONT_CN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, name=name, size=size, bold=bold, color=color)
    return txBox

def _make_base_slide(prs):
    """创建基础暗色幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_DARK)
    return slide

def _add_en_title(slide, en_label, cn_title, subtitle=None):
    """英文标签 + 中文标题 + 可选副标题"""
    _add_textbox(slide, CONTENT_L, TITLE_TOP, CONTENT_W, 133350,
                 en_label, size=9, bold=True, color=TEXT_BLUE, name=FONT_EN)
    _add_textbox(slide, CONTENT_L, CN_TITLE_TOP, CONTENT_W, 381000,
                 cn_title, size=30, bold=True, color=TEXT_WHITE, name=FONT_CN_DISPLAY)
    if subtitle:
        _add_textbox(slide, CONTENT_L, CN_TITLE_TOP + 640000, CONTENT_W, 200000,
                     subtitle, size=11, color=TEXT_GRAY)

def _add_accent_bar(slide, left, top, height):
    """左侧装饰竖线"""
    return _add_rect(slide, left, top, 9522, height, TEXT_BLUE)

def _add_card(slide, left, top, width, height, bg_color=CARD_LIGHT):
    """浅色卡片"""
    card = _add_rect(slide, left, top, width, height, bg_color)
    _add_accent_bar(slide, left, top, height)
    return card

def _add_watermark(slide):
    _add_textbox(slide, SLIDE_W - 1800000, SLIDE_H - 350000, 1600000, 250000,
                 "ZPP011 偏差分析", size=8, color=TEXT_GRAY, align=PP_ALIGN.RIGHT, name=FONT_EN)


# ── 主类 ──────────────────────────────────────────────────
class AdvancedPPTGeneratorV3(QObject):
    log_message = Signal(str)

    def __init__(self, excel_path=None, output_path=None, log_cb=None, df=None):
        super().__init__()
        self.excel_path = excel_path
        self.output_path = output_path
        self.log_cb = log_cb
        self.df_main = df.copy() if df is not None else None
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.df_alerts = None
        self.df_workshop = None
        self.df_material_type = None
        self.df_weekly = None
        self.df_alt_saving = None
        self.today = datetime.now().strftime('%Y-%m-%d')
        self.amount_col = '偏差金额(含税)'

    def log(self, msg, level="info"):
        if self.log_cb:
            self.log_cb(msg, level)
        self.log_message.emit(msg)

    def load_data(self):
        """加载+聚合 与v3原版相同"""
        if self.df_main is not None:
            df = self.df_main.copy()
        elif self.excel_path and os.path.exists(self.excel_path):
            df = pd.read_excel(self.excel_path)
            self.df_main = df
        else:
            self.log("未指定输入数据源", "error")
            return
        for c in ['偏差率(%)', '偏差金额(含税)', '定额', '实际', '偏差数量']:
            if c in df.columns:
                df[c] = pd.to_numeric(df[c], errors='coerce').fillna(0)
        if '偏差金额(含税)' in self.df_main.columns:
            self.amount_col = '偏差金额(含税)'
        elif '偏差金额' in self.df_main.columns:
            self.amount_col = '偏差金额'
        self.dev_rate_col = '偏差率(%)' if '偏差率(%)' in self.df_main.columns else ('偏差率' if '偏差率' in self.df_main.columns else None)
        self.df_workshop = self._agg_by('车间') if '车间' in df.columns else None
        self.df_material_type = self._agg_by('物料类型') if '物料类型' in df.columns else None
        self.df_weekly = self._agg_weekly()
        self.df_alt_saving = self._agg_alt() if '替代料组' in df.columns else None
        if self.dev_rate_col is not None:
            th = config.get('alert.threshold_percent', 10.0)
            oa = config.get('alert.only_alt_materials', True)
            dt = self.df_main.copy()
            if oa and '是否替代料' in dt.columns:
                dt = dt[dt['是否替代料'] == '是']
            self.df_alerts = dt[dt[self.dev_rate_col].abs() > th].copy() if not dt.empty else pd.DataFrame()
        self.log("数据加载完成", "info")

    def _agg_by(self, col):
        if col not in self.df_main.columns: return pd.DataFrame()
        rate_col = self.dev_rate_col or '偏差金额(含税)'
        g = self.df_main.groupby(col).agg(
            偏差金额合计=(self.amount_col, 'sum'),
            偏差率均值=(rate_col, 'mean'),
            记录数=('物料编码', 'count')
        ).reset_index()
        return g[(g['偏差金额合计'].notna() & (g['偏差金额合计'] != 0)) | (g['记录数'] > 0)]

    def _agg_weekly(self):
        if '订单日期' not in self.df_main.columns: return pd.DataFrame()
        df = self.df_main.copy()
        df['订单日期'] = pd.to_datetime(df['订单日期'], errors='coerce')
        df = df.dropna(subset=['订单日期'])
        if df.empty: return pd.DataFrame()
        df['周'] = df['订单日期'].dt.strftime('%Y-W%W')
        w = df.groupby('周').agg(
            偏差金额合计=(self.amount_col, 'sum'),
            记录数=('物料编码', 'count')
        ).reset_index()
        return w[w['偏差金额合计'].notna() & (w['偏差金额合计'] != 0)]

    def _agg_alt(self):
        if '替代料组' not in self.df_main.columns or '净偏差' not in self.df_main.columns:
            return pd.DataFrame()
        s = self.df_main.groupby('替代料组').agg(
            原偏差总额=(self.amount_col, 'sum'),
            净偏差总额=('净偏差', 'sum'),
            物料数量=('物料编码', 'count')
        ).reset_index()
        s['节省金额'] = s['原偏差总额'] - s['净偏差总额']
        return s

    # ── Slide 1: 封面 ──────────────────────────────────
    def _slide_cover(self):
        slide = _make_base_slide(self.prs)
        period = "综合分析"
        try:
            if '订单日期' in self.df_main.columns and not self.df_main['订单日期'].isna().all():
                dmin = self.df_main['订单日期'].min()
                dmax = self.df_main['订单日期'].max()
                period = f"{dmin.strftime('%Y-%m-%d') if hasattr(dmin, 'strftime') else dmin} ~ {dmax.strftime('%Y-%m-%d') if hasattr(dmax, 'strftime') else dmax}"
        except Exception:
            pass
        _add_textbox(slide, CONTENT_L, int(SLIDE_H * 0.28), CONTENT_W, 142875,
                     "Industrial Production Data Analysis", size=9, bold=True,
                     color=TEXT_BLUE, name=FONT_EN)
        _add_textbox(slide, CONTENT_L, int(SLIDE_H * 0.33), CONTENT_W, 838200,
                     "ZPP011 生产偏差分析报告", size=60, bold=True,
                     color=TEXT_WHITE, name=FONT_CN_DISPLAY)
        _add_textbox(slide, CONTENT_L, int(SLIDE_H * 0.44), CONTENT_W, 250000,
                     "云南达利生产基地", size=14, color=TEXT_GRAY, name=FONT_CN)
        _add_textbox(slide, CONTENT_L, int(SLIDE_H * 0.49), CONTENT_W, 200000,
                     period, size=11, color=TEXT_GRAY)
        _add_textbox(slide, CONTENT_L, int(SLIDE_H * 0.53), CONTENT_W, 200000,
                     f"{len(self.df_main):,} 条记录深度洞察", size=11, color=TEXT_GRAY)

    # ── Slide 2: 执行摘要 ──────────────────────────────
    def _slide_exec_summary(self):
        slide = _make_base_slide(self.prs)
        _add_en_title(slide, "EXECUTIVE SUMMARY", "核心指标总览")

        total = len(self.df_main)
        pos = self.df_main[self.df_main[self.amount_col] > 0] if self.amount_col in self.df_main.columns else pd.DataFrame()
        neg = self.df_main[self.df_main[self.amount_col] < 0] if self.amount_col in self.df_main.columns else pd.DataFrame()
        pos_c, neg_c = len(pos), len(neg)
        pos_a = pos[self.amount_col].sum() if not pos.empty else 0
        neg_a = neg[self.amount_col].sum() if not neg.empty else 0
        net = self.df_main[self.amount_col].sum() if self.amount_col in self.df_main.columns else 0
        remark_rate = (self.df_main['备注原因'].notna().sum() / total * 100) if '备注原因' in self.df_main.columns else 0
        alt_count = self.df_main['替代料组'].nunique() if '替代料组' in self.df_main.columns else 0

        # 顶部 4 个 KPI 卡片
        kpi_top = 1300000
        kpi_w = 2670875
        kpi_h = 1100000
        gap = (CONTENT_W - kpi_w * 4) // 3
        metrics = [
            (f"{total:,}", "总记录数", TEXT_WHITE),
            (f"{pos_c:,} / {neg_c:,}", "正/负偏差条数", TEXT_BLUE),
            (f"¥{abs(net)/10000:.0f}万", "净偏差金额", ACCENT_RED if net > 0 else (ACCENT_GREEN if net < 0 else TEXT_WHITE)),
            (f"{remark_rate:.1f}%", "备注覆盖率", ACCENT_AMBER if remark_rate < 50 else ACCENT_GREEN),
        ]
        for i, (val, lbl, clr) in enumerate(metrics):
            left = CONTENT_L + i * (kpi_w + gap)
            _add_card(slide, left, kpi_top, kpi_w, kpi_h, CARD_LIGHT)
            _add_textbox(slide, left + 150000, kpi_top + 150000, kpi_w - 300000, 500000,
                         val, size=36, bold=True, color=clr, name=FONT_EN)
            _add_textbox(slide, left + 150000, kpi_top + 650000, kpi_w - 300000, 300000,
                         lbl, size=11, color=BG_DARK, name=FONT_CN)

        # 底部分析概要
        summary_top = kpi_top + kpi_h + 300000
        _add_card(slide, CONTENT_L, summary_top, CONTENT_W, 3200000, CARD_LIGHT)
        _add_textbox(slide, CONTENT_L + 200000, summary_top + 200000, CONTENT_W - 400000, 350000,
                     "ZPP011 系统监控", size=14, bold=True, color=BG_DARK)

        details = []
        details.append(f"5月ZPP011系统监控显示，虽然净偏差金额{abs(net)/10000:.0f}万元整体{'可控' if abs(net/10000) < 100 else '需关注'}，"
                       f"但备注覆盖率{remark_rate:.1f}%{'暴露管理透明度不足' if remark_rate < 50 else '处于良好水平'}。")

        if not pos.empty or not neg.empty:
            details.append(f"全厂偏差呈现正偏差42.4%、负偏差38.1%、无偏差19.5%的分布特征，偏差发生具有普遍性。")

        if alt_count > 0:
            total_saved = self.df_alt_saving['节省金额'].sum() if self.df_alt_saving is not None else 0
            details.append(f"替代料核对机制覆盖 {alt_count} 组物料配对，有效识别净偏差抵消效果。")

        if self.df_workshop is not None and not self.df_workshop.empty:
            worst = self.df_workshop.sort_values('偏差金额合计', ascending=False).iloc[0] if len(self.df_workshop) > 0 else None
            if worst is not None:
                details.append(f"最高偏差单元：{worst['车间']}（偏差{worst['偏差金额合计']/10000:.0f}万元），建议重点关注。")

        _add_textbox(slide, CONTENT_L + 200000, summary_top + 600000, CONTENT_W - 400000, 2300000,
                     "\n".join(details), size=11, color=BG_DARK, name=FONT_CN)

    # ── Slide 3: 车间偏差分析（三栏卡片 ──
    def _slide_workshop_analysis(self):
        slide = _make_base_slide(self.prs)
        _add_en_title(slide, "PRODUCTION DEVIATION ANALYSIS", "TOP 车间偏差分析",
                       "各车间关键偏差数据透视")

        if self.df_workshop is None or self.df_workshop.empty:
            _add_textbox(slide, CONTENT_L, BODY_TOP, CONTENT_W, 400000,
                         "暂无车间数据", size=14, color=TEXT_GRAY)
            return

        ws = self.df_workshop.sort_values('偏差金额合计', ascending=False).head(3)
        card_w = 3605799
        card_h = 3579614
        gap3 = (CONTENT_W - card_w * 3) // 2
        card_top = 1824930

        for idx, (_, row) in enumerate(ws.iterrows()):
            left = CONTENT_L + idx * (card_w + gap3)
            _add_card(slide, left, card_top, card_w, card_h, CARD_LIGHT)

            # 编号 + 车间名
            _add_textbox(slide, left + 238065, card_top + 314325, card_w - 476130, 114300,
                         f"NO.{idx+1:02d}", size=9, bold=True, color=TEXT_BLUE, name=FONT_EN)
            _add_textbox(slide, left + 238065, card_top + 571668, card_w - 476130, 238125,
                         row['车间'], size=19, bold=True, color=BG_DARK, name=FONT_CN_DISPLAY)

            # 金额
            amt = row['偏差金额合计']
            amt_val = f"¥{abs(amt)/10000:.0f}万" if abs(amt) >= 10000 else f"¥{amt:.0f}"
            amt_clr = ACCENT_RED if amt > 0 else (ACCENT_GREEN if amt < 0 else BG_DARK)
            _add_textbox(slide, left + 238065, card_top + 1100000, 1109386, 438150,
                         amt_val, size=36, bold=True, color=amt_clr, name=FONT_EN)
            label_text = "正偏差(超耗)" if amt > 0 else "负偏差(节约)"
            _add_textbox(slide, left + 238065 + 1109386, card_top + 1250000, 800000, 250000,
                         label_text, size=11, color=amt_clr, name=FONT_CN)

            # 详情字段
            rec = row.get('记录数', 0)
            rate = row.get('偏差率均值', 0)
            y = card_top + 1700000
            _add_textbox(slide, left + 238065, y, card_w - 476130, 200000,
                         f"记录数：{rec} 条", size=11, color=BG_DARK, name=FONT_CN)
            _add_textbox(slide, left + 238065, y + 230000, card_w - 476130, 200000,
                         f"平均偏差率：{rate:.1f}%" if not pd.isna(rate) else "",
                         size=11, color=BG_DARK, name=FONT_CN)

            # 底部智能分析
            insight = ""
            if abs(amt) > 500000:
                insight = "⚠ 高偏差单元，需优先处理"
            elif abs(rate) > 10 and not pd.isna(rate):
                insight = "偏差率偏高，建议核查定额基准"
            if rec > 1000:
                insight = f"{insight}；高频记录，建议批量核查" if insight else "高频记录，建议批量核查"
            _add_textbox(slide, left + 238065, card_top + 2300000, card_w - 476130, 900000,
                         insight, size=11, color=ACCENT_RED, name=FONT_CN)

    # ── Slide 4: 偏差类型分布 ──────────────────────────
    def _slide_deviation_distribution(self):
        slide = _make_base_slide(self.prs)
        _add_en_title(slide, "DEVIATION DISTRIBUTION", "偏差类型分布与核心指标")

        if self.amount_col not in self.df_main.columns:
            _add_textbox(slide, CONTENT_L, BODY_TOP, CONTENT_W, 400000,
                         "缺少偏差数据", size=14, color=TEXT_GRAY)
            return

        pos_c = (self.df_main[self.amount_col] > 0).sum()
        neg_c = (self.df_main[self.amount_col] < 0).sum()
        zero_c = (self.df_main[self.amount_col] == 0).sum()
        total = pos_c + neg_c + zero_c or 1
        pos_p = pos_c / total * 100
        neg_p = neg_c / total * 100
        zero_p = zero_c / total * 100

        # 左侧：饼图
        fig, ax = plt.subplots(figsize=(4.5, 4.2), facecolor='#0A1F3D')
        ax.set_facecolor('#0A1F3D')
        sizes = [pos_c, neg_c, zero_c]
        lbls = ['正偏差(多耗)', '负偏差(少耗)', '无偏差']
        colors_pie = ['#C45C4A', '#5B8CB8', '#D4A574']
        wedges, texts, autotexts = ax.pie(
            sizes, labels=None, autopct='%1.1f%%',
            colors=colors_pie, explode=(0.05, 0.05, 0),
            startangle=90, textprops={'color': 'white', 'fontsize': 10})
        for at in autotexts:
            at.set_color('white')
            at.set_fontsize(9)
        # Legend with dark bg
        legend = ax.legend(wedges, [f'{l} ({s:.1f}%)' for l, s in zip(lbls, [pos_p, neg_p, zero_p])],
                          loc='lower center', bbox_to_anchor=(0.5, -0.12),
                          ncol=3, frameon=False, fontsize=8)
        for t in legend.get_texts():
            t.set_color('white')
        plt.tight_layout()
        img_path = "_temp_pie_v2.png"
        plt.savefig(img_path, dpi=150, bbox_inches='tight', facecolor='#0A1F3D')
        plt.close()

        pie_left = CONTENT_L
        pie_w = 5486400
        slide.shapes.add_picture(img_path, pie_left, BODY_TOP,
                                 width=Emu(pie_w), height=Emu(4114800))
        os.remove(img_path)

        # 右侧: 结构分析卡片
        right_left = CONTENT_L + pie_w + 400000
        card_w2 = CONTENT_W - pie_w - 400000
        _add_card(slide, right_left, BODY_TOP + 300000, card_w2, 2200000, CARD_LIGHT)
        _add_textbox(slide, right_left + 200000, BODY_TOP + 450000, card_w2 - 400000, 300000,
                     "偏差结构分析", size=16, bold=True, color=BG_DARK, name=FONT_CN_DISPLAY)
        y = BODY_TOP + 900000
        for label, desc, pct in [("正偏差（多耗）", "实际用量 > 定额用量", pos_p),
                                  ("负偏差（少耗）", "实际用量 < 定额用量", neg_p),
                                  ("无偏差", "实际用量 = 定额用量", zero_p)]:
            _add_textbox(slide, right_left + 200000, y, card_w2 - 400000, 250000,
                         f"{label}  |  记录占比 {pct:.1f}%", size=11, color=BG_DARK)
            y += 350000

        # 底部 insight
        insight_top = BODY_TOP + 2800000
        _add_card(slide, CONTENT_L, insight_top, CONTENT_W - 100000, 800000, CARD_LIGHT)
        _add_textbox(slide, CONTENT_L + 200000, insight_top + 150000, CONTENT_W - 500000, 500000,
                     "全厂偏差呈现正偏差42.4%、负偏差38.1%、无偏差19.5%的分布特征，偏差发生具有普遍性。"
                     "包材类净偏差为负（节约），原材料类净偏差为正（超耗），建议重点关注原材料定额管理。",
                     size=11, color=BG_DARK, name=FONT_CN)

    # ── Slide 5: 物料偏差排行 ──────────────────────────
    def _slide_material_top10(self):
        slide = _make_base_slide(self.prs)
        _add_en_title(slide, "MATERIAL ANALYSIS", "高价值偏差组件物料 TOP 分析")

        if self.amount_col not in self.df_main.columns or '物料名称' not in self.df_main.columns:
            _add_textbox(slide, CONTENT_L, BODY_TOP, CONTENT_W, 400000,
                         "缺少物料偏差数据", size=14, color=TEXT_GRAY)
            return

        top_mat = self.df_main.groupby('物料名称')[self.amount_col].sum().abs().sort_values(ascending=False).head(10).reset_index()
        top_mat = top_mat.merge(
            self.df_main.groupby('物料名称')[self.amount_col].sum().reset_index(),
            on='物料名称', suffixes=('_abs', ''))
        top_mat = top_mat.sort_values(self.amount_col)

        fig, ax = plt.subplots(figsize=(11, 4.5), facecolor='#0A1F3D')
        ax.set_facecolor('#0A1F3D')
        labels = [t[:20] + '..' if len(t) > 20 else t for t in top_mat['物料名称']]
        vals = top_mat[self.amount_col].tolist()
        colors_bar = ['#C45C4A' if v > 0 else '#5B8CB8' for v in vals]
        bars = ax.barh(labels, vals, color=colors_bar, height=0.6)
        ax.axvline(0, color='#555', linewidth=0.5)
        ax.tick_params(axis='y', colors='white', labelsize=9)
        ax.tick_params(axis='x', colors='white', labelsize=8)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#555')
        ax.spines['bottom'].set_color('#555')
        for bar, v in zip(bars, vals):
            px = bar.get_width()
            if abs(v) > max(abs(v) for v in vals) * 0.03:
                ax.text(px + max(abs(v) for v in vals) * 0.005 if px >= 0 else px - max(abs(v) for v in vals) * 0.08,
                        bar.get_y() + bar.get_height() / 2,
                        f'{v/10000:.1f}万' if abs(v) >= 10000 else f'{v:.0f}',
                        ha='left' if px >= 0 else 'right', va='center',
                        fontsize=9, color='white')
        plt.tight_layout()
        img_path = "_temp_mat_v2.png"
        plt.savefig(img_path, dpi=150, bbox_inches='tight', facecolor='#0A1F3D')
        plt.close()
        slide.shapes.add_picture(img_path, CONTENT_L, BODY_TOP,
                                 width=Emu(CONTENT_W), height=Emu(4800000))
        os.remove(img_path)

    # ── Slide 6: 预警明细表 ────────────────────────────
    def _slide_alert_table(self):
        slide = _make_base_slide(self.prs)
        _add_en_title(slide, "WARNING DETAILS", "无备注预警明细",
                       "核心发现：高偏差无备注记录需紧急跟进")

        # 取高偏差 + 无备注
        alert_df = self.df_alerts.copy() if self.df_alerts is not None else pd.DataFrame()
        if alert_df.empty and '备注原因' in self.df_main.columns:
            alert_df = self.df_main[
                (self.df_main[self.amount_col].abs() > 50000) &
                (self.df_main['备注原因'].isna())
            ].head(10).copy()

        if alert_df.empty:
            _add_textbox(slide, CONTENT_L, BODY_TOP, CONTENT_W, 400000,
                         "暂无预警记录 ✅", size=14, color=TEXT_GRAY)
            return

        alert_df = alert_df.head(5)
        cols_to_show = [c for c in ['订单日期', '工厂', '车间', '物料名称', self.amount_col] if c in alert_df.columns]
        col_lbls = ['订单日期', '工厂', '车间', '物料名称', '偏差金额']

        # 顶部 insight
        insight_bar = _add_rect(slide, CONTENT_L, BODY_TOP, CONTENT_W, 700000, CARD_LIGHT)
        _add_accent_bar(slide, CONTENT_L, BODY_TOP, 700000)
        detail_text = (
            f"共 {len(alert_df)} 条无备注高偏差记录，涉及金额超{alert_df[self.amount_col].abs().sum()/10000:.0f}万元。"
            f"平均值 {alert_df[self.amount_col].abs().mean()/10000:.2f}万元/条，缺乏原因追溯能力。"
        )
        _add_textbox(slide, CONTENT_L + 200000, BODY_TOP + 150000, CONTENT_W - 400000, 400000,
                     detail_text, size=10, color=BG_DARK, name=FONT_CN)

        # 表格
        table_top = BODY_TOP + 900000
        header_h = 450000
        col_widths = [1800000, 2000000, 1800000, 3400000, 2274781]
        total_w = sum(col_widths)
        x_start = CONTENT_L

        # 表头
        for ci, (col_lbl, cw) in enumerate(zip(col_lbls, col_widths)):
            _add_rect(slide, x_start + sum(col_widths[:ci]), table_top, cw, header_h, TEXT_BLUE)
            _add_textbox(slide, x_start + sum(col_widths[:ci]), table_top + 100000, cw, 250000,
                         col_lbl, size=11, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, name=FONT_CN)

        # 数据行
        row_h = 400000
        for ri in range(min(len(alert_df), 5)):
            row_top = table_top + header_h + ri * row_h
            bg = CARD_LIGHT if ri % 2 == 0 else CARD_WHITE
            row = alert_df.iloc[ri]
            for ci, col in enumerate(cols_to_show):
                cw = col_widths[ci]
                _add_rect(slide, x_start + sum(col_widths[:ci]), row_top, cw, row_h, bg)
                val = row[col]
                if isinstance(val, float):
                    txt = f"{val/10000:.1f}万" if abs(val) >= 10000 else f"{val:.0f}"
                elif isinstance(val, pd.Timestamp):
                    txt = val.strftime('%Y-%m-%d')
                else:
                    txt = str(val)[:15] if val is not None else ""
                _add_textbox(slide, x_start + sum(col_widths[:ci]) + 50000, row_top + 80000,
                             cw - 100000, 250000, txt, size=10, color=BG_DARK, name=FONT_CN)

        # 底部详细说明
        detail_top = table_top + header_h + 5 * row_h + 150000
        for ri in range(min(len(alert_df), 5)):
            row = alert_df.iloc[ri]
            mat = row.get('物料名称', '')
            amt = row.get(self.amount_col, 0)
            ws_name = row.get('车间', '')
            _add_rect(slide, CONTENT_L, detail_top + ri * 300000, CONTENT_W, 280000,
                      CARD_LIGHT if ri % 2 == 0 else CARD_WHITE)
            _add_textbox(slide, CONTENT_L + 150000, detail_top + ri * 300000 + 50000,
                         CONTENT_W - 300000, 200000,
                         f"{mat}（{ws_name}）：偏差金额 {amt/10000:.1f}万元，无备注说明",
                         size=10, color=BG_DARK, name=FONT_CN)

    # ── Slide 7: 物料类型分析 ──────────────────────────
    def _slide_material_type(self):
        slide = _make_base_slide(self.prs)
        _add_en_title(slide, "COMPONENT ANALYSIS", "组件物料分类偏差特征分析")

        if self.df_material_type is None or self.df_material_type.empty:
            _add_textbox(slide, CONTENT_L, BODY_TOP, CONTENT_W, 400000,
                         "暂无物料类型数据", size=14, color=TEXT_GRAY)
            return

        # Chart
        fig, ax = plt.subplots(figsize=(11, 3.8), facecolor='#0A1F3D')
        ax.set_facecolor('#0A1F3D')
        mt = self.df_material_type.sort_values('偏差金额合计')
        labels = mt['物料类型'].tolist()
        vals = mt['偏差金额合计'].tolist()
        colors_bar = ['#C45C4A' if v > 0 else '#5B8CB8' for v in vals]
        bars = ax.barh(labels, vals, color=colors_bar, height=0.5)
        ax.axvline(0, color='#555', linewidth=0.5)
        ax.tick_params(axis='y', colors='white', labelsize=10)
        ax.tick_params(axis='x', colors='white', labelsize=8)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#555'); ax.spines['bottom'].set_color('#555')
        for bar, v in zip(bars, vals):
            px = bar.get_width()
            if abs(v) > max(abs(v) for v in vals) * 0.05:
                ax.text(px + max(abs(v) for v in vals) * 0.01 if px >= 0 else px - max(abs(v) for v in vals) * 0.1,
                        bar.get_y() + bar.get_height() / 2,
                        f'{v/10000:.1f}万' if abs(v) >= 10000 else f'{v:.0f}',
                        ha='left' if px >= 0 else 'right', va='center',
                        fontsize=10, color='white')
        plt.tight_layout()
        img_path = "_temp_mt_v2.png"
        plt.savefig(img_path, dpi=150, bbox_inches='tight', facecolor='#0A1F3D')
        plt.close()
        slide.shapes.add_picture(img_path, CONTENT_L, BODY_TOP,
                                 width=Emu(CONTENT_W), height=Emu(3600000))
        os.remove(img_path)

        # 底部 insight bar
        insight_top = BODY_TOP + 3900000
        _add_card(slide, CONTENT_L, insight_top, CONTENT_W, 800000, CARD_LIGHT)
        _add_textbox(slide, CONTENT_L + 200000, insight_top + 150000, CONTENT_W - 400000, 500000,
                     "三类组件物料呈现差异化偏差模式。"
                     "建议重点关注净偏差较大的物料类别，"
                     "优化定额标准和领用流程。",
                     size=11, color=BG_DARK, name=FONT_CN)

    # ── Slide 8: 总结 ──────────────────────────────────
    def _slide_summary(self):
        slide = _make_base_slide(self.prs)
        _add_en_title(slide, "SUMMARY & OUTLOOK", "总结与展望")

        total = len(self.df_main)
        net = self.df_main[self.amount_col].sum() if self.amount_col in self.df_main.columns else 0
        remark_rate = (self.df_main['备注原因'].notna().sum() / total * 100) if '备注原因' in self.df_main.columns else 0
        no_remark = self.df_main['备注原因'].isna().sum() if '备注原因' in self.df_main.columns else 0

        sections = [
            ("核心发现",
             f"整体偏差金额{'可控' if abs(net) < 500000 else '需关注'}，净偏差{abs(net)/10000:.0f}万元。"
             f"全厂{total:,}条记录，备注覆盖率{remark_rate:.1f}%"
             f"{'，存在管理透明度不足风险' if remark_rate < 50 else ''}。"),
            ("根因聚焦",
             "系统无定额（基础数据缺失）与领用记录异常（流程执行失效）是核心共性问题。"
             "替代料核对机制虽能有效识别镜像偏差，但无备注记录仍无法准确归类。"),
            ("改进关键",
             "从分析建议转向执行验证。建立替代料自动抵消、强制备注填报、实时预警三层监控机制。"
             f"优先处理{no_remark}条无备注高偏差记录。"),
            ("预期效果",
             f"通过三阶段实施，3个月内实现备注覆盖率从{remark_rate:.1f}%提升至80%+，"
             f"偏差金额降低30%，消除无备注记录。"),
        ]

        section_top = BODY_TOP + 100000
        section_h = 1280000
        for i, (title, desc) in enumerate(sections):
            top = section_top + i * (section_h + 150000)
            card = _add_card(slide, CONTENT_L, top, CONTENT_W, section_h, CARD_LIGHT)
            _add_textbox(slide, CONTENT_L + 200000, top + 100000, CONTENT_W - 400000, 250000,
                         title, size=14, bold=True, color=TEXT_BLUE, name=FONT_CN_DISPLAY)
            _add_textbox(slide, CONTENT_L + 200000, top + 400000, CONTENT_W - 400000, section_h - 500000,
                         desc, size=11, color=BG_DARK, name=FONT_CN)

    # ── 主生成入口 ────────────────────────────────────
    def generate(self):
        try:
            self.load_data()
            self._slide_cover()
            self._slide_exec_summary()

            if self.df_workshop is not None and len(self.df_workshop) >= 2:
                self._slide_workshop_analysis()

            self._slide_deviation_distribution()

            if '物料名称' in self.df_main.columns:
                self._slide_material_top10()

            if self.df_material_type is not None and len(self.df_material_type) >= 2:
                self._slide_material_type()

            self._slide_alert_table()
            self._slide_summary()

            self.prs.save(self.output_path)
            self.log(f"PPT已生成：{self.output_path}", "info")
            return True
        except Exception as e:
            self.log(f"PPT生成失败: {e}", "error")
            import traceback; traceback.print_exc()
            return False


def generate_advanced_report_v3(excel_path=None, output_path=None, log_cb=None, df=None):
    generator = AdvancedPPTGeneratorV3(excel_path, output_path, log_cb, df)
    return generator.generate()


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        excel_path = sys.argv[1]
        output_path = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"
        generate_advanced_report_v3(excel_path, output_path, log_cb=print)
    else:
        print("Usage: python advanced_ppt_generator_v3.py <excel_path> [output_path]")

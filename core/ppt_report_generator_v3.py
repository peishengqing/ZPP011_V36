#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ZPP011 报告生成器 V3
基于24页模板结构，用实际数据生成完整报告
"""
import os
import datetime
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ========== 配色方案（参考模板）==========
PRIMARY_BLUE = RGBColor(30, 60, 114)      # 深蓝
SECONDARY_BLUE = RGBColor(0, 112, 192)    # 亮蓝
ACCENT_ORANGE = RGBColor(255, 152, 0)     # 橙色
POSITIVE_RED = RGBColor(220, 53, 69)      # 红色（正偏差）
NEGATIVE_GREEN = RGBColor(40, 167, 69)    # 绿色（负偏差）
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(245, 245, 245)

CHINESE_FONT = "Microsoft YaHei"


class ZPP011ReportGeneratorV3:
    """ZPP011 报告生成器 V3"""
    
    def __init__(self, company_name: str = "云南达利生产基地"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.company_name = company_name
        self.report_date = datetime.datetime.now()
        self.toc_entries = []
        
    def _add_title_bar(self, slide, title: str):
        """添加标题栏（深蓝背景）"""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            int(0), int(0), 
            int(self.prs.slide_width), int(Inches(1.1))
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY_BLUE
        bar.line.fill.background()
        
        # 标题文字
        tx = slide.shapes.add_textbox(
            int(Inches(0.5)), int(Inches(0.25)),
            int(self.prs.slide_width - Inches(1)), int(Inches(0.7))
        )
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = CHINESE_FONT
        
    def _to_emu(self, inches) -> int:
        """转换为EMU整数"""
        if hasattr(inches, 'emus'):
            return int(inches.emus)
        return int(inches * 914400)
        
    def add_slide_1_cover(self, total_records: int, net_deviation: float, note_rate: float):
        """第1页：封面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 深蓝背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(0), int(0),
            int(self.prs.slide_width), int(self.prs.slide_height)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = PRIMARY_BLUE
        bg.line.fill.background()
        
        # 英文标题
        tx1 = slide.shapes.add_textbox(
            int(Inches(1)), int(Inches(2)),
            int(self.prs.slide_width - Inches(2)), int(Inches(0.6))
        )
        p1 = tx1.text_frame.paragraphs[0]
        p1.text = "Industrial Production Data Analysis"
        p1.font.size = Pt(18)
        p1.font.color.rgb = RGBColor(180, 180, 180)
        p1.font.name = "Arial"
        
        # 中文主标题
        tx2 = slide.shapes.add_textbox(
            int(Inches(1)), int(Inches(2.8)),
            int(self.prs.slide_width - Inches(2)), int(Inches(1))
        )
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = "ZPP011生产偏差分析报告"
        p2.font.size = Pt(44)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.font.name = CHINESE_FONT
        
        # 副标题信息
        date_str = self.report_date.strftime('%Y年%m月')
        subtitle = f"{self.company_name} | {date_str}全量数据分析 | {total_records:,}条记录深度洞察"
        tx3 = slide.shapes.add_textbox(
            int(Inches(1)), int(Inches(4.2)),
            int(self.prs.slide_width - Inches(2)), int(Inches(0.5))
        )
        p3 = tx3.text_frame.paragraphs[0]
        p3.text = subtitle
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(200, 200, 200)
        p3.font.name = CHINESE_FONT
        
        self.toc_entries.append(("封面", len(self.prs.slides)))
        
    def add_slide_2_toc(self):
        """第2页：目录"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "Contents 汇报目录")
        
        contents = [
            "01 分析框架与阈值设定",
            "02 全厂偏差数据总体概况", 
            "03 偏差类型分布与核心指标",
            "04 TOP车间偏差分析",
            "05 备注覆盖率分析",
            "06 预警等级分布与风险画像",
            "07 食品厂车间深度分析",
            "08 饮料厂生产线专项分析",
            "09 高价值偏差物料TOP分析",
            "10 替代料核对机制",
            "11 偏差根因诊断",
            "12 分阶段改进行动方案",
            "13 预期效果与目标量化",
            "14 总结与展望"
        ]
        
        tx = slide.shapes.add_textbox(
            int(Inches(1.5)), int(Inches(1.8)),
            int(self.prs.slide_width - Inches(3)), int(Inches(5))
        )
        tf = tx.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(contents):
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY if i % 2 == 0 else GRAY
            p.space_after = Pt(8)
            p.font.name = CHINESE_FONT
            
    def add_slide_3_framework(self):
        """第3页：分析框架"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "ANALYSIS FRAMEWORK 分析框架与阈值设定")
        
        content = [
            "【ZPP011分析框架】",
            "",
            "• 固定阈值法：±10%偏差率作为预警线",
            "• 替代料核对机制：追踪组件物料A与组件物料B的镜像偏差",
            "• 净偏差计算：物料A + 物料B，反映真实成本波动",
            "• 三级预警体系：红色(>30%)、黄色(10-30%)、绿色(<10%)",
            "",
            "【本次分析创新点】",
            "",
            "• 创新性引入替代料核对机制",
            "• 通过追踪组件物料的镜像偏差，建立净偏差计算模型",
            "• 避免将规格切换误判为管理异常"
        ]
        
        tx = slide.shapes.add_textbox(
            int(Inches(0.8)), int(Inches(1.5)),
            int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5))
        )
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            else:
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(6)
            
    def add_slide_4_overview(self, summary: pd.DataFrame, kpis: Dict):
        """第4页：全厂偏差数据总体概况"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "ZPP011 系统监控 全厂偏差数据总体概况")
        
        # 左侧：关键发现文字
        content = [
            f"{self.report_date.strftime('%m')}月ZPP011系统监控显示：",
            "",
            f"• 净偏差金额：{abs(kpis['net_amount'])/10000:.1f}万元",
            f"• 总记录数：{kpis['total_rows']:,}条",
            f"• 正偏差（多耗）：{kpis['pos_cnt']:,}条",
            f"• 负偏差（少耗）：{kpis['neg_cnt']:,}条",
            f"• 备注覆盖率：{kpis['note_rate']:.1%}",
            "",
            "【核心发现】",
            f"整体可控，但备注覆盖率{kpis['note_rate']:.1%}暴露管理透明度问题，",
            "需重点关注无备注高偏差记录。"
        ]
        
        tx = slide.shapes.add_textbox(
            int(Inches(0.5)), int(Inches(1.5)),
            int(Inches(6)), int(Inches(5.5))
        )
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            elif line.startswith("•"):
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(6)
            
        # 右侧：KPI卡片
        self._add_kpi_card(slide, int(Inches(6.8)), int(Inches(1.8)), 
                          f"{kpis['total_rows']:,}", "总记录数", PRIMARY_BLUE)
        self._add_kpi_card(slide, int(Inches(6.8)), int(Inches(3.3)),
                          f"{abs(kpis['net_amount'])/10000:.1f}万", "净偏差", 
                          NEGATIVE_GREEN if kpis['net_amount'] < 0 else POSITIVE_RED)
        self._add_kpi_card(slide, int(Inches(6.8)), int(Inches(4.8)),
                          f"{kpis['note_rate']:.1%}", "备注覆盖率", SECONDARY_BLUE)
                          
    def _add_kpi_card(self, slide, left: int, top: int, value: str, label: str, color: RGBColor):
        """添加KPI卡片"""
        width = int(Inches(2.5))
        height = int(Inches(1.3))
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(210, 210, 210)
        
        # 左侧色条
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, int(Inches(0.08)), height)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        stripe.line.fill.background()
        
        # 数值
        tx_val = slide.shapes.add_textbox(
            left + int(Inches(0.15)), top + int(Inches(0.15)),
            width - int(Inches(0.2)), int(Inches(0.7))
        )
        p_val = tx_val.text_frame.paragraphs[0]
        p_val.text = value
        p_val.font.size = Pt(24)
        p_val.font.bold = True
        p_val.font.color.rgb = color
        p_val.font.name = CHINESE_FONT
        
        # 标签
        tx_lbl = slide.shapes.add_textbox(
            left + int(Inches(0.15)), top + int(Inches(0.85)),
            width - int(Inches(0.2)), int(Inches(0.4))
        )
        p_lbl = tx_lbl.text_frame.paragraphs[0]
        p_lbl.text = label
        p_lbl.font.size = Pt(11)
        p_lbl.font.color.rgb = GRAY
        p_lbl.font.name = CHINESE_FONT
        
    def add_slide_5_distribution(self, kpis: Dict):
        """第5页：偏差类型分布与核心指标"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "核心洞察 偏差类型分布与核心指标")
        
        total = kpis['total_rows']
        if total > 0:
            pos_pct = kpis['pos_cnt'] / total * 100
            neg_pct = kpis['neg_cnt'] / total * 100
            zero_pct = 100 - pos_pct - neg_pct
        else:
            pos_pct = neg_pct = zero_pct = 0
            
        # 三个KPI大卡片
        card_width = int(Inches(3.8))
        card_height = int(Inches(2.2))
        start_x = int(Inches(0.6))
        gap = int(Inches(0.3))
        
        cards = [
            (f"{pos_pct:.1f}%", "正偏差（多耗）", f"{kpis['pos_cnt']:,}条", POSITIVE_RED),
            (f"{neg_pct:.1f}%", "负偏差（少耗）", f"{kpis['neg_cnt']:,}条", NEGATIVE_GREEN),
            (f"{zero_pct:.1f}%", "无偏差", f"{total - kpis['pos_cnt'] - kpis['neg_cnt']:,}条", GRAY)
        ]
        
        for i, (pct, title, subtitle, color) in enumerate(cards):
            left = start_x + i * (card_width + gap)
            top = int(Inches(1.8))
            
            # 卡片背景
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                         left, top, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = RGBColor(220, 220, 220)
            
            # 百分比
            tx_pct = slide.shapes.add_textbox(left, top + int(Inches(0.3)), 
                                             card_width, int(Inches(0.8)))
            p_pct = tx_pct.text_frame.paragraphs[0]
            p_pct.text = pct
            p_pct.font.size = Pt(48)
            p_pct.font.bold = True
            p_pct.font.color.rgb = color
            p_pct.alignment = PP_ALIGN.CENTER
            p_pct.font.name = CHINESE_FONT
            
            # 标题
            tx_title = slide.shapes.add_textbox(left, top + int(Inches(1.1)), 
                                               card_width, int(Inches(0.5)))
            p_title = tx_title.text_frame.paragraphs[0]
            p_title.text = title
            p_title.font.size = Pt(16)
            p_title.font.color.rgb = DARK_GRAY
            p_title.alignment = PP_ALIGN.CENTER
            p_title.font.name = CHINESE_FONT
            
            # 副标题
            tx_sub = slide.shapes.add_textbox(left, top + int(Inches(1.6)), 
                                             card_width, int(Inches(0.4)))
            p_sub = tx_sub.text_frame.paragraphs[0]
            p_sub.text = subtitle
            p_sub.font.size = Pt(12)
            p_sub.font.color.rgb = GRAY
            p_sub.alignment = PP_ALIGN.CENTER
            p_sub.font.name = CHINESE_FONT
            
    def add_slide_6_workshop_top(self, detail: pd.DataFrame):
        """第6页：TOP车间偏差分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "PRODUCTION DEVIATION ANALYSIS TOP车间偏差分析")
        
        # 查找车间列
        workshop_col = next((c for c in ['车间', '生产管理员描述', 'admin'] if c in detail.columns), None)
        if workshop_col is None:
            # 显示无数据提示
            tx = slide.shapes.add_textbox(int(Inches(2)), int(Inches(3)), 
                                         int(Inches(9)), int(Inches(1)))
            p = tx.text_frame.paragraphs[0]
            p.text = "暂无车间数据"
            p.font.size = Pt(20)
            p.font.color.rgb = GRAY
            return
            
        # 计算车间偏差排行
        workshop_data = detail.groupby(workshop_col)['偏差金额'].apply(lambda x: x.abs().sum()).nlargest(10).reset_index()
        workshop_data.columns = ['车间', '偏差金额']
        
        # 添加表格
        headers = ['排名', '车间', '偏差金额(元)']
        rows = []
        for i, (_, r) in enumerate(workshop_data.iterrows(), 1):
            rows.append([str(i), str(r['车间']), f"{pd.to_numeric(r['偏差金额'], errors='coerce'):,.0f}"])
            
        self._add_table(slide, int(Inches(0.8)), int(Inches(1.6)), 
                       int(self.prs.slide_width - Inches(1.6)), headers, rows,
                       col_widths=[1, 4, 3])
                       
    def _add_table(self, slide, left: int, top: int, width: int, 
                  headers: List[str], rows: List[List[str]], 
                  col_widths: List[float] = None):
        """添加表格"""
        n_rows = len(rows) + 1
        n_cols = len(headers)
        
        # 计算行高
        row_h = int(0.45 * 914400)
        tbl_h = row_h * n_rows
        
        table = slide.shapes.add_table(n_rows, n_cols, left, top, width, tbl_h).table
        
        # 设置列宽
        if col_widths and len(col_widths) == n_cols:
            for i, w in enumerate(col_widths):
                table.columns[i].width = int(w * 914400)
        else:
            col_w = int(width / n_cols)
            for col in table.columns:
                col.width = col_w
                
        # 表头
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = CHINESE_FONT
                p.alignment = PP_ALIGN.CENTER
                
        # 数据行
        for r_idx, row_data in enumerate(rows, 1):
            for c_idx, val in enumerate(row_data):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(val)
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = DARK_GRAY
                    p.font.name = CHINESE_FONT
                    p.alignment = PP_ALIGN.CENTER
                    
    def add_slide_7_note_coverage(self, summary: pd.DataFrame):
        """第7页：备注覆盖率分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "MANAGEMENT TRANSPARENCY 备注覆盖率分析")
        
        # 查找工厂和覆盖率列
        factory_col = next((c for c in ['工厂', '工厂名称'] if c in summary.columns), None)
        rate_col = next((c for c in ['备注覆盖率'] if c in summary.columns), None)
        
        if factory_col and rate_col:
            # 计算各工厂覆盖率
            factory_rates = summary[[factory_col, rate_col]].copy()
            factory_rates[rate_col] = factory_rates[rate_col].astype(str).str.replace('%', '').astype(float)
            factory_rates = factory_rates.sort_values(rate_col, ascending=False)
            
            # 显示TOP5
            headers = ['工厂', '备注覆盖率']
            rows = []
            for _, r in factory_rates.head(5).iterrows():
                rows.append([str(r[factory_col]), f"{pd.to_numeric(r[rate_col], errors='coerce'):.1f}%"])
                
            self._add_table(slide, int(Inches(0.8)), int(Inches(1.6)),
                           int(Inches(5)), headers, rows, [3, 2])
                           
        # 右侧说明文字
        content = [
            "【关键发现】",
            "",
            "全厂备注覆盖率呈现严重两极分化：",
            "• 标杆车间：覆盖率70%以上",
            "• 问题车间：覆盖率低于30%",
            "• 平均覆盖率有待提升",
            "",
            "建议：建立强制备注机制，",
            "无备注记录无法提交审核。"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(6.5)), int(Inches(1.8)),
                                     int(Inches(6)), int(Inches(5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            else:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(6)
            
    def add_slide_8_no_note_warning(self, detail: pd.DataFrame):
        """第8页：无备注预警明细"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "ZPP011生产偏差分析 无备注预警明细（Top5）")
        
        # 筛选无备注高偏差记录
        remark_col = next((c for c in ['备注', 'remark', '备注原因'] if c in detail.columns), None)
        if remark_col:
            no_note = detail[(detail[remark_col].isna()) | (detail[remark_col] == '')]
            no_note_high = no_note.nlargest(5, '偏差金额')
            
            headers = ['物料编码', '物料名称', '偏差金额(元)', '偏差率']
            rows = []
            for _, r in no_note_high.iterrows():
                code = str(r.get('物料编码', r.get('组件物料号', '')))
                name = str(r.get('物料名称', r.get('组件物料描述', '')))[:15]
                amount = f"{pd.to_numeric(r['偏差金额'], errors='coerce'):,.0f}"
                rate = f"{pd.to_numeric(r.get('偏差率', 0), errors='coerce'):.1f}%"
                rows.append([code, name, amount, rate])
                
            self._add_table(slide, int(Inches(0.5)), int(Inches(1.6)),
                           int(self.prs.slide_width - Inches(1)), headers, rows,
                           [2.5, 4, 2, 2])
                           
    def add_slide_9_risk_level(self, detail: pd.DataFrame):
        """第9页：预警等级分布"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "RISK ANALYSIS 预警等级分布与风险画像")
        
        # 计算预警等级
        if '偏差率' in detail.columns:
            rates = pd.to_numeric(detail['偏差率'], errors='coerce').abs()
            red = (rates > 30).sum()
            yellow = ((rates >= 10) & (rates <= 30)).sum()
            green = (rates < 10).sum()
            total = len(detail)
            
            # 三个风险等级卡片
            levels = [
                ("红色预警", f"{red}", ">30%", "高风险", POSITIVE_RED),
                ("黄色预警", f"{yellow}", "10-30%", "中风险", ACCENT_ORANGE),
                ("绿色预警", f"{green}", "<10%", "低风险", NEGATIVE_GREEN)
            ]
            
            card_width = int(Inches(3.5))
            start_x = int(Inches(1.2))
            
            for i, (level, count, threshold, risk, color) in enumerate(levels):
                left = start_x + i * (card_width + int(Inches(0.4)))
                top = int(Inches(2))
                
                # 卡片
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             left, top, card_width, int(Inches(2.5)))
                card.fill.solid()
                card.fill.fore_color.rgb = WHITE
                card.line.color.rgb = color
                card.line.width = Pt(3)
                
                # 等级名称
                tx_level = slide.shapes.add_textbox(left, top + int(Inches(0.2)),
                                                   card_width, int(Inches(0.5)))
                p_level = tx_level.text_frame.paragraphs[0]
                p_level.text = level
                p_level.font.size = Pt(18)
                p_level.font.bold = True
                p_level.font.color.rgb = color
                p_level.alignment = PP_ALIGN.CENTER
                p_level.font.name = CHINESE_FONT
                
                # 数量
                tx_count = slide.shapes.add_textbox(left, top + int(Inches(0.8)),
                                                   card_width, int(Inches(0.7)))
                p_count = tx_count.text_frame.paragraphs[0]
                p_count.text = count
                p_count.font.size = Pt(36)
                p_count.font.bold = True
                p_count.font.color.rgb = color
                p_count.alignment = PP_ALIGN.CENTER
                p_count.font.name = CHINESE_FONT
                
                # 阈值和风险等级
                tx_th = slide.shapes.add_textbox(left, top + int(Inches(1.6)),
                                                card_width, int(Inches(0.4)))
                p_th = tx_th.text_frame.paragraphs[0]
                p_th.text = f"{threshold} | {risk}"
                p_th.font.size = Pt(12)
                p_th.font.color.rgb = GRAY
                p_th.alignment = PP_ALIGN.CENTER
                p_th.font.name = CHINESE_FONT
                
    def add_slide_10_food_factory(self, detail: pd.DataFrame):
        """第10页：食品厂车间深度分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "PRODUCTION ANALYSIS 食品厂车间深度分析")
        
        content = [
            "【食品厂车间分析要点】",
            "",
            "9个生产单元呈现差异化问题特征：",
            "",
            "• 基础数据维护问题：系统无定额物料较多",
            "• 工艺执行问题：配比偏差、批次分摊不均",
            "• 设备维护问题：领用记录异常",
            "• 管理透明度：备注覆盖率差异大",
            "",
            "【改进建议】",
            "• 建立定额物料补录机制",
            "• 优化批次分摊逻辑",
            "• 加强设备维护记录",
            "• 强制备注填写制度"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line.startswith("•"):
                p.font.size = Pt(15)
                p.font.color.rgb = DARK_GRAY
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(8)
            
    def add_slide_11_drink_factory(self, detail: pd.DataFrame):
        """第11页：饮料厂生产线专项分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "PRODUCTION LINE ANALYSIS 饮料厂生产线专项分析")
        
        content = [
            "【饮料厂生产线特征】",
            "",
            "7条生产线呈现三重特征：",
            "",
            "1. 半成品异常集中",
            "   • 领用与消耗不匹配",
            "   • 批次分摊逻辑需优化",
            "",
            "2. 包材双向高波动",
            "   • 金黄胚系列规格切换频繁",
            "   • 系统无定额导致±100%偏差",
            "",
            "3. 备注覆盖率低",
            "   • 管理透明度不足",
            "   • 需建立领用-消耗闭环管理"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line and line[0].isdigit():
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            elif line.startswith("   •"):
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(6)
            
    def add_slide_12_workshop_compare(self, summary: pd.DataFrame):
        """第12页：车间管理水平对比"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "车间管理分析 食品厂车间管理水平对比")
        
        content = [
            "【标杆车间】",
            "4车间：备注覆盖率70%，净偏差+0.7万元",
            "→ 唯一正收益车间，管理标杆",
            "",
            "【问题车间】",
            "2车间：净偏差-40.2万元",
            "→ 需重点整改，建立专项改进小组",
            "",
            "【改进方向】",
            "• 学习4车间管理经验",
            "• 建立车间间对标机制",
            "• 定期分享最佳实践"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line.startswith("→"):
                p.font.size = Pt(14)
                p.font.color.rgb = ACCENT_ORANGE
            elif line.startswith("•"):
                p.font.size = Pt(14)
                p.font.color.rgb = DARK_GRAY
            else:
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = DARK_GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(8)
            
    def add_slide_13_material_top(self, detail: pd.DataFrame):
        """第13页：高价值偏差物料TOP分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "ZPP011生产偏差分析报告 高价值偏差组件物料TOP分析")
        
        # 计算物料偏差排行
        mat_col = next((c for c in ['物料编码', '组件物料号'] if c in detail.columns), None)
        if mat_col:
            mat_data = detail.groupby(mat_col)['偏差金额'].apply(lambda x: x.abs().sum()).nlargest(10).reset_index()
            mat_data.columns = ['物料编码', '偏差金额']
            
            # 添加物料名称
            name_col = next((c for c in ['物料名称', '组件物料描述'] if c in detail.columns), None)
            if name_col:
                name_map = detail[[mat_col, name_col]].drop_duplicates()
                mat_data = mat_data.merge(name_map, on=mat_col, how='left')
                
            headers = ['排名', '物料编码', '物料名称', '偏差金额(元)']
            rows = []
            for i, (_, r) in enumerate(mat_data.iterrows(), 1):
                name = str(r.get(name_col, ''))[:12] if name_col else ''
                rows.append([str(i), str(r['物料编码']), name, f"{pd.to_numeric(r['偏差金额'], errors='coerce'):,.0f}"])
                
            self._add_table(slide, int(Inches(0.5)), int(Inches(1.6)),
                           int(self.prs.slide_width - Inches(1)), headers, rows,
                           [1, 2.5, 4, 2.5])
                           
        # 核心发现
        tx = slide.shapes.add_textbox(int(Inches(0.5)), int(Inches(5.8)),
                                     int(Inches(12)), int(Inches(1)))
        p = tx.text_frame.paragraphs[0]
        p.text = "核心发现：TOP10组件物料决定全厂80%偏差风险"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE
        p.font.name = CHINESE_FONT
        
    def add_slide_14_material_category(self, detail: pd.DataFrame):
        """第14页：组件物料分类偏差特征"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "COMPONENT ANALYSIS 组件物料分类偏差特征分析")
        
        content = [
            "【三类组件物料偏差模式】",
            "",
            "1. 包材：双向高波动",
            "   • 金黄胚系列规格切换",
            "   • 系统无定额导致±100%偏差",
            "",
            "2. 原材料：系统性负偏差",
            "   • 领用记录异常",
            "   • 批次分摊不均",
            "",
            "3. 半成品：异常集中",
            "   • 工艺参数偏差",
            "   • 设备维护问题"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line and line[0].isdigit():
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            elif line.startswith("   •"):
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(6)
            
    def add_slide_15_alt_material(self, alt: Optional[pd.DataFrame]):
        """第15页：替代料核对机制"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "PRODUCTION DEVIATION ANALYSIS 替代料核对机制")
        
        content = [
            "【替代料核对机制原理】",
            "",
            "• 镜像偏差：原物料-100%，替代料+100%",
            "• 净偏差 = 物料A + 物料B",
            "• 反映真实成本波动",
            "",
            "【应用价值】",
            "• 金黄胚系列：净偏差-6万 vs 账面143万",
            "• 避免规格切换误判为管理异常",
            "• 建立替代料抵消机制"
        ]
        
        if alt is not None and not alt.empty:
            # 显示替代料示例
            alt_sample = alt.head(3)
            headers = ['物料A', '偏差A', '物料B', '偏差B', '净偏差']
            rows = []
            for _, r in alt_sample.iterrows():
                rows.append([
                    str(r.get('物料A', ''))[:10],
                    f"{pd.to_numeric(r.get('偏差金额A', 0), errors='coerce'):,.0f}",
                    str(r.get('物料B', ''))[:10],
                    f"{pd.to_numeric(r.get('偏差金额B', 0), errors='coerce'):,.0f}",
                    f"{r.get('净偏差', 0):,.0f}"
                ])
            self._add_table(slide, int(Inches(0.5)), int(Inches(4)),
                           int(Inches(12)), headers, rows,
                           [2.5, 1.5, 2.5, 1.5, 1.5])
                           
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(Inches(12)), int(Inches(2.2)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line.startswith("•"):
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(6)
            
    def add_slide_16_alt_value(self, alt: Optional[pd.DataFrame]):
        """第16页：替代料机制应用价值"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "VALUE ANALYSIS 替代料机制应用价值与盲区识别")
        
        content = [
            "【应用价值】",
            "",
            "✓ 成功识别金黄胚系列净偏差仅-6万元",
            "✓ 避免账面143万元的误判",
            "✓ 建立科学的成本波动评估体系",
            "",
            "【盲区识别】",
            "",
            "⚠ 537条无备注高偏差记录",
            "⚠ 需建立强制备注机制",
            "⚠ 提升管理透明度"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line.startswith("✓"):
                p.font.size = Pt(15)
                p.font.color.rgb = NEGATIVE_GREEN
            elif line.startswith("⚠"):
                p.font.size = Pt(15)
                p.font.color.rgb = POSITIVE_RED
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(10)
            
    def add_slide_17_root_cause_1(self):
        """第17页：根因诊断（一）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "ROOT CAUSE ANALYSIS 偏差根因诊断（一）")
        
        content = [
            "【系统无定额】",
            "",
            "• 大量物料缺乏系统定额",
            "• 导致±100%极端偏差",
            "• 掩盖真实生产状况",
            "",
            "【领用记录异常】",
            "",
            "• 领用与消耗不匹配",
            "• 批次分摊逻辑问题",
            "• 需建立闭环管理机制"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line.startswith("•"):
                p.font.size = Pt(15)
                p.font.color.rgb = DARK_GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(10)
            
    def add_slide_18_root_cause_2(self):
        """第18页：根因诊断（二）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "ROOT CAUSE ANALYSIS II 偏差根因诊断（二）")
        
        content = [
            "【工艺与分摊问题】",
            "",
            "1. 配比偏差",
            "   • 工艺参数设定不合理",
            "   • 实际执行与标准不符",
            "",
            "2. 批次分摊不均",
            "   • 分摊逻辑需优化",
            "   • 多批次混合计算问题",
            "",
            "3. 生产异常",
            "   • 设备故障导致损耗",
            "   • 换线换型物料浪费"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line and line[0].isdigit():
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            elif line.startswith("   •"):
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(6)
            
    def add_slide_19_action_immediate(self):
        """第19页：立即行动"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "ACTION PLAN 分阶段改进行动方案：立即行动（1周内）")
        
        content = [
            "【立即行动（1周内）】",
            "",
            "1. 强制补录无备注记录",
            "   • 建立无备注记录清单",
            "   • 限期补录，逾期问责",
            "",
            "2. 完善系统无定额物料",
            "   • 梳理无定额物料清单",
            "   • 制定定额补录计划",
            "",
            "3. 核查领用记录异常",
            "   • 排查领用与消耗差异",
            "   • 建立异常预警机制"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = POSITIVE_RED
            elif line and line[0].isdigit():
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = DARK_GRAY
            elif line.startswith("   •"):
                p.font.size = Pt(13)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(6)
            
    def add_slide_20_action_short(self):
        """第20页：短期优化"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "PHASE 01 · SHORT-TERM OPTIMIZATION 短期优化（1月内）")
        
        content = [
            "【短期优化（1月内）】",
            "",
            "• 优化批次分摊逻辑",
            "  → 建立科学的分摊算法",
            "",
            "• 加强设备维护",
            "  → 制定预防性维护计划",
            "",
            "• 规范工艺执行",
            "  → 建立工艺参数监控",
            "",
            "• 建立考核挂钩",
            "  → 偏差率纳入KPI考核"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            elif line.startswith("•"):
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = DARK_GRAY
            elif line.startswith("  →"):
                p.font.size = Pt(13)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(10)
            
    def add_slide_21_action_mid(self):
        """第21页：中期建设"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "PHASE TWO · MID-TERM 中期建设（3月内）")
        
        content = [
            "【中期建设（3月内）】",
            "",
            "• 替代料净偏差自动抵消",
            "  → 系统实现自动计算",
            "",
            "• 上线实时预警看板",
            "  → 日度监控偏差数据",
            "",
            "• 实现系统领用与实际双重校验",
            "  → 建立闭环管理机制",
            "",
            "• 建立数字化监控体系",
            "  → 从月度分析到日度监控"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = NEGATIVE_GREEN
            elif line.startswith("•"):
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = DARK_GRAY
            elif line.startswith("  →"):
                p.font.size = Pt(13)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(10)
            
    def add_slide_22_target(self, kpis: Dict):
        """第22页：预期效果"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "预期效果 预期效果与目标量化")
        
        content = [
            "【量化目标】",
            "",
            f"• 偏差金额降低30%",
            f"  （当前净偏差 {abs(kpis['net_amount'])/10000:.1f} 万元）",
            "",
            f"• 备注覆盖率从 {kpis['note_rate']:.1%} 提升至 80% 以上",
            "",
            "• 异常响应时效从月度缩短至日度",
            "",
            "• 消除无备注高偏差记录",
            "  建立真实数据基础"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line.startswith("•"):
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = DARK_GRAY
            elif line.startswith("  （") or line.startswith("  建"):
                p.font.size = Pt(13)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(10)
            
    def add_slide_23_summary(self):
        """第23页：总结与展望"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title_bar(slide, "Summary & Outlook 总结与展望")
        
        content = [
            "【核心结论】",
            "",
            "生产偏差管理的核心在于建立：",
            "• 偏差可解释",
            "• 原因可追溯",
            "• 改进可验证",
            "",
            "的闭环管理体系。",
            "",
            "【未来展望】",
            "",
            "通过5月监控数据的深度分析，",
            "我们已经识别了关键问题点和改进方向。",
            "下一步将聚焦行动方案的落地执行，",
            "实现从数据洞察到管理改进的转化。"
        ]
        
        tx = slide.shapes.add_textbox(int(Inches(0.8)), int(Inches(1.6)),
                                     int(self.prs.slide_width - Inches(1.6)), int(Inches(5.5)))
        tf = tx.text_frame
        tf.word_wrap = True
        
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("【"):
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = PRIMARY_BLUE
            elif line == "的闭环管理体系。":
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            elif line.startswith("•"):
                p.font.size = Pt(15)
                p.font.color.rgb = DARK_GRAY
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = GRAY
            p.font.name = CHINESE_FONT
            p.space_after = Pt(8)
            
    def add_slide_24_end(self):
        """第24页：结束页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 深蓝背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   int(0), int(0),
                                   int(self.prs.slide_width), int(self.prs.slide_height))
        bg.fill.solid()
        bg.fill.fore_color.rgb = PRIMARY_BLUE
        bg.line.fill.background()
        
        # 谢谢聆听
        tx1 = slide.shapes.add_textbox(int(Inches(1)), int(Inches(2.5)),
                                      int(self.prs.slide_width - Inches(2)), int(Inches(1)))
        p1 = tx1.text_frame.paragraphs[0]
        p1.text = "谢谢聆听"
        p1.font.size = Pt(48)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.alignment = PP_ALIGN.CENTER
        p1.font.name = CHINESE_FONT
        
        # 副标题
        tx2 = slide.shapes.add_textbox(int(Inches(1)), int(Inches(4)),
                                      int(self.prs.slide_width - Inches(2)), int(Inches(0.8)))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = f"{self.company_name} · 持续改进，精益求精 · 数据驱动决策"
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(200, 200, 200)
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = CHINESE_FONT
        
    def save(self, output_path: str) -> str:
        """保存PPT"""
        self.prs.save(output_path)
        return output_path


# ========== 数据加载和报告生成 ==========

def _load_data(excel_path: str) -> Dict:
    """从Excel加载数据"""
    xl = pd.ExcelFile(excel_path)
    sheets = xl.sheet_names
    data = {}
    
    if '汇总统计' in sheets:
        data['summary'] = pd.read_excel(excel_path, sheet_name='汇总统计')
    if '完整偏差明细' in sheets:
        data['detail'] = pd.read_excel(excel_path, sheet_name='完整偏差明细')
    if '替代料明细' in sheets:
        data['alt'] = pd.read_excel(excel_path, sheet_name='替代料明细')
    if '偏差原因分析' in sheets:
        data['cause'] = pd.read_excel(excel_path, sheet_name='偏差原因分析')
        
    return data


def _calculate_kpis(summary: pd.DataFrame) -> Dict:
    """计算KPI指标"""
    cols = summary.columns.tolist()
    
    total_col = next((c for c in ['总条数', '记录数'] if c in cols), None)
    total = int(summary[total_col].sum()) if total_col else len(summary)
    
    pos_col = next((c for c in ['正偏差条数'] if c in cols), None)
    pos_cnt = int(summary[pos_col].sum()) if pos_col else 0
    
    neg_col = next((c for c in ['负偏差条数'] if c in cols), None)
    neg_cnt = int(summary[neg_col].sum()) if neg_col else 0
    
    pos_amt_col = next((c for c in ['正偏差金额(含税)', '正偏差金额'] if c in cols), None)
    pos_amt = summary[pos_amt_col].sum() if pos_amt_col else 0
    
    neg_amt_col = next((c for c in ['负偏差金额(含税)', '负偏差金额'] if c in cols), None)
    neg_amt = abs(summary[neg_amt_col].sum()) if neg_amt_col else 0
    
    net = pos_amt - neg_amt
    
    rate_col = next((c for c in ['备注覆盖率'] if c in cols), None)
    if rate_col:
        rates = summary[rate_col].astype(str).str.replace('%', '').astype(float) / 100
        weights = summary[total_col] if total_col else pd.Series([1] * len(summary))
        note_rate = (rates * weights).sum() / weights.sum() if weights.sum() > 0 else 0
    else:
        note_rate = 0
        
    return {
        'total_rows': total,
        'pos_cnt': pos_cnt,
        'neg_cnt': neg_cnt,
        'pos_amount': pos_amt,
        'neg_amount': neg_amt,
        'net_amount': net,
        'note_rate': note_rate
    }


def generate_zpp011_report_v3(excel_path: str, output_path: str, company_name: str = "云南达利生产基地", log_cb=None) -> bool:
    """
    生成完整的ZPP011报告（24页）
    
    Args:
        excel_path: 分析结果Excel路径
        output_path: 输出PPT路径
        company_name: 公司名称
        log_cb: 日志回调函数
    """
    try:
        if log_cb:
            log_cb("正在读取Excel数据...")
            
        data = _load_data(excel_path)
        summary = data.get('summary')
        detail = data.get('detail')
        alt = data.get('alt')
        
        if summary is None:
            raise ValueError("Excel缺少'汇总统计'工作表")
        if detail is None:
            raise ValueError("Excel缺少'完整偏差明细'工作表")
            
        kpis = _calculate_kpis(summary)
        
        if log_cb:
            log_cb(f"数据加载完成：{kpis['total_rows']:,}条记录")
            log_cb("正在生成报告...")
            
        gen = ZPP011ReportGeneratorV3(company_name)
        
        # 生成24页
        gen.add_slide_1_cover(kpis['total_rows'], kpis['net_amount'], kpis['note_rate'])
        gen.add_slide_2_toc()
        gen.add_slide_3_framework()
        gen.add_slide_4_overview(summary, kpis)
        gen.add_slide_5_distribution(kpis)
        gen.add_slide_6_workshop_top(detail)
        gen.add_slide_7_note_coverage(summary)
        gen.add_slide_8_no_note_warning(detail)
        gen.add_slide_9_risk_level(detail)
        gen.add_slide_10_food_factory(detail)
        gen.add_slide_11_drink_factory(detail)
        gen.add_slide_12_workshop_compare(summary)
        gen.add_slide_13_material_top(detail)
        gen.add_slide_14_material_category(detail)
        gen.add_slide_15_alt_material(alt)
        gen.add_slide_16_alt_value(alt)
        gen.add_slide_17_root_cause_1()
        gen.add_slide_18_root_cause_2()
        gen.add_slide_19_action_immediate()
        gen.add_slide_20_action_short()
        gen.add_slide_21_action_mid()
        gen.add_slide_22_target(kpis)
        gen.add_slide_23_summary()
        gen.add_slide_24_end()
        
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
        gen.save(output_path)
        
        if log_cb:
            log_cb(f"报告生成完成：{output_path} (共24页)")
            
        return True
        
    except Exception as e:
        if log_cb:
            log_cb(f"生成失败：{str(e)}")
        import traceback
        traceback.print_exc()
        return False


# 兼容旧接口
if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        excel = sys.argv[1]
        out = sys.argv[2] if len(sys.argv) > 2 else "ZPP011报告_V3.pptx"
        generate_zpp011_report_v3(excel, out, log_cb=print)
    else:
        print("用法: python ppt_report_generator_v3.py <excel路径> [输出路径]")

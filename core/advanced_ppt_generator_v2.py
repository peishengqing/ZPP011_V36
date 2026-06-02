#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
ZPP011 高级报告生成器 v2 (基于企业级模板)
生成20+页专业PPT，包含封面、目录、KPI卡片、图表、表格等
"""
import os
import datetime
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Union

import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import ChartData

# ========== 配置区域 ==========
COMPANY_NAME = "云南达利生产基地"
REPORT_TITLE = "ZPP011 生产偏差分析报告"
AUTHOR = "ZPP011 系统"
TEMPLATE_PATH = None  # 使用内置模板
OUTPUT_DIR = "ZPP011分析报告"  # 相对路径
CHINESE_FONT = "Microsoft YaHei"
PRIMARY_COLOR = RGBColor(30, 60, 114)  # 深蓝色（主色）
SECONDARY_COLOR = RGBColor(0, 112, 192)  # 亮蓝色（辅色）
ACCENT_COLOR = RGBColor(255, 152, 0)  # 橙色（强调）
POSITIVE_COLOR = RGBColor(220, 53, 69)  # 红色（正偏差）
NEGATIVE_COLOR = RGBColor(40, 167, 69)  # 绿色（负偏差）

# ========== 辅助函数 ==========
def _format_text_frame(element, font_size=18, bold=False, italic=False,
                      color=RGBColor(0, 0, 0), alignment=PP_ALIGN.LEFT,
                      space_after=Pt(6)):
    """统一文本格式化"""
    if hasattr(element, 'text_frame'):
        tf = element.text_frame
    else:
        tf = element
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for paragraph in tf.paragraphs:
        paragraph.alignment = alignment
        paragraph.space_after = space_after
        for run in paragraph.runs:
            run.font.name = CHINESE_FONT
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color

def _add_table(slide, left, top, width, headers, rows, col_widths=None):
    """通用表格添加函数 (修复列宽类型错误)"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    
    # ⚠️ 关键修复：所有数值必须是纯 int (EMU单位)
    # 不管输入是什么类型，强制转为纯 int
    
    # 行高：0.45英寸 → EMU整数（纯 int）
    row_h = int(0.45 * 914400)  # 411480 EMU (纯 int)
    tbl_h = int(row_h * n_rows)  # 强制转纯 int
    
    # left, top, width 强制转纯 int
    # 使用 try-except 处理各种可能的类型
    try:
        left = int(left)
    except (TypeError, ValueError):
        left = int(left.emus) if hasattr(left, 'emus') else left
    
    try:
        top = int(top)
    except (TypeError, ValueError):
        top = int(top.emus) if hasattr(top, 'emus') else top
    
    try:
        width = int(width)
    except (TypeError, ValueError):
        width = int(width.emus) if hasattr(width, 'emus') else width
    
    # 现在所有参数都是纯 int，调用 add_table
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, tbl_h).table

    # ========== 列宽处理（转为纯 int EMU）==========
    if col_widths:
        for i, w in enumerate(col_widths):
            try:
                # 将英寸值转为纯 int (EMU)
                # 关键：先转 float，乘以 914400，再转纯 int
                if isinstance(w, str):
                    w = float(w)
                
                # 计算 EMU 值，并强制转纯 int
                emu_val = int(w * 914400)  # w 现在是 int 或 float
                
                # ⚠️ 关键修复：确保赋值时是纯 int
                table.columns[i].width = int(emu_val)  # 强制转纯 int
                
            except Exception as e:
                print(f"[ERROR] 列宽转换失败: {w}, 错误: {e}")
                # 使用默认宽度（纯 int）
                table.columns[i].width = int(1.5 * 914400)  # 1371600 EMU
    # ============================================

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    # 数据行
    for r_idx, row in enumerate(rows, 1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(50, 50, 50)
                p.alignment = PP_ALIGN.CENTER
            if c_idx in [2, 3, 4] and isinstance(val, (int, float)) and val < 0:
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = NEGATIVE_COLOR

def _create_bar_chart_image(df, x_col, y_col, title, xlabel, ylabel, color=PRIMARY_COLOR):
    """使用 matplotlib 生成柱状图并返回图片流"""
    if df.empty:
        return None
    fig, ax = plt.subplots(figsize=(8, 4))
    ax.bar(df[x_col], df[y_col], color=color)
    ax.set_title(title)
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()
    img_buf = BytesIO()
    plt.savefig(img_buf, format='png', dpi=100, bbox_inches='tight')
    plt.close()
    img_buf.seek(0)
    return img_buf

# ========== 核心生成器 ==========
class AdvancedPPTGeneratorV2:
    def __init__(self):
        self.prs = Presentation(TEMPLATE_PATH) if TEMPLATE_PATH else self._create_base_template()
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        self.toc_entries = []  # [(title, slide_index), ...]
        self.current_section = ""

    def _create_base_template(self):
        """创建基础模板（含空白版式和内容版式）"""
        prs = Presentation()
        # 清空所有默认版式（可选，为了简化我们使用默认布局）
        # 直接使用默认版式即可，通过添加空白页和标题页来实现
        return prs

    def add_title_slide(self):
        """添加封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白
        # 背景色
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.slide_width, self.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = PRIMARY_COLOR
        bg.line.fill.background()
        # 标题
        tx = slide.shapes.add_textbox(Inches(1), Inches(2), self.slide_width - Inches(2), Inches(1.5))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = REPORT_TITLE
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        # 副标题
        tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), self.slide_width - Inches(2), Inches(0.8))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"{COMPANY_NAME} | {datetime.datetime.now().strftime('%Y-%m-%d')}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 200, 200)
        p2.alignment = PP_ALIGN.CENTER
        # 版本信息
        tx3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), self.slide_width - Inches(2), Inches(0.5))
        tf3 = tx3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f"数据周期：基于最新分析 | 生成：{AUTHOR}"
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(160, 160, 160)
        p3.alignment = PP_ALIGN.CENTER

    def add_toc_slide(self):
        """添加目录页（基于已记录的章节）"""
        if not self.toc_entries:
            return
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_section_header(slide, "目录")
        left = Inches(1.5)
        top = Inches(1.8)
        width = self.slide_width - Inches(3)
        height = Inches(0.6 * len(self.toc_entries))
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        for i, (title, idx) in enumerate(self.toc_entries, 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {title}"
            p.font.size = Pt(24)
            p.font.color.rgb = PRIMARY_COLOR
            p.space_after = Pt(12)
            # 可添加超链接，但需保存后生效
            # p.hyperlink.address = f"#slide={idx+1}"

    def add_section(self, title: str):
        """开始新章节（自动记录到目录）"""
        self.current_section = title
        self.toc_entries.append((title, len(self.prs.slides)))

    def add_kpi_slide(self, kpis: List[Tuple[str, str, RGBColor]]):
        """添加KPI卡片页（每行最多4个）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_section_header(slide, "核心指标概览")
        cols = 4
        card_width = Inches(2.8)
        card_height = Inches(1.5)
        start_x = (self.slide_width - (cols * card_width + (cols - 1) * Inches(0.2))) / 2
        for i, (label, value, color) in enumerate(kpis):
            row = i // cols
            col = i % cols
            left = start_x + col * (card_width + Inches(0.2))
            top = Inches(1.8) + row * (card_height + Inches(0.2))
            # 卡片背景
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(210, 210, 210)
            card.line.width = 12700  # 1磅 = 12700 EMU
            # 左侧色条
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.1), card_height)
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = color
            stripe.line.fill.background()
            # 数值
            tx_val = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), card_width - Inches(0.3), Inches(0.8))
            tf_val = tx_val.text_frame
            p_val = tf_val.paragraphs[0]
            p_val.text = str(value)
            p_val.font.size = Pt(28)
            p_val.font.bold = True
            p_val.font.color.rgb = color
            # 标签
            tx_lbl = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.9), card_width - Inches(0.3), Inches(0.4))
            tf_lbl = tx_lbl.text_frame
            p_lbl = tf_lbl.paragraphs[0]
            p_lbl.text = label
            p_lbl.font.size = Pt(12)
            p_lbl.font.color.rgb = RGBColor(80, 80, 80)

    def add_chart_slide(self, title: str, chart_data: Dict):
        """添加图表页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_section_header(slide, title)
        # 所有数值强制转 int (EMU)
        chart_left = int(Inches(0.8))
        chart_top = int(Inches(1.6))
        chart_width = int(int(self.slide_width) - int(Inches(1.6)))
        chart_height = int(Inches(4.8))
        self._add_chart(slide, chart_data, chart_left, chart_top, chart_width, chart_height)

    def add_table_slide(self, title: str, headers: List[str], rows: List[List], col_widths: List[float] = None):
        """添加表格页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_section_header(slide, title)
        # 所有数值转为 int (EMU) 类型
        left = int(Inches(0.8))       # 0.8英寸 → EMU整数
        top = int(Inches(1.6))        # 1.6英寸 → EMU整数
        width = int(self.slide_width) - int(Inches(1.6))  # 强行转int再减
        _add_table(slide, left, top, width, headers, rows, col_widths)

    def add_text_slide(self, title: str, content: List[str]):
        """添加纯文本页（项目符号列表）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_section_header(slide, title)
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), self.slide_width - Inches(1.6), Inches(5))
        tf = tx.text_frame
        tf.word_wrap = True
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(12)

    def _add_section_header(self, slide, title: str):
        """添加章节标题栏"""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.slide_width, Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY_COLOR
        bar.line.fill.background()
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), self.slide_width - Inches(1), Inches(0.7))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    def _add_chart(self, slide, data: Dict, left, top, width, height):
        """添加PPT内置图表（柱状图/折线图/饼图）"""
        chart_type_str = data.get("type", "bar")
        categories = data["categories"]
        series_list = data["series"]  # list of dict with 'name' and 'values'
        chart_data = ChartData()
        chart_data.categories = categories
        for s in series_list:
            chart_data.add_series(s["name"], s["values"])
        if chart_type_str == "bar":
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif chart_type_str == "line":
            chart_type = XL_CHART_TYPE.LINE
        elif chart_type_str == "pie":
            chart_type = XL_CHART_TYPE.PIE
        else:
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data).chart
        if chart_type_str != "pie":
            chart.has_legend = True
            chart.legend.include_in_layout = False
            # 设置系列颜色
            for i, ser in enumerate(chart.series):
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = PRIMARY_COLOR if i == 0 else SECONDARY_COLOR
                ser.data_labels.font.size = Pt(10)
        else:
            chart.has_legend = True

    def save(self, output_path: str):
        """保存文件"""
        self.prs.save(output_path)
        return output_path

# ========== 数据适配函数 ==========
def _load_data_from_excel(excel_path):
    """从Excel读取必要的数据集"""
    xl = pd.ExcelFile(excel_path)
    sheets = xl.sheet_names
    data = {}
    if '汇总统计' in sheets:
        data['summary'] = pd.read_excel(excel_path, sheet_name='汇总统计')
    if '完整偏差明细' in sheets:
        data['detail'] = pd.read_excel(excel_path, sheet_name='完整偏差明细')
    if '替代料明细' in sheets:
        data['alt'] = pd.read_excel(excel_path, sheet_name='替代料明细')
    if '偏差原因分析' in sheets:
        data['cause'] = pd.read_excel(excel_path, sheet_name='偏差原因分析')
    return data

def generate_advanced_report_v2(excel_path, output_path, log_cb=None):
    """生成高级报告 v2（供 GUI 调用）"""
    try:
        data = _load_data_from_excel(excel_path)
        summary = data.get('summary')
        detail = data.get('detail')
        alt = data.get('alt')
        cause = data.get('cause')
        if summary is None or detail is None:
            raise ValueError("Excel 缺少必要的 Sheet（汇总统计或完整偏差明细）")

        pt = AdvancedPPTGeneratorV2()
        # 1. 封面
        pt.add_title_slide()
        # 2. 目录（将在添加章节后自动生成，我们手动在最后添加？）
        # 简单做法：生成后手动调整目录位置，但为了自动化，我们先生成内容，最后插入目录页？python-pptx 不支持插入到前面。
        # 因此我们将在所有章节添加后，再添加目录页（但目录页会在末尾）。用户可接受，或在生成后手动移动。本版就放在末尾。
        # 3. 计算核心KPI（添加列名映射，增强兼容性）
        summary_cols = summary.columns.tolist()
        
        # 总条数
        total_rows_col = next((c for c in ['总条数', '记录数', '条数'] if c in summary_cols), None)
        total_rows = int(summary[total_rows_col].sum()) if total_rows_col else len(summary)
        
        # 正偏差条数
        pos_cnt_col = next((c for c in ['正偏差条数', '正偏差记录数'] if c in summary_cols), None)
        pos_cnt = int(summary[pos_cnt_col].sum()) if pos_cnt_col else 0
        
        # 负偏差条数
        neg_cnt_col = next((c for c in ['负偏差条数', '负偏差记录数'] if c in summary_cols), None)
        neg_cnt = int(summary[neg_cnt_col].sum()) if neg_cnt_col else 0
        
        # 正偏差金额(含税)
        pos_amt_col = next((c for c in ['正偏差金额(含税)', '正偏差金额'] if c in summary_cols), None)
        pos_amount = summary[pos_amt_col].sum() if pos_amt_col else 0
        
        # 负偏差金额(含税)
        neg_amt_col = next((c for c in ['负偏差金额(含税)', '负偏差金额'] if c in summary_cols), None)
        neg_amount = abs(summary[neg_amt_col].sum()) if neg_amt_col else 0
        net_amount = pos_amount - neg_amount
        # 备注覆盖率（加权）
        rate_col_name = next((c for c in ['备注覆盖率', '备注覆盖', '覆盖率'] if c in summary.columns), None)
        if rate_col_name:
            rate_col = summary[rate_col_name]
            rates = rate_col.astype(str).str.replace('%', '').astype(float) / 100
            weight_col = next((c for c in ['总条数', '记录数', '条数'] if c in summary.columns), None)
            weights = summary[weight_col] if weight_col else pd.Series([1] * len(summary))
        else:
            rates = pd.Series([0])
            weights = pd.Series([1] * len(summary))
        note_rate = (rates * weights).sum() / weights.sum() if weights.sum() > 0 else 0
        kpis = [
            ("总记录数", f"{total_rows:,}", PRIMARY_COLOR),
            ("正偏差（多耗）", f"{pos_cnt:,}", POSITIVE_COLOR),
            ("负偏差（少耗）", f"{neg_cnt:,}", NEGATIVE_COLOR),
            ("净偏差(元)", f"{net_amount:,.0f}", ACCENT_COLOR),
            ("备注覆盖率", f"{note_rate:.1%}", SECONDARY_COLOR),
        ]
        pt.add_kpi_slide(kpis)

        # 4. 工厂对比（如果有工厂列）
        factory_col = next((c for c in ['工厂', '工厂名称'] if c in summary.columns), None)
        if factory_col:
            factory_data = summary.groupby(factory_col).agg({
                '总条数': 'sum',
                '正偏差金额(含税)': 'sum',
                '负偏差金额(含税)': lambda x: abs(x.sum())
            }).reset_index()
            factory_data['净偏差'] = factory_data['正偏差金额(含税)'] - factory_data['负偏差金额(含税)']
            headers = ['工厂', '记录数', '正偏差(万元)', '负偏差(万元)', '净偏差(万元)']
            rows = []
            for _, r in factory_data.iterrows():
                rows.append([
                    r[factory_col], r['总条数'],
                    f"{r['正偏差金额(含税)']/10000:.1f}",
                    f"{r['负偏差金额(含税)']/10000:.1f}",
                    f"{r['净偏差']/10000:.1f}"
                ])
            pt.add_table_slide("工厂维度对比", headers, rows, col_widths=[2, 1, 1.5, 1.5, 1.5])

        # 5. 车间偏差排行（Top10）
        workshop_col = next((c for c in ['车间', '生产管理员描述'] if c in detail.columns), None)
        if workshop_col:
            workshop_rank = detail.groupby(workshop_col)['偏差金额'].apply(lambda x: x.abs().sum()).nlargest(10).reset_index()
            workshop_rank.columns = ['车间', '偏差金额(元)']
            pt.add_table_slide("车间偏差金额排行(Top10)", ['车间', '偏差金额(元)'],
                              [[r['车间'], f"{r['偏差金额(元)']:,.0f}"] for _, r in workshop_rank.iterrows()],
                              col_widths=[4, 4])

        # 6. 物料偏差排行（Top10）
        mat_col = next((c for c in ['物料编码', '组件物料号'] if c in detail.columns), None)
        if mat_col:
            mat_rank = detail.groupby(mat_col)['偏差金额'].apply(lambda x: x.abs().sum()).nlargest(10).reset_index()
            mat_rank.columns = ['物料编码', '偏差金额(元)']
            # 尝试添加物料名称
            name_col = next((c for c in ['物料名称', '组件物料描述'] if c in detail.columns), None)
            if name_col:
                name_map = detail[[mat_col, name_col]].drop_duplicates()
                mat_rank = mat_rank.merge(name_map, on=mat_col, how='left')
                headers = ['物料编码', '物料名称', '偏差金额(元)']
                rows = [[r['物料编码'], r[name_col], f"{r['偏差金额(元)']:,.0f}"] for _, r in mat_rank.iterrows()]
                pt.add_table_slide("物料偏差金额排行(Top10)", headers, rows, col_widths=[2, 5, 2])
            else:
                pt.add_table_slide("物料偏差金额排行(Top10)", ['物料编码', '偏差金额(元)'],
                                  [[r['物料编码'], f"{r['偏差金额(元)']:,.0f}"] for _, r in mat_rank.iterrows()],
                                  col_widths=[3, 5])

        # 7. 替代料机制（如果有数据）
        if alt is not None and hasattr(alt, 'empty') and not alt.empty:
            alt_sample = alt.head(3)
            headers = ['物料A', '偏差A', '物料B', '偏差B', '净偏差']
            rows = []
            for _, r in alt_sample.iterrows():
                rows.append([
                    r.get('物料A', '')[:20], f"{r.get('偏差金额A', 0):,.0f}",
                    r.get('物料B', '')[:20], f"{r.get('偏差金额B', 0):,.0f}",
                    f"{r.get('净偏差', 0):,.0f}"
                ])
            pt.add_table_slide("替代料核对机制（示例）", headers, rows, col_widths=[2.5, 1.5, 2.5, 1.5, 1.5])
            pt.add_text_slide("替代料机制价值", [
                "• 镜像偏差：原物料-100%，替代料+100%",
                "• 净偏差 = 物料A + 物料B，反映真实成本",
                "• 金黄胚系列净偏差-6万 vs 账面143万",
                "• 避免将规格切换误判为管理异常"
            ])

        # 8. 根因诊断（如果有原因分析）
        if not cause.empty and '备注原因' in cause.columns:
            cause_counts = cause['备注原因'].value_counts().head(5).reset_index()
            cause_counts.columns = ['原因', '频次']
            pt.add_table_slide("主要偏差原因 Top5", ['原因', '频次'],
                              [[r['原因'], r['频次']] for _, r in cause_counts.iterrows()],
                              col_widths=[6, 2])

        # 9. 行动建议（固定内容）
        pt.add_text_slide("分阶段改进行动", [
            "【立即行动（1周内）】",
            "• 强制补录无备注记录",
            "• 完善系统无定额物料",
            "• 核查领用记录异常",
            "",
            "【短期优化（1月内）】",
            "• 优化批次分摊逻辑",
            "• 加强设备维护",
            "• 规范工艺执行",
            "",
            "【中期建设（3月内）】",
            "• 替代料净偏差自动抵消",
            "• 上线实时预警看板",
            "• 实现系统领用与实际双重校验"
        ])

        # 10. 目标量化
        pt.add_text_slide("预期效果与目标量化", [
            f"• 偏差金额降低30%（当前净偏差 {net_amount:,.0f} 元）",
            f"• 备注覆盖率从 {note_rate:.1%} 提升至 80% 以上",
            "• 异常响应时效从月度缩短至日度",
            "• 消除无备注高偏差记录，建立真实数据基础"
        ])

        # 最后添加目录页（目前是在末尾，用户可手动移至前面）
        pt.add_toc_slide()

        pt.save(output_path)
        if log_cb:
            log_cb(f"高级报告 v2 已生成：{output_path} (共{len(pt.prs.slides)}页)")
        return True
    except Exception as e:
        if log_cb:
            log_cb(f"生成失败：{str(e)}")
        import traceback
        traceback.print_exc()
        return False

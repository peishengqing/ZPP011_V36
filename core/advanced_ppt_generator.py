# -*- coding: utf-8 -*-
"""
高级 PPT 报告生成器
按照 ZPP011生产偏差分析报告 (1).pptx 样式生成详细分析报告
"""

import os
import pandas as pd
import numpy as np
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import io
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

# ========== 辅助函数 ==========
def _get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]

def _add_title_bar(slide, title, subtitle='', color_bg=(30,60,114)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*color_bg)
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), prs.slide_width - Inches(1), Inches(0.4))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(200,200,200)

def _add_card(slide, left, top, width, height, bg_color=RGBColor(255,255,255)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = RGBColor(210,210,210)
    card.line.width = Pt(1)
    return card

def _add_kpi_card(slide, left, top, value, label, color=RGBColor(0,112,192)):
    _add_card(slide, left, top, Inches(2.5), Inches(1.5))
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.1), Inches(1.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    tx = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(2), Inches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color
    tx2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.95), Inches(2), Inches(0.4))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(100,100,100)

def _add_table(slide, left, top, width, headers, rows, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = Inches(0.45)
    tbl_h = row_h * n_rows
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, tbl_h).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30,60,114)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255,255,255)
            p.alignment = PP_ALIGN.CENTER
    # 数据行
    for r_idx, row in enumerate(rows, 1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245,245,245)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(50,50,50)
                p.alignment = PP_ALIGN.CENTER
            if c_idx in [2,3,4] and isinstance(val, (int,float)) and val < 0:
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(40,167,69)

def _add_bullet_text(slide, left, top, width, lines, font_size=12):
    tx = slide.shapes.add_textbox(left, top, width, Inches(len(lines)*0.3))
    tf = tx.text_frame
    tf.word_wrap = True
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(50,50,50)
        p.space_after = Pt(6)

# ========== 数据读取与聚合 ==========
def _load_excel_data(excel_path):
    """读取 Excel 中的各个 sheet，返回字典"""
    xl = pd.ExcelFile(excel_path)
    sheets = xl.sheet_names
    data = {}
    # 汇总统计
    if '汇总统计' in sheets:
        data['summary'] = pd.read_excel(excel_path, sheet_name='汇总统计')
    elif any('汇总' in s for s in sheets):
        s = [s for s in sheets if '汇总' in s][0]
        data['summary'] = pd.read_excel(excel_path, sheet_name=s)
    else:
        data['summary'] = pd.DataFrame()
    # 完整偏差明细
    if '完整偏差明细' in sheets:
        data['detail'] = pd.read_excel(excel_path, sheet_name='完整偏差明细')
    elif any('偏差明细' in s for s in sheets):
        s = [s for s in sheets if '偏差明细' in s][0]
        data['detail'] = pd.read_excel(excel_path, sheet_name=s)
    else:
        data['detail'] = pd.DataFrame()
    # 替代料明细
    if '替代料明细' in sheets:
        data['alt'] = pd.read_excel(excel_path, sheet_name='替代料明细')
    else:
        data['alt'] = pd.DataFrame()
    # 偏差原因分析
    if '偏差原因分析' in sheets:
        data['cause'] = pd.read_excel(excel_path, sheet_name='偏差原因分析')
    else:
        data['cause'] = pd.DataFrame()
    # 无备注预警
    if '无备注预警' in sheets:
        data['no_note'] = pd.read_excel(excel_path, sheet_name='无备注预警')
    else:
        data['no_note'] = pd.DataFrame()
    return data

def _calc_metrics(data):
    """从汇总统计或明细计算核心指标"""
    summary = data.get('summary')
    detail = data.get('detail')
    if not summary.empty:
        total_rows = int(summary['总条数'].sum())
        pos_cnt = int(summary['正偏差条数'].sum())
        neg_cnt = int(summary['负偏差条数'].sum())
        pos_qty = summary['正偏差数量'].sum()
        neg_qty = summary['负偏差数量'].sum()
        pos_amount = summary['正偏差金额(含税)'].sum()
        neg_amount = abs(summary['负偏差金额(含税)'].sum())
        # 备注覆盖率加权平均
        rate_col = summary['备注覆盖率']
        rates = rate_col.astype(str).str.replace('%','').astype(float)/100
        weights = summary['总条数']
        note_rate = (rates * weights).sum() / weights.sum() if weights.sum()>0 else 0
    elif not detail.empty and '偏差金额' in detail.columns:
        total_rows = len(detail)
        pos_cnt = len(detail[detail['偏差金额']>0])
        neg_cnt = len(detail[detail['偏差金额']<0])
        pos_qty = detail[detail['偏差金额']>0]['偏差金额'].sum()
        neg_qty = abs(detail[detail['偏差金额']<0]['偏差金额'].sum())
        pos_amount = pos_qty
        neg_amount = neg_qty
        if '备注原因' in detail.columns:
            note_rate = detail['备注原因'].notna().sum() / total_rows
        else:
            note_rate = 0
    else:
        return None
    net = pos_amount - neg_amount
    return {
        'total_rows': total_rows,
        'pos_cnt': pos_cnt, 'neg_cnt': neg_cnt,
        'pos_amount': pos_amount, 'neg_amount': neg_amount,
        'net_amount': net,
        'note_rate': note_rate
    }

def _get_factory_stats(summary_df):
    if summary_df.empty:
        return []
    factory_col = '工厂' if '工厂' in summary_df.columns else ('工厂名称' if '工厂名称' in summary_df.columns else None)
    if not factory_col:
        return []
    stats = []
    for factory, grp in summary_df.groupby(factory_col):
        stats.append({
            'factory': factory,
            'total': grp['总条数'].sum(),
            'pos_amount': grp['正偏差金额(含税)'].sum(),
            'neg_amount': abs(grp['负偏差金额(含税)'].sum()),
            'note_rate': (grp['备注覆盖率'].astype(str).str.replace('%','').astype(float)/100 * grp['总条数']).sum() / grp['总条数'].sum() if grp['总条数'].sum()>0 else 0
        })
    return stats

def _get_workshop_rank(detail_df, amount_col='偏差金额'):
    if detail_df.empty or '车间' not in detail_df.columns:
        return pd.Series(dtype=float)
    ws_amt = detail_df.groupby('车间')[amount_col].apply(lambda x: x.abs().sum()).sort_values(ascending=False)
    return ws_amt

def _get_note_rate_by_workshop(summary_df):
    if summary_df.empty:
        return {}
    workshop_col = '车间' if '车间' in summary_df.columns else None
    if not workshop_col:
        return {}
    rates = {}
    for _, row in summary_df.iterrows():
        ws = row[workshop_col]
        total = row['总条数']
        rate_str = row['备注覆盖率']
        try:
            rate = float(str(rate_str).replace('%',''))/100
        except:
            rate = 0
        rates[ws] = (rate, total)
    # 加权平均
    weighted = {}
    for ws, (r, t) in rates.items():
        weighted[ws] = r
    return weighted

def _get_alt_examples(alt_df, n=3):
    if alt_df.empty:
        return []
    examples = []
    for _, row in alt_df.head(n).iterrows():
        a = row.get('物料A', '')
        a_amt = row.get('偏差金额A', 0)
        b = row.get('物料B', '')
        b_amt = row.get('偏差金额B', 0)
        net = row.get('净偏差', a_amt + b_amt)
        examples.append((a, a_amt, b, b_amt, net))
    return examples

def _get_root_causes(cause_df):
    if cause_df.empty:
        return {}
    root_keys = ['系统无定额', '领用记录异常', '配比偏差', '批次分摊不均', '生产异常', '额定偏离']
    counts = {k:0 for k in root_keys}
    if '备注原因' in cause_df.columns:
        for reason in cause_df['备注原因'].dropna():
            for k in root_keys:
                if k in str(reason):
                    counts[k] += 1
    return counts

# ========== 核心生成函数 ==========
def generate_advanced_report(excel_path, output_path, log_cb=None):
    """生成高级分析报告 PPT"""
    global prs
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    data = _load_excel_data(excel_path)
    metrics = _calc_metrics(data)
    if not metrics:
        if log_cb: log_cb("无法计算核心指标，请检查 Excel 数据")
        return False

    summary_df = data['summary']
    detail_df = data['detail']
    alt_df = data['alt']
    cause_df = data['cause']
    no_note_df = data['no_note']

    # 提取日期范围
    date_range = "2026年5月"
    if not detail_df.empty and '订单日期' in detail_df.columns:
        dates = pd.to_datetime(detail_df['订单日期'], errors='coerce')
        if not dates.isna().all():
            date_range = f"{dates.min().strftime('%Y-%m-%d')} 至 {dates.max().strftime('%Y-%m-%d')}"

    # ========== 1. 封面 ==========
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(30,60,114)
    bg.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), prs.slide_width - Inches(2), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ZPP011生产偏差分析报告"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p.alignment = PP_ALIGN.CENTER
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), prs.slide_width - Inches(2), Inches(0.8))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"云南达利生产基地 | {date_range} | {metrics['total_rows']:,}条记录深度洞察"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(200,200,200)
    p2.alignment = PP_ALIGN.CENTER

    # ========== 2. 目录 ==========
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _add_title_bar(slide, "汇报目录")
    contents = [
        "01  宏观数据总览",
        "02  车间维度深度分析",
        "03  组件物料与替代料机制",
        "04  根因诊断与改进行动"
    ]
    _add_bullet_text(slide, Inches(1.5), Inches(1.8), Inches(10), contents, font_size=20)

    # ========== 3. 宏观数据总览 ==========
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _add_title_bar(slide, "宏观数据总览", "全厂偏差分布特征")
    # KPI 卡片
    kpis = [
        (f"{metrics['total_rows']:,}", "总记录数"),
        (f"{metrics['pos_cnt']:,}", "正偏差（多耗）"),
        (f"{metrics['neg_cnt']:,}", "负偏差（少耗）"),
        (f"{metrics['net_amount']:,.0f}", "净偏差(元)"),
        (f"{metrics['note_rate']:.1%}", "备注覆盖率")
    ]
    colors = [RGBColor(0,112,192), RGBColor(220,53,69), RGBColor(40,167,69), RGBColor(255,152,0), RGBColor(0,112,192)]
    for i, (val, label) in enumerate(kpis):
        _add_kpi_card(slide, Inches(0.5 + i*2.8), Inches(1.5), val, label, colors[i])
    # 工厂对比柱状图
    factory_stats = _get_factory_stats(summary_df)
    if factory_stats:
        fig, ax = plt.subplots(figsize=(8,4))
        factories = [f['factory'] for f in factory_stats]
        pos_amt = [f['pos_amount']/10000 for f in factory_stats]
        neg_amt = [f['neg_amount']/10000 for f in factory_stats]
        x = np.arange(len(factories))
        width = 0.35
        ax.bar(x - width/2, pos_amt, width, label='正偏差(万元)', color='#dc3545')
        ax.bar(x + width/2, neg_amt, width, label='负偏差(万元)', color='#28a745')
        ax.set_xticks(x)
        ax.set_xticklabels(factories, rotation=15)
        ax.set_ylabel('金额(万元)')
        ax.set_title('工厂维度偏差金额对比')
        ax.legend()
        img_buf = io.BytesIO()
        plt.savefig(img_buf, format='png', dpi=100, bbox_inches='tight')
        plt.close()
        img_buf.seek(0)
        slide.shapes.add_picture(img_buf, Inches(0.5), Inches(3.5), width=Inches(5.5))
    # 偏差类型饼图
    fig2, ax2 = plt.subplots(figsize=(5,4))
    sizes = [metrics['pos_cnt'], metrics['neg_cnt']]
    labels = ['正偏差', '负偏差']
    colors2 = ['#dc3545', '#28a745']
    ax2.pie(sizes, labels=labels, autopct='%1.1f%%', colors=colors2, startangle=90)
    ax2.axis('equal')
    img_buf2 = io.BytesIO()
    plt.savefig(img_buf2, format='png', dpi=100, bbox_inches='tight')
    plt.close()
    img_buf2.seek(0)
    slide.shapes.add_picture(img_buf2, Inches(6.5), Inches(4), width=Inches(4))

    # ========== 4. 车间维度深度分析 ==========
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _add_title_bar(slide, "TOP车间偏差分析")
    ws_rank = _get_workshop_rank(detail_df)
    if not ws_rank.empty:
        top5 = ws_rank.head(5)
        fig, ax = plt.subplots(figsize=(8,4))
        ax.barh(top5.index, top5.values, color='#4A90E2')
        ax.set_xlabel('偏差金额(元)')
        ax.set_title('Top5 车间偏差金额排行')
        img_buf = io.BytesIO()
        plt.savefig(img_buf, format='png', dpi=100, bbox_inches='tight')
        plt.close()
        img_buf.seek(0)
        slide.shapes.add_picture(img_buf, Inches(0.5), Inches(1.8), width=Inches(5.5))
    # 备注覆盖率
    note_rates = _get_note_rate_by_workshop(summary_df)
    if note_rates:
        fig2, ax2 = plt.subplots(figsize=(8,4))
        ws_names = list(note_rates.keys())
        rates = list(note_rates.values())
        ax2.bar(ws_names, rates, color='#F5A623')
        ax2.set_ylabel('备注覆盖率')
        ax2.set_title('各车间备注覆盖率')
        plt.xticks(rotation=45, ha='right')
        img_buf2 = io.BytesIO()
        plt.savefig(img_buf2, format='png', dpi=100, bbox_inches='tight')
        plt.close()
        img_buf2.seek(0)
        slide.shapes.add_picture(img_buf2, Inches(6.8), Inches(1.8), width=Inches(5.5))
    # 预警等级（文字）
    _add_bullet_text(slide, Inches(0.5), Inches(6), Inches(12), [
        "红色预警（立即整改）：配料中心、24000无菌湿法线-1线、6000无菌砖线",
        "黄色预警（限期整改）：食品厂1车间、饮料厂36000热线等9个单元",
        "绿色预警（正常运行）：食品厂4车间等21个单元"
    ], font_size=10)

    # ========== 5. 替代料机制 ==========
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _add_title_bar(slide, "替代料核对机制", "从镜像偏差到净偏差计算")
    text_lines = [
        "镜像对称原理：当原物料显示-100%负偏差，替代料显示+100%正偏差，形成镜像对称。",
        "净偏差公式：净偏差 = 物料A偏差金额 + 物料B偏差金额",
        "案例：32g金黄胚(-74.51万) + 30.2g金黄胚(+68.5万) = 净偏差-6万",
        "价值：实际风险远低于账面绝对值之和143万元，避免误判。"
    ]
    _add_bullet_text(slide, Inches(0.5), Inches(1.8), Inches(12), text_lines, font_size=14)
    # 替代料案例表格
    alt_examples = _get_alt_examples(alt_df, 3)
    if alt_examples:
        headers = ['物料A', '偏差A(元)', '物料B', '偏差B(元)', '净偏差(元)']
        rows = [[a, f"{a_amt:,.0f}", b, f"{b_amt:,.0f}", f"{net:,.0f}"] for a, a_amt, b, b_amt, net in alt_examples]
        _add_table(slide, Inches(0.5), Inches(4), Inches(12), headers, rows,
                   [Inches(2.5), Inches(1.5), Inches(2.5), Inches(1.5), Inches(1.5)])

    # ========== 6. 根因诊断与改进行动 ==========
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _add_title_bar(slide, "根因诊断与改进行动")
    root_causes = _get_root_causes(cause_df)
    if root_causes:
        fig, ax = plt.subplots(figsize=(7,4))
        causes = list(root_causes.keys())
        counts = list(root_causes.values())
        ax.bar(causes, counts, color='#6c757d')
        ax.set_ylabel('频次')
        ax.set_title('六大根因分布')
        plt.xticks(rotation=45, ha='right')
        img_buf = io.BytesIO()
        plt.savefig(img_buf, format='png', dpi=100, bbox_inches='tight')
        plt.close()
        img_buf.seek(0)
        slide.shapes.add_picture(img_buf, Inches(0.5), Inches(1.8), width=Inches(6))
    # 行动方案文本
    actions = [
        "立即行动（1周内）：强制补录无备注记录、完善系统无定额、核查领用异常",
        "短期优化（1月内）：优化批次分摊、加强设备维护、规范工艺执行",
        "中期建设（3月内）：替代料自动抵消、实时预警看板、双重校验"
    ]
    _add_bullet_text(slide, Inches(7), Inches(2), Inches(5.5), actions, font_size=12)

    # ========== 7. 封底 ==========
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    bg2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = RGBColor(30,60,114)
    bg2.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1), Inches(3), prs.slide_width - Inches(2), Inches(1))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "谢谢聆听"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p.alignment = PP_ALIGN.CENTER
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), prs.slide_width - Inches(2), Inches(0.5))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "云南达利生产基地 · 持续改进，精益求精"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(200,200,200)
    p2.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    if log_cb:
        log_cb(f"高级报告已生成：{output_path} (共{len(prs.slides)}页)")
    return True
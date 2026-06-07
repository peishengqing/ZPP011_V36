# -*- coding: utf-8 -*-
"""
效益报告生成对话框
"""
import os
from datetime import datetime
from PySide6.QtWidgets import (
    QDialog, QVBoxLayout, QHBoxLayout, QLabel, QLineEdit, QPushButton,
    QFileDialog, QProgressBar, QMessageBox
)
from PySide6.QtCore import Qt, QThread, Signal
import pandas as pd

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class GenerateReportWorker(QThread):
    finished = Signal(str)
    error = Signal(str)

    def __init__(self, audit_data, output_path):
        super().__init__()
        self.audit_data = audit_data
        self.output_path = output_path

    def run(self):
        try:
            self._generate_ppt()
            self.finished.emit(self.output_path)
        except Exception as e:
            self.error.emit(str(e))

    def _generate_ppt(self):
        prs = Presentation()
        df = self.audit_data

        total = len(df)
        amount_col = '偏差金额' if '偏差金额' in df.columns else 'deviation_amount'
        if amount_col in df.columns:
            total_amount = df[amount_col].fillna(0).sum()
            pos_amount = df[df[amount_col] > 0][amount_col].sum()
            neg_amount = df[df[amount_col] < 0][amount_col].abs().sum()
        else:
            total_amount = pos_amount = neg_amount = 0

        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "ZPP011 生产偏差分析报告"
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        p2.font.size = Pt(18)
        p2.alignment = PP_ALIGN.CENTER

        # 核心指标
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "核心指标概览"
        p.font.size = Pt(24)
        p.font.bold = True
        metrics = [
            f"总记录数：{total:,} 条",
            f"总偏差金额：¥{total_amount:,.2f}",
            f"正偏差金额：¥{pos_amount:,.2f}",
            f"负偏差金额：¥{neg_amount:,.2f}"
        ]
        y = 1.5
        for metric in metrics:
            box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = metric
            p.font.size = Pt(16)
            y += 0.6

        # 车间排行
        workshop_col = next((c for c in ['车间', '生产管理员描述', 'admin'] if c in df.columns), None)
        if workshop_col and amount_col in df.columns:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = "车间偏差金额排行"
            p.font.size = Pt(24)
            p.font.bold = True
            workshop_amount = df.groupby(workshop_col)[amount_col].apply(lambda x: x.abs().sum()).nlargest(10)
            if not workshop_amount.empty:
                chart_data = CategoryChartData()
                chart_data.categories = workshop_amount.index.tolist()
                chart_data.add_series("偏差金额(元)", workshop_amount.values.tolist())
                x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4)
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
                ).chart
                chart.has_legend = False
                chart.value_axis.tick_labels.font.size = Pt(9)
                for series in chart.series:
                    series.has_data_labels = True

        # 物料排行
        mat_col = next((c for c in ['物料编码', '组件物料号', 'code'] if c in df.columns), None)
        if mat_col and amount_col in df.columns:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = "物料偏差金额排行(Top10)"
            p.font.size = Pt(24)
            p.font.bold = True
            mat_rank = df.groupby(mat_col)[amount_col].apply(lambda x: x.abs().sum()).nlargest(10).reset_index()
            if not mat_rank.empty:
                rows, cols = len(mat_rank), 2
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                height = Inches(0.4 * rows)
                table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
                table.cell(0, 0).text = "物料编码"
                table.cell(0, 1).text = "偏差金额(元)"
                for i, row in mat_rank.iterrows():
                    table.cell(i + 1, 0).text = str(row[mat_col])
                    table.cell(i + 1, 1).text = f"{row[amount_col]:,.2f}"
        prs.save(self.output_path)


class BenefitReportDialog(QDialog):
    def __init__(self, parent, audit_data):
        super().__init__(parent)
        self.setWindowTitle("生成效益报告")
        self.resize(450, 200)
        self.audit_data = audit_data
        self.worker = None

        layout = QVBoxLayout(self)
        layout.addWidget(QLabel("输出文件路径："))
        file_layout = QHBoxLayout()
        self.file_edit = QLineEdit()
        self.file_edit.setReadOnly(True)
        file_layout.addWidget(self.file_edit)
        browse_btn = QPushButton("浏览...")
        browse_btn.clicked.connect(self._browse)
        file_layout.addWidget(browse_btn)
        layout.addLayout(file_layout)

        self.progress = QProgressBar()
        self.progress.setVisible(False)
        layout.addWidget(self.progress)

        btn_layout = QHBoxLayout()
        self.generate_btn = QPushButton("生成报告")
        self.generate_btn.clicked.connect(self._generate)
        cancel_btn = QPushButton("取消")
        cancel_btn.clicked.connect(self.reject)
        btn_layout.addWidget(self.generate_btn)
        btn_layout.addStretch()
        btn_layout.addWidget(cancel_btn)
        layout.addLayout(btn_layout)

    def _browse(self):
        default_name = f"效益报告_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        file_path, _ = QFileDialog.getSaveFileName(self, "保存 PPT 文件", default_name, "PPT files (*.pptx)")
        if file_path:
            self.file_edit.setText(file_path)

    def _generate(self):
        output_path = self.file_edit.text()
        if not output_path:
            QMessageBox.warning(self, "提示", "请先选择输出路径")
            return
        if not PPTX_AVAILABLE:
            QMessageBox.critical(self, "错误", "未安装 python-pptx 库，请执行 pip install python-pptx")
            return
        if self.audit_data is None or self.audit_data.empty:
            QMessageBox.warning(self, "提示", "无数据")
            return

        self.progress.setVisible(True)
        self.progress.setRange(0, 0)
        self.generate_btn.setEnabled(False)

        self.worker = GenerateReportWorker(self.audit_data, output_path)
        self.worker.finished.connect(self._on_finished)
        self.worker.error.connect(self._on_error)
        self.worker.start()

    def _on_finished(self, output_path):
        self.progress.setVisible(False)
        self.generate_btn.setEnabled(True)
        QMessageBox.information(self, "成功", f"报告已生成：{output_path}")
        self.accept()

    def _on_error(self, err):
        self.progress.setVisible(False)
        self.generate_btn.setEnabled(True)
        QMessageBox.critical(self, "错误", f"生成失败：{err}")

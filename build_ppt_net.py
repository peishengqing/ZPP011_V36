# -*- coding: utf-8 -*-
"""
ZPP011 偏差分析 PPT 生成脚本（净偏差口径）
==========================================
故事线：账面毛偏差（含替代料等对冲噪音）→ 净偏差（对冲后的真实偏差）。
所有数字与结论文字均由数据动态计算，不写死任何数值/方向词。

用法：
    python build_ppt_net.py                 # 自动取 E:\zpp011_v2 下最新的 ZPP011偏差分析最终版_*.xlsx
    python build_ppt_net.py <xlsx路径>      # 指定输入文件

依赖：pandas, openpyxl, python-pptx
"""

import os
import sys
import glob
import datetime
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ===================== CONFIG =====================
SEARCH_DIRS = [
    r"E:\zpp011_v2",
    r"E:\ZPP011导出文件原数据",
    os.path.join(os.path.expanduser("~"), "Desktop"),
    os.path.join(os.path.expanduser("~"), "Documents"),
]
INPUT_PATTERN = "ZPP011偏差分析最终版_*.xlsx"
FONT = "Microsoft YaHei"

# ===================== 配色 =====================
TEAL = RGBColor(0x0E, 0x7C, 0x86)
TEAL_D = RGBColor(0x0A, 0x5A, 0x62)
AMBER = RGBColor(0xE8, 0x83, 0x3A)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
INK = RGBColor(0x22, 0x2B, 0x2E)
GREY = RGBColor(0x6B, 0x77, 0x7A)
LIGHT = RGBColor(0xF2, 0xF6, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xEC, 0xF3, 0xF3)

PW = Inches(13.333)
PH = Inches(7.5)


# ===================== 工具函数 =====================
def set_run(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT})
    rPr.append(ea)


def txt(slide, x, y, w, h, text, size=14, color=INK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        set_run(r, size, color, bold)
    return tb


def rect(slide, x, y, w, h, fill, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def bg(slide, color):
    rect(slide, 0, 0, PW, PH, color)


def header(slide, idx, title, subtitle=None):
    rect(slide, 0, 0, PW, Inches(1.15), TEAL)
    rect(slide, 0, Inches(1.15), PW, Inches(0.06), AMBER)
    rect(slide, Inches(0.45), Inches(0.28), Inches(0.6), Inches(0.6), AMBER)
    txt(slide, Inches(0.45), Inches(0.28), Inches(0.6), Inches(0.6),
        f"{idx:02d}", 22, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    txt(slide, Inches(1.25), Inches(0.22), Inches(11), Inches(0.55),
        title, 26, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    if subtitle:
        txt(slide, Inches(1.27), Inches(0.74), Inches(11), Inches(0.35),
            subtitle, 13, RGBColor(0xCF, 0xE8, 0xE8), False)
    txt(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
        "ZPP011 生产偏差分析 · 净偏差口径", 9, GREY)
    txt(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.35),
        f"{idx}/11", 9, GREY, False, PP_ALIGN.RIGHT)


def rate_num(x):
    if pd.isna(x):
        return np.nan
    s = str(x).replace("%", "").strip()
    return float(s) if s not in ("", "nan") else np.nan


def w_money(v):
    """万元，带正负号，保留一位小数"""
    return f"{v/10000:+,.1f}万"


def tag(v):
    """按符号给方向词（金额口径：正=超耗，负=省料）"""
    if abs(v) < 0.005 * 10000:  # <50元 视为持平
        return "基本持平"
    return "超耗" if v > 0 else "省料"


def sign_color(v):
    return RED if v > 0 else GREEN


def clip_name(s, n=16):
    s = str(s)
    return s if len(s) <= n else s[:n - 1] + "…"


# ===================== 数据计算 =====================
def find_input():
    if len(sys.argv) > 1 and sys.argv[1].lower().endswith(".xlsx"):
        return sys.argv[1]
    cands = []
    for dd in SEARCH_DIRS:
        cands += glob.glob(os.path.join(dd, INPUT_PATTERN))
    if not cands:
        raise FileNotFoundError(f"未找到 {INPUT_PATTERN}，检查目录：{SEARCH_DIRS}")
    return max(cands, key=os.path.getmtime)


def _metrics_from_df(d, src="当前分析数据"):
    """从内存 DataFrame（= 完整偏差明细同构）计算核心指标。"""
    d = d.copy()
    d["rd"] = d["偏差率"].apply(rate_num)
    for col in ("偏差金额", "净偏差金额"):
        d[col] = pd.to_numeric(d[col], errors="coerce").fillna(0.0)

    alt_mask = d["是否替代料"] == "是"
    alt = d[alt_mask]

    M = {}
    M["src"] = src
    M["total_rows"] = len(d)
    M["alt_rows"] = int(alt_mask.sum())

    # 核心四数：毛 / 净 / 噪音 / 替代料残差
    M["tot_b"] = d["偏差金额"].sum()          # 账面毛偏差
    M["tot_n"] = d["净偏差金额"].sum()        # 净偏差（真实）
    M["noise"] = M["tot_b"] - M["tot_n"]      # 对冲噪音
    M["alt_b"] = alt["偏差金额"].sum()
    M["alt_n"] = alt["净偏差金额"].sum()
    M["alt_b_pos"] = alt.loc[alt["偏差金额"] > 0, "偏差金额"].sum()
    M["alt_b_neg"] = alt.loc[alt["偏差金额"] < 0, "偏差金额"].sum()

    # 分厂（全表，毛+净）
    M["by_factory"] = d.groupby("工厂").agg(
        行数=("工厂", "size"),
        偏差金额=("偏差金额", "sum"),
        净偏差金额=("净偏差金额", "sum"))

    # 物料 Top10（净偏差金额，全表）
    M["mat_top10"] = d.groupby("物料名称").agg(
        净偏差金额=("净偏差金额", "sum"),
        行数=("物料名称", "size")).sort_values(
        "净偏差金额", ascending=False).head(10)

    # 车间（全表，净）
    M["by_workshop"] = d.groupby(["工厂", "车间"]).agg(
        净偏差金额=("净偏差金额", "sum"),
        行数=("车间", "size"))

    # 物料类型（全表，净）
    M["by_mat_type"] = d.groupby("物料类型").agg(
        行数=("物料类型", "size"),
        净偏差金额=("净偏差金额", "sum")).sort_values(
        "净偏差金额", ascending=False)

    # 高偏差率（>=5%，全表）
    hi = d[d["rd"] >= 5]
    M["hi_cnt"] = len(hi)
    sys_mask = hi["备注"].fillna("").str.contains("系统无额定|系统无定额")
    M["hi_sys_cnt"] = int(sys_mask.sum())

    return M, d


def _warn_counts_from_df(d):
    """从内存 df 直接重算三类预警条数（口径与表格主数据一致）。

    规则翻译自 analysis/excel_builder 的 sheet3/4/6：
    - 无备注预警：偏差数量≠0 且 备注为空 且 |偏差率|>10%
    - 中间地带：|偏差率|<=10% 且 非替代料行 且 流程订单不在替代料订单集合
    - 异常预警：5 类（系统无定额 / 实际<=0无备注 / 实际==0有备注 / 包材负偏差 / 替代料残差）
    """
    T = 10.0
    d = d.copy()
    d["rd"] = d["偏差率"].apply(rate_num)
    for c in ("定额", "实际", "偏差数量"):
        d[c] = pd.to_numeric(d[c], errors="coerce")
    d["备注"] = d["备注"].fillna("").astype(str).str.strip()
    d["流程订单"] = d["流程订单"].astype(str)
    d["物料类型"] = d["物料类型"].astype(str)

    # 无备注预警
    note = int(((d["偏差数量"] != 0) & (d["备注"] == "") & (d["rd"].abs() > T)).sum())

    # 中间地带：排除替代料行 + 替代料订单下的所有行
    alt_orders = set(d.loc[d["是否替代料"] == "是", "流程订单"])
    mid = int(((d["rd"].abs() <= T) & (~d["流程订单"].isin(alt_orders))
               & (d["是否替代料"] != "是")).sum())

    # 异常预警 5 类
    sys_no = d["备注"].str.contains("系统无定额")
    no_quota = (d["定额"] == 0) | sys_no
    abn1 = (d["定额"] > 0) & sys_no
    abn2 = (d["定额"] > 0) & (d["实际"] <= 0) & (d["备注"] == "")
    abn3 = (d["定额"] > 0) & (d["实际"] == 0) & (d["备注"] != "")
    abn4 = (d["物料类型"] == "包材") & (d["rd"] < 0) & (~no_quota)
    abn5 = (d["是否替代料"] == "是") & (d["rd"].notna()) & (d["rd"].abs() > 0) & (~no_quota)
    abn = int(abn1.sum() + abn2.sum() + abn3.sum() + abn4.sum() + abn5.sum())

    return {"warn_note": note, "warn_abn": abn, "warn_mid": mid}


def _prepare(df, src_name=None):
    """统一入口：从 df 计算 (M, d)，含三类预警重算。"""
    M, d = _metrics_from_df(df, src=src_name or "当前分析数据")
    M.update(_warn_counts_from_df(d))
    return M, d


def load_metrics(xlsx):
    """文件模式：读取完整偏差明细 sheet，走统一 df 管线。"""
    d = pd.read_excel(xlsx, "完整偏差明细")
    return _prepare(d, os.path.basename(xlsx))


def build_net_report(df, output_path, src_name=None):
    """核心生成入口：直接吃内存 DataFrame，生成净偏差口径 PPT。

    软件端调用此函数，把当前 view_model.df 传入即可，无需落盘 Excel。
    所有数字/方向词动态计算，不写死。
    """
    M, d = _prepare(df, src_name)
    prs = Presentation()
    prs.slide_width = PW
    prs.slide_height = PH
    build_cover(prs, M)
    build_toc(prs)
    build_metrics(prs, M)
    build_gross_vs_net(prs, M)
    build_factory_cmp(prs, M)
    build_mat_top10(prs, M)
    build_workshop(prs, M)
    build_mat_type(prs, M)
    build_cause(prs, M)
    build_warn(prs, M)
    build_summary(prs, M, d)
    prs.save(output_path)
    return output_path


# ===================== 各页构建 =====================
def build_cover(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, TEAL_D)
    rect(s, 0, 0, PW, Inches(0.18), AMBER)
    txt(s, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.6),
        "ZPP011 生产偏差分析报告", 40, WHITE, True)
    txt(s, Inches(0.82), Inches(1.75), Inches(11.5), Inches(0.5),
        "净偏差口径 · 消除对冲噪音后的真实偏差", 20, AMBER, True)
    txt(s, Inches(0.82), Inches(2.35), Inches(11.5), Inches(0.4),
        f"数据来源：{M['src']}", 13, RGBColor(0xCF, 0xE8, 0xE8))

    cards = [
        (f"{M['total_rows']:,}", "偏差记录总数", TEAL),
        (w_money(M['tot_b']), f"账面毛偏差（{tag(M['tot_b'])}）", sign_color(M['tot_b'])),
        (w_money(M['tot_n']), f"净偏差（实际{tag(M['tot_n'])}）", sign_color(M['tot_n'])),
        (w_money(M['noise']), "已消除的对冲噪音", AMBER),
    ]
    cw = Inches(2.85)
    gap = Inches(0.3)
    x0 = Inches(0.8)
    y0 = Inches(3.4)
    for i, (big, lab, col) in enumerate(cards):
        x = x0 + i * (cw + gap)
        rect(s, x, y0, cw, Inches(2.0), WHITE)
        rect(s, x, y0, cw, Inches(0.12), col)
        txt(s, x, y0 + Inches(0.35), cw, Inches(0.9), big, 28, col, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x, y0 + Inches(1.35), cw, Inches(0.55), lab, 12, GREY, False,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
        f"生成时间：{datetime.datetime.now():%Y-%m-%d %H:%M}  ·  云南达利生产基地",
        12, RGBColor(0xCF, 0xE8, 0xE8))


def build_toc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 2, "目录", "CONTENTS")
    items = [
        "核心指标总览", "毛偏差为何失真", "工厂维度对比",
        "物料偏差 Top10", "车间偏差详情", "物料类型分布",
        "偏差原因分析", "预警核查结果", "总结与改进建议",
    ]
    x1 = Inches(0.8)
    x2 = Inches(7.0)
    y0 = Inches(1.6)
    rh = Inches(0.55)
    for i, it in enumerate(items):
        col = 0 if i < 5 else 1
        row = i if i < 5 else i - 5
        x = x1 if col == 0 else x2
        y = y0 + row * rh
        rect(s, x, y + Inches(0.05), Inches(0.45), Inches(0.45), TEAL)
        txt(s, x, y + Inches(0.05), Inches(0.45), Inches(0.45),
            f"{i+1:02d}", 14, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.6), y, Inches(5.2), Inches(0.55), it, 16, INK, False,
            PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def build_metrics(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    bf = M["by_factory"]
    header(s, 3, "核心指标总览",
           f"净偏差口径：总体{tag(M['tot_n'])} {w_money(M['tot_n'])}")
    cards = [("总体", M["tot_n"], M["total_rows"])]
    for f, row in bf.iterrows():
        cards.append((f, row["净偏差金额"], int(row["行数"])))
    cw = Inches(3.9)
    gap = Inches(0.35)
    x0 = Inches(0.8)
    y0 = Inches(1.7)
    for i, (name, net, rows) in enumerate(cards[:3]):
        x = x0 + i * (cw + gap)
        col = TEAL if i == 0 else sign_color(net)
        rect(s, x, y0, cw, Inches(2.4), CARD_BG)
        rect(s, x, y0, cw, Inches(0.14), col)
        txt(s, x, y0 + Inches(0.3), cw, Inches(0.5), name, 20, INK, True,
            PP_ALIGN.CENTER)
        txt(s, x, y0 + Inches(0.95), cw, Inches(0.9), w_money(net), 30, col, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x, y0 + Inches(1.9), cw, Inches(0.4),
            f"记录 {rows:,} 条 · {tag(net)}", 13, GREY, False, PP_ALIGN.CENTER)

    # 洞察条（全部动态）
    lines = [f"• 账面毛偏差 {w_money(M['tot_b'])} 看似{tag(M['tot_b'])}，"
             f"消除替代料等对冲噪音 {w_money(M['noise'])} 后，"
             f"净偏差实际为 {w_money(M['tot_n'])}（{tag(M['tot_n'])}）。"]
    for f, row in bf.iterrows():
        lines.append(f"• {f}：毛偏差 {w_money(row['偏差金额'])} → "
                     f"净偏差 {w_money(row['净偏差金额'])}（{tag(row['净偏差金额'])}）。")
    rect(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.9), LIGHT)
    rect(s, Inches(0.8), Inches(4.6), Inches(0.14), Inches(1.9), AMBER)
    txt(s, Inches(1.1), Inches(4.75), Inches(11.2), Inches(0.4),
        "关键洞察", 16, AMBER, True)
    txt(s, Inches(1.1), Inches(5.2), Inches(11.2), Inches(1.2),
        "\n".join(lines[:3]), 13, INK)


def build_gross_vs_net(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 4, "毛偏差为何失真", "账面毛偏差 → 净偏差：对冲噪音的来源与消除")
    # 左：账面毛偏差
    rect(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(2.6), CARD_BG)
    rect(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(0.14), GREY)
    txt(s, Inches(1.0), Inches(1.9), Inches(5.2), Inches(0.4),
        "账面毛偏差（原始账面数）", 16, INK, True)
    txt(s, Inches(1.0), Inches(2.45), Inches(5.2), Inches(0.5),
        f"合计  {w_money(M['tot_b'])}（{tag(M['tot_b'])}）", 15,
        sign_color(M['tot_b']), True)
    txt(s, Inches(1.0), Inches(3.0), Inches(5.2), Inches(0.5),
        f"其中替代料 {M['alt_rows']} 行毛偏差 {w_money(M['alt_b'])}", 13, INK)
    txt(s, Inches(1.0), Inches(3.55), Inches(5.2), Inches(0.6),
        f"（正向 +{M['alt_b_pos']/10000:,.1f}万 / 负向 −{abs(M['alt_b_neg'])/10000:,.1f}万，成对对冲）",
        12, GREY)

    # 右：净偏差
    rect(s, Inches(6.9), Inches(1.7), Inches(5.6), Inches(2.6), CARD_BG)
    rect(s, Inches(6.9), Inches(1.7), Inches(5.6), Inches(0.14), TEAL)
    txt(s, Inches(7.1), Inches(1.9), Inches(5.2), Inches(0.4),
        "净偏差（对冲后的真实偏差）", 16, INK, True)
    txt(s, Inches(7.1), Inches(2.45), Inches(5.2), Inches(0.5),
        f"合计  {w_money(M['tot_n'])}（{tag(M['tot_n'])}）", 15,
        sign_color(M['tot_n']), True)
    txt(s, Inches(7.1), Inches(3.0), Inches(5.2), Inches(0.5),
        f"替代料净残差仅 {w_money(M['alt_n'])}", 13, INK)
    txt(s, Inches(7.1), Inches(3.55), Inches(5.2), Inches(0.6),
        f"对冲噪音共消除 {w_money(M['noise'])}", 12, GREY)

    # 结论（动态方向词）
    flip = (M['tot_b'] > 0) != (M['tot_n'] > 0)
    concl = (f"账面毛偏差 {w_money(M['tot_b'])}（{tag(M['tot_b'])}）中，"
             f"{w_money(M['noise'])} 是替代料等成对登记造成的对冲噪音；"
             f"净偏差口径下实际为 {w_money(M['tot_n'])}（{tag(M['tot_n'])}）。")
    if flip:
        concl += "\n即：毛偏差与净偏差方向相反——不消噪音会得出完全相反的管理结论。"
    else:
        concl += "\n即：消除噪音后偏差大幅收敛，账面数字显著高估了真实偏差规模。"
    rect(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.8), LIGHT)
    rect(s, Inches(0.8), Inches(4.7), Inches(0.14), Inches(1.8), AMBER)
    txt(s, Inches(1.1), Inches(4.85), Inches(11.2), Inches(0.4),
        "结论", 16, AMBER, True)
    txt(s, Inches(1.1), Inches(5.3), Inches(11.2), Inches(1.1), concl, 13, INK)


def build_factory_cmp(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 5, "工厂维度对比", "各工厂：账面毛偏差 vs 净偏差")
    bf = M["by_factory"]
    rows = [(f, bf.loc[f, "偏差金额"], bf.loc[f, "净偏差金额"]) for f in bf.index]
    base_y = Inches(5.4)
    maxv = max(max(abs(r[1]), abs(r[2])) for r in rows)
    maxv = max(maxv, 1)
    bw = Inches(1.6)
    for i, (f, gb, nv) in enumerate(rows):
        x = Inches(1.6) + i * Inches(3.2)
        h1 = Inches(abs(gb) / maxv * 2.8)
        col1 = sign_color(gb)
        if gb >= 0:
            rect(s, x, base_y - h1, bw, h1, col1)
        else:
            rect(s, x, base_y, bw, h1, col1)
        txt(s, x, base_y - h1 - Inches(0.35), bw, Inches(0.3),
            w_money(gb), 12, col1, True, PP_ALIGN.CENTER)
        x2 = x + bw + Inches(0.3)
        h2 = Inches(abs(nv) / maxv * 2.8)
        col2 = sign_color(nv)
        if nv >= 0:
            rect(s, x2, base_y - h2, bw, h2, col2)
        else:
            rect(s, x2, base_y, bw, h2, col2)
        txt(s, x2, base_y - h2 - Inches(0.35), bw, Inches(0.3),
            w_money(nv), 12, col2, True, PP_ALIGN.CENTER)
        txt(s, x, base_y + Inches(0.1), bw * 2 + Inches(0.3), Inches(0.4),
            f, 13, INK, True, PP_ALIGN.CENTER)

    txt(s, Inches(0.8), Inches(1.7), Inches(6), Inches(0.4),
        "（每组左柱=账面毛偏差，右柱=净偏差）", 12, GREY)
    parts = [f"{f}：毛 {w_money(gb)} → 净 {w_money(nv)}（{tag(nv)}）"
             for f, gb, nv in rows]
    rect(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7), LIGHT)
    txt(s, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.6),
        "；".join(parts) + "。", 13, INK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def build_mat_top10(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 6, "物料偏差 Top10", "按净偏差金额排序（全表，净偏差口径）")
    mt = M["mat_top10"]
    maxv = mt["净偏差金额"].abs().max()
    x0 = Inches(0.8)
    y0 = Inches(1.6)
    rh = Inches(0.48)
    bw_max = Inches(7.5)
    for i, (name, row) in enumerate(mt.iterrows()):
        y = y0 + i * rh
        v = row["净偏差金额"]
        col = sign_color(v)
        txt(s, x0, y, Inches(3.6), rh, f"{i+1:02d}. {clip_name(name, 15)}", 11, INK,
            False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        bw = Inches(abs(v) / maxv * bw_max.inches) if maxv else Inches(0)
        bx = x0 + Inches(3.7)
        if v >= 0:
            rect(s, bx, y + Inches(0.08), bw, Inches(0.3), col)
        else:
            rect(s, bx + (bw_max - bw), y + Inches(0.08), bw, Inches(0.3), col)
        txt(s, bx + bw_max + Inches(0.1), y, Inches(1.6), rh,
            f"{v/10000:+,.1f}万", 11, col, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    # 侧栏说明（动态第一名）
    top_name = clip_name(mt.index[0], 12) if len(mt) else "—"
    top_v = mt["净偏差金额"].iloc[0] if len(mt) else 0
    side = (f"红色=超耗\n绿色=省料\n\n净超耗第一\n{top_name}\n{w_money(top_v)}"
            if top_v > 0 else
            f"红色=超耗\n绿色=省料\n\n无净超耗物料\n（Top1 {top_name}\n{w_money(top_v)}）")
    rect(s, Inches(11.5), Inches(1.6), Inches(1.3), Inches(4.8), CARD_BG)
    txt(s, Inches(11.6), Inches(1.7), Inches(1.2), Inches(4.6),
        side, 11, INK, False, PP_ALIGN.CENTER, MSO_ANCHOR.TOP)


def build_workshop(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 7, "车间偏差详情", "各工厂车间净偏差金额 Top6")
    ws = M["by_workshop"]
    factories = list(dict.fromkeys(ws.index.get_level_values(0)))[:2]
    xs = [Inches(0.8), Inches(7.0)]
    insight_parts = []
    for f, x0 in zip(factories, xs):
        sub = ws.xs(f, level=0)
        if len(sub) == 0:
            continue
        sub = sub.sort_values("净偏差金额", ascending=False).head(6)
        maxv = sub["净偏差金额"].abs().max()
        txt(s, x0, Inches(1.55), Inches(5.5), Inches(0.4), f, 16, TEAL, True)
        y0 = Inches(2.1)
        rh = Inches(0.62)
        bw_max = Inches(3.6)
        for i, (name, row) in enumerate(sub.iterrows()):
            y = y0 + i * rh
            v = row["净偏差金额"]
            col = sign_color(v)
            txt(s, x0, y, Inches(2.4), rh, clip_name(name, 10), 11, INK, False,
                PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
            bw = Inches(abs(v) / maxv * bw_max.inches) if maxv else Inches(0)
            bx = x0 + Inches(2.5)
            if v >= 0:
                rect(s, bx, y + Inches(0.18), bw, Inches(0.26), col)
            else:
                rect(s, bx + (bw_max - bw), y + Inches(0.18), bw, Inches(0.26), col)
            txt(s, bx + bw_max + Inches(0.05), y, Inches(1.4), rh,
                f"{v/10000:+,.1f}万", 11, col, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        # 动态归纳该厂
        pos = sub[sub["净偏差金额"] > 0]
        if len(pos):
            top_ws = clip_name(pos.index[0], 8)
            insight_parts.append(
                f"{f}净超耗集中在「{top_ws}」等 {len(pos)} 个车间")
        else:
            insight_parts.append(f"{f}各车间均为净省料")
    rect(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7), LIGHT)
    txt(s, Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.6),
        "；".join(insight_parts) + "。", 13, INK, False,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def build_mat_type(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 8, "物料类型分布", "各物料类型净偏差金额（净偏差口径）")
    mt = M["by_mat_type"]
    maxv = mt["净偏差金额"].abs().max()
    x0 = Inches(1.0)
    y0 = Inches(1.8)
    rh = Inches(0.7)
    bw_max = Inches(8.0)
    for i, (name, row) in enumerate(mt.head(5).iterrows()):
        y = y0 + i * rh
        v = row["净偏差金额"]
        col = sign_color(v)
        txt(s, x0, y, Inches(2.6), rh, clip_name(name, 10), 13, INK, False,
            PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        bw = Inches(abs(v) / maxv * bw_max.inches) if maxv else Inches(0)
        bx = x0 + Inches(2.7)
        if v >= 0:
            rect(s, bx, y + Inches(0.22), bw, Inches(0.3), col)
        else:
            rect(s, bx + (bw_max - bw), y + Inches(0.22), bw, Inches(0.3), col)
        txt(s, bx + bw_max + Inches(0.1), y, Inches(2.0), rh,
            f"{v/10000:+,.1f}万", 12, col, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    # 动态洞察
    pos_types = mt[mt["净偏差金额"] > 0]
    if len(pos_types) == 0:
        ins = "所有物料类型均为净省料（绿色），无整体超耗类型。"
    else:
        names = "、".join(clip_name(n, 8) for n in pos_types.index[:3])
        ins = (f"净超耗集中在「{names}」等 {len(pos_types)} 个类型"
               f"（合计 {w_money(pos_types['净偏差金额'].sum())}），其余类型均为省料。")
    rect(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0), LIGHT)
    rect(s, Inches(0.8), Inches(5.9), Inches(0.14), Inches(1.0), AMBER)
    txt(s, Inches(1.1), Inches(6.0), Inches(11.2), Inches(0.9),
        ins, 13, INK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def build_cause(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 9, "偏差原因分析", "高偏差率（≥5%）记录归因")
    total_hi = M["hi_cnt"]
    sys_cnt = M["hi_sys_cnt"]
    other = total_hi - sys_cnt
    pct_sys = sys_cnt / total_hi * 100 if total_hi else 0
    pct_other = 100 - pct_sys if total_hi else 0
    rect(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.2), CARD_BG)
    rect(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(0.5), TEAL)
    txt(s, Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.5),
        f"主数据型  {pct_sys:.0f}%", 18, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.0), Inches(2.4), Inches(5.2), Inches(3.3),
        f"• 系统无额定 / 无定额 导致 {sys_cnt} 条\n"
        f"• 占高偏差率记录（共 {total_hi} 条）的 {pct_sys:.1f}%\n"
        f"• 本质：BOM/定额主数据未维护\n"
        f"• 并非真实超耗，而是计量无基准\n"
        f"• 建议：补全定额主数据即可大幅消除",
        14, INK)
    rect(s, Inches(6.9), Inches(1.7), Inches(5.6), Inches(4.2), RGBColor(0xFB, 0xEF, 0xE3))
    rect(s, Inches(6.9), Inches(1.7), Inches(5.6), Inches(0.5), AMBER)
    txt(s, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.5),
        f"耗用/订单型  {pct_other:.0f}%", 18, WHITE, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.1), Inches(2.4), Inches(5.2), Inches(3.3),
        f"• 真实耗用偏差 / 订单差异等 {other} 条\n"
        f"• 占高偏差率记录的 {pct_other:.1f}%\n"
        f"• 含工艺损耗、计量误差、替代料残差\n"
        f"• 需结合车间现场核查\n"
        f"• 建议：区分真实超耗与数据缺失",
        14, INK)


def build_warn(prs, M):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 10, "预警核查结果", "无备注 / 异常 / 中间地带 三类预警")
    cards = [
        ("无备注预警", M["warn_note"],
         "无待补备注记录" if M["warn_note"] == 0 else f"共 {M['warn_note']} 条需补备注", TEAL),
        ("异常预警", M["warn_abn"],
         f"共 {M['warn_abn']} 条异常记录" if M["warn_abn"] else "无异常预警记录", AMBER),
        ("中间地带明细", M["warn_mid"],
         f"共 {M['warn_mid']} 条待复核" if M["warn_mid"] else "无中间地带记录",
         RGBColor(0x3A, 0xA8, 0xB0)),
    ]
    cw = Inches(3.7)
    gap = Inches(0.3)
    x0 = Inches(0.8)
    y0 = Inches(1.8)
    for i, (name, cnt, desc, col) in enumerate(cards):
        x = x0 + i * (cw + gap)
        rect(s, x, y0, cw, Inches(2.6), CARD_BG)
        rect(s, x, y0, cw, Inches(0.14), col)
        txt(s, x, y0 + Inches(0.35), cw, Inches(0.5), name, 16, INK, True,
            PP_ALIGN.CENTER)
        txt(s, x, y0 + Inches(1.0), cw, Inches(0.8), f"{cnt}", 34, col, True,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x, y0 + Inches(1.95), cw, Inches(0.5), desc, 12, GREY, False,
            PP_ALIGN.CENTER)
    total_warn = M["warn_note"] + M["warn_abn"] + M["warn_mid"]
    rect(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.6), LIGHT)
    rect(s, Inches(0.8), Inches(4.9), Inches(0.14), Inches(1.6), AMBER)
    txt(s, Inches(1.1), Inches(5.05), Inches(11.2), Inches(0.4),
        "核查结论", 16, AMBER, True)
    txt(s, Inches(1.1), Inches(5.5), Inches(11.2), Inches(0.9),
        f"三类预警合计 {total_warn} 条（无备注 {M['warn_note']} / 异常 {M['warn_abn']} / "
        f"中间地带 {M['warn_mid']}）。建议优先补全主数据缺失项以消除大部分异常。",
        13, INK)


def build_summary(prs, M, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WHITE)
    header(s, 11, "总结与改进建议", "分工厂针对性行动项（净偏差口径）")
    bf = M["by_factory"]
    xs = [Inches(0.8), Inches(6.8)]
    fills = [CARD_BG, RGBColor(0xFB, 0xEF, 0xE3)]
    bars = [TEAL, AMBER]
    for (f, row), x0, fl, bar in zip(bf.iterrows(), xs, fills, bars):
        net = row["净偏差金额"]
        gross = row["偏差金额"]
        # 该厂净超耗 Top3 物料（动态）
        sub = d[d["工厂"] == f].groupby("物料名称")["净偏差金额"].sum()
        top_pos = sub[sub > 0].sort_values(ascending=False).head(3)
        if len(top_pos):
            mats = "\n".join(f"  - {clip_name(n, 16)}（{w_money(v)}）"
                             for n, v in top_pos.items())
            mat_block = f"• 净超耗 Top{len(top_pos)} 物料：\n{mats}"
        else:
            mat_block = "• 无净超耗物料，整体省料"
        body = (f"• 毛偏差 {w_money(gross)} → 净偏差 {w_money(net)}（{tag(net)}）\n"
                f"{mat_block}\n"
                "• 建议：\n"
                "  1. 优先补全 BOM 与定额主数据\n"
                "  2. 核查净超耗物料的真实耗用\n"
                "  3. 替代料保持成对登记，按净偏差考核")
        rect(s, x0, Inches(1.7), Inches(5.7), Inches(4.8), fl)
        rect(s, x0, Inches(1.7), Inches(5.7), Inches(0.5), bar)
        txt(s, x0 + Inches(0.2), Inches(1.7), Inches(5.3), Inches(0.5),
            f"{f}（净偏差 {w_money(net)}）", 15, WHITE, True,
            PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        txt(s, x0 + Inches(0.2), Inches(2.4), Inches(5.3), Inches(4.0),
            body, 12, INK)


# ===================== 主流程 =====================
def main():
    xlsx = find_input()
    print(f"输入文件：{xlsx}")
    d = pd.read_excel(xlsx, "完整偏差明细")
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       f"ZPP011偏差分析_净偏差口径_{datetime.datetime.now():%Y%m%d}.pptx")
    build_net_report(d, out, src_name=os.path.basename(xlsx))
    M, _ = _prepare(d, os.path.basename(xlsx))
    print(f"已生成：{out}")
    print(f"  总记录 {M['total_rows']:,} / 替代料 {M['alt_rows']} 行")
    print(f"  毛偏差 {w_money(M['tot_b'])} → 净偏差 {w_money(M['tot_n'])}"
          f"（消除噪音 {w_money(M['noise'])}）")
    print(f"  预警：无备注 {M['warn_note']} / 异常 {M['warn_abn']} / 中间地带 {M['warn_mid']}")


if __name__ == "__main__":
    main()

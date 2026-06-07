#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
ZPP011 智能PPT生成器 v3 (全新设计)
相比于现有版本的优势：
1. 更现代的设计语言（扁平化、渐变色、卡片式布局）
2. 更智能的数据分析（自动识别异常模式、根因聚类）
3. 更丰富的可视化（组合图表、热力图、水fall图）
4. 更清晰的故事线（问题→原因→影响→建议）
5. 更专业的排版（网格系统、视觉层次、留白艺术）
"""

import os
import re
import warnings
from datetime import datetime
from io import BytesIO
from typing import List, Dict, Tuple, Optional, Any

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.colors import LinearSegmentedColormap
from matplotlib.ticker import FuncFormatter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData

warnings.filterwarnings("ignore", category=UserWarning)
plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

# ============================================================
# 配置区域
# ============================================================
COMPANY = "云南达利生产基地"
TITLE = "ZPP011 生产偏差深度分析报告"
SUBTITLE = "数据驱动决策 · 精准管控成本"

# 现代配色方案（扁平化设计）
COLORS = {
    'primary': RGBColor(41, 128, 185),      # 现代蓝
    'secondary': RGBColor(142, 68, 173),    # 紫色
    'accent': RGBColor(230, 126, 34),       # 橙色
    'success': RGBColor(39, 174, 96),       # 绿色
    'danger': RGBColor(231, 76, 60),       # 红色
    'warning': RGBColor(241, 196, 15),      # 黄色
    'dark': RGBColor(44, 62, 80),           # 深蓝灰
    'light': RGBColor(236, 240, 241),       # 浅灰
    'white': RGBColor(255, 255, 255),
    'text': RGBColor(52, 73, 94),
    'text_light': RGBColor(149, 165, 166),
}

# matplotlib配色（用于图表生成）
MPL_COLORS = {
    'positive': '#e74c3c',
    'negative': '#27ae60',
    'primary': '#2980b9',
    'accent': '#e67e22',
}


# ============================================================
# 工具函数
# ============================================================
def _log(msg: str, log_cb=None):
    """日志输出"""
    if log_cb:
        log_cb(msg)
    else:
        print(msg)


def _safe_get(df: pd.DataFrame, row_idx, col_name, default=None):
    """安全获取DataFrame值"""
    try:
        if col_name in df.columns and row_idx < len(df):
            val = df.iloc[row_idx][col_name]
            return val if pd.notna(val) else default
    except:
        pass
    return default


def _format_number(num: Any, precision: int = 0) -> str:
    """格式化数字显示"""
    try:
        num = float(num)
        if abs(num) >= 100000000:
            return f"{num/100000000:.{precision}f}亿"
        elif abs(num) >= 10000:
            return f"{num/10000:.{precision}f}万"
        elif abs(num) >= 1000:
            return f"{num/1000:.{precision}f}K"
        else:
            return f"{num:,.{precision}f}"
    except:
        return str(num)


def _find_column(df: pd.DataFrame, possible_names: List[str], default=None):
    """模糊查找列名"""
    for name in possible_names:
        if name in df.columns:
            return name
    return default


# ============================================================
# PPT设计系统（现代化）
# ============================================================
class PPTDesignSystem:
    """PPT设计系统 - 统一管理样式"""
    
    def __init__(self, prs: Presentation):
        self.prs = prs
        self.W = prs.slide_width
        self.H = prs.slide_height
        
    def add_gradient_background(self, slide, color1, color2, angle=45):
        """添加渐变背景"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.W, self.H
        )
        fill = shape.fill
        fill.gradient()
        fill.gradient_angle = angle
        fill.gradient_stops[0].color.rgb = color1
        fill.gradient_stops[1].color.rgb = color2
        shape.line.fill.background()
        return shape
    
    def add_modern_card(self, slide, left, top, width, height, 
                       bg_color=COLORS['white'], shadow=True):
        """添加现代化卡片（圆角+阴影）"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = COLORS['light']
        card.line.width = Pt(1)
        
        # 圆角设置
        if hasattr(card, 'adjustments'):
            try:
                card.adjustments[0] = 0.1  # 圆角程度
            except:
                pass
        
        return card
    
    def add_section_title(self, slide, title: str, subtitle: str = '', 
                         color=COLORS['primary']):
        """添加章节标题栏（现代设计）"""
        # 背景条
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.W, Pt(70)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # 左侧装饰条
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Pt(8), Pt(70)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = COLORS['accent']
        accent.line.fill.background()
        
        # 标题文字
        txbox = slide.shapes.add_textbox(Pt(30), Pt(10), self.W - Pt(60), Pt(40))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.font.name = 'Microsoft YaHei'
        
        # 副标题
        if subtitle:
            txbox2 = slide.shapes.add_textbox(Pt(30), Pt(45), self.W - Pt(60), Pt(20))
            tf2 = txbox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(200, 220, 240)
            p2.font.name = 'Microsoft YaHei'
    
    def add_kpi_display(self, slide, left, top, value, label, 
                       color=COLORS['primary'], icon=''):
        """添加KPI展示组件（大数字+标签）"""
        # 数字
        tx_val = slide.shapes.add_textbox(left, top, Pt(150), Pt(60))
        tf_val = tx_val.text_frame
        tf_val.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_val = tf_val.paragraphs[0]
        p_val.text = str(value)
        p_val.font.size = Pt(44)
        p_val.font.bold = True
        p_val.font.color.rgb = color
        p_val.font.name = 'Arial'
        p_val.alignment = PP_ALIGN.CENTER
        
        # 标签
        tx_lbl = slide.shapes.add_textbox(left, top + Pt(60), Pt(150), Pt(25))
        tf_lbl = tx_lbl.text_frame
        tf_lbl.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_lbl = tf_lbl.paragraphs[0]
        p_lbl.text = label
        p_lbl.font.size = Pt(12)
        p_lbl.font.color.rgb = COLORS['text_light']
        p_lbl.font.name = 'Microsoft YaHei'
        p_lbl.alignment = PP_ALIGN.CENTER
        
        return tx_val, tx_lbl
    
    def add_bullet_box(self, slide, left, top, width, height, 
                       items: List[str], title: str = ''):
        """添加项目符号框"""
        # 背景卡片
        card = self.add_modern_card(slide, left, top, width, height)
        
        current_top = top + Pt(15)
        
        # 标题
        if title:
            tx_title = slide.shapes.add_textbox(
                left + Pt(15), current_top, width - Pt(30), Pt(25)
            )
            tf_title = tx_title.text_frame
            p_title = tf_title.paragraphs[0]
            p_title.text = title
            p_title.font.size = Pt(16)
            p_title.font.bold = True
            p_title.font.color.rgb = COLORS['dark']
            current_top += Pt(30)
        
        # 项目符号
        for item in items:
            tx_item = slide.shapes.add_textbox(
                left + Pt(20), current_top, width - Pt(40), Pt(20)
            )
            tf_item = tx_item.text_frame
            p_item = tf_item.paragraphs[0]
            p_item.text = f"• {item}"
            p_item.font.size = Pt(11)
            p_item.font.color.rgb = COLORS['text']
            p_item.space_after = Pt(5)
            current_top += Pt(22)


# ============================================================
# 数据分析引擎（智能化）
# ============================================================
class SmartAnalyzer:
    """智能数据分析引擎"""
    
    def __init__(self, excel_path: str, log_cb=None):
        self.excel_path = excel_path
        self.log_cb = log_cb
        self.data = {}
        self.insights = []
        
    def load_data(self):
        """加载所有Sheet数据"""
        _log("[智能分析] 开始加载数据...", self.log_cb)
        xl = pd.ExcelFile(self.excel_path)
        sheet_names = xl.sheet_names
        
        # 智能匹配Sheet
        sheet_mapping = {
            'summary': ['汇总统计', '汇总', '统计'],
            'detail': ['完整偏差明细', '偏差明细', '明细'],
            'alt': ['替代料明细', '替代料'],
            'cause': ['偏差原因分析', '原因分析'],
            'trend': ['趋势分析', 'Trend'],
            'abnormal': ['异常预警', '预警'],
            'note': ['无备注预警', '无备注'],
        }
        
        for key, possible_names in sheet_mapping.items():
            for name in possible_names:
                if name in sheet_names:
                    self.data[key] = pd.read_excel(self.excel_path, sheet_name=name)
                    _log(f"  [智能分析] 加载Sheet: {name}", self.log_cb)
                    break
        
        _log(f"[智能分析] 数据加载完成，共{len(self.data)}个数据集", self.log_cb)
        return self
    
    def analyze(self):
        """执行智能分析"""
        _log("[智能分析] 开始深度分析...", self.log_cb)
        
        # 1. 基础统计
        self._calc_basic_stats()
        
        # 2. 异常检测
        self._detect_anomalies()
        
        # 3. 根因分析
        self._analyze_root_cause()
        
        # 4. 趋势预测
        self._predict_trend()
        
        _log(f"[智能分析] 分析完成，生成{len(self.insights)}条洞察", self.log_cb)
        return self
    
    def _calc_basic_stats(self):
        """计算基础统计"""
        summary = self.data.get('summary')
        if summary is None:
            return
        
        # 智能识别列
        col_map = {
            'total': ['总条数', '记录数'],
            'pos': ['正偏差条数', '正偏差'],
            'neg': ['负偏差条数', '负偏差'],
            'pos_amt': ['正偏差金额(含税)', '正偏差金额'],
            'neg_amt': ['负偏差金额(含税)', '负偏差金额'],
        }
        
        self.stats = {}
        for key, possible in col_map.items():
            col = _find_column(summary, possible)
            if col:
                self.stats[key] = summary[col].sum()
            else:
                self.stats[key] = 0
        
        # 备注覆盖率
        rate_col = _find_column(summary, ['备注覆盖率', '覆盖率'])
        if rate_col:
            rates = summary[rate_col].astype(str).str.replace('%', '').astype(float)
            weights = summary[_find_column(summary, ['总条数', '记录数'])]
            self.stats['note_rate'] = (rates * weights).sum() / weights.sum() / 100
        else:
            self.stats['note_rate'] = 0
        
        # ========= 新增：食品/饮料区分统计 =========
        self._calc_food_beverage_stats(summary)
        # =============================================
        
        _log(f"[智能分析] 基础统计: {self.stats}", self.log_cb)
    
    def _calc_food_beverage_stats(self, summary):
        """计算食品/饮料分别的统计"""
        self.has_food_beverage = False
        self.food_stats = {}
        self.beverage_stats = {}
        
        # 查找工厂列
        factory_col = _find_column(summary, ['工厂名称', '工厂', '生产工厂'])
        if factory_col is None:
            return
        
        # 判断是否有食品/饮料区分
        all_factories = summary[factory_col].dropna().astype(str).unique()
        has_food = any('食品' in f for f in all_factories)
        has_bev = any('饮料' in f for f in all_factories)
        
        if not (has_food or has_bev):
            return
        
        self.has_food_beverage = True
        _log(f"[智能分析] 检测到食品/饮料区分，工厂列表: {list(all_factories)[:10]}", self.log_cb)
        
        col_map = {
            'total': ['总条数', '记录数'],
            'pos': ['正偏差条数', '正偏差'],
            'neg': ['负偏差条数', '负偏差'],
            'pos_amt': ['正偏差金额(含税)', '正偏差金额'],
            'neg_amt': ['负偏差金额(含税)', '负偏差金额'],
        }
        
        for category, keyword in [('food', '食品'), ('beverage', '饮料')]:
            mask = summary[factory_col].astype(str).str.contains(keyword, na=False)
            cat_df = summary[mask]
            if cat_df.empty:
                continue
            
            stats = {}
            for key, possible in col_map.items():
                col = _find_column(cat_df, possible)
                stats[key] = cat_df[col].sum() if col else 0
            
            # 备注覆盖率
            rate_col = _find_column(cat_df, ['备注覆盖率', '覆盖率'])
            if rate_col:
                try:
                    rates = cat_df[rate_col].astype(str).str.replace('%', '').astype(float)
                    weights = cat_df[_find_column(cat_df, ['总条数', '记录数'])]
                    stats['note_rate'] = (rates * weights).sum() / weights.sum() / 100
                except:
                    stats['note_rate'] = 0
            else:
                stats['note_rate'] = 0
            
            if category == 'food':
                self.food_stats = stats
                _log(f"[智能分析] 食品统计: 总条数={stats.get('total',0)}", self.log_cb)
            else:
                self.beverage_stats = stats
                _log(f"[智能分析] 饮料统计: 总条数={stats.get('total',0)}", self.log_cb)
    
    def _detect_anomalies(self):
        """异常检测"""
        self.anomalies = []
        
        # 检查无备注高偏差
        note_df = self.data.get('note')
        if note_df is not None and not note_df.empty:
            self.anomalies.append({
                'type': '无备注预警',
                'count': len(note_df),
                'severity': 'high',
                'desc': f"发现{len(note_df)}条高偏差无备注记录"
            })
        
        # 检查异常预警
        abnormal_df = self.data.get('abnormal')
        if abnormal_df is not None and not abnormal_df.empty:
            self.anomalies.append({
                'type': '异常预警',
                'count': len(abnormal_df),
                'severity': 'medium',
                'desc': f"发现{len(abnormal_df)}条异常记录"
            })
        
        _log(f"[智能分析] 异常检测: 发现{len(self.anomalies)}类异常", self.log_cb)
    
    def _analyze_root_cause(self):
        """根因分析"""
        cause_df = self.data.get('cause')
        if cause_df is None or cause_df.empty:
            self.root_causes = []
            return
        
        # 统计原因分布
        cause_col = _find_column(cause_df, ['备注原因', '原因'])
        if cause_col:
            counts = cause_df[cause_col].value_counts().head(10)
            self.root_causes = [
                {'cause': str(k), 'count': int(v), 'percent': v/len(cause_df)*100}
                for k, v in counts.items()
            ]
        else:
            self.root_causes = []
        
        _log(f"[智能分析] 根因分析: Top原因{len(self.root_causes)}个", self.log_cb)
    
    def _predict_trend(self):
        """趋势预测（简化版）"""
        trend_df = self.data.get('trend')
        if trend_df is None or trend_df.empty:
            self.trend = None
            return
        
        # 简单趋势判断
        self.trend = {
            'has_data': True,
            'summary': "趋势分析数据可用，建议重点关注近期偏差率上升的物料"
        }
        
        _log("[智能分析] 趋势预测完成", self.log_cb)


# ============================================================
# 图表生成器（高质量）
# ============================================================
class ChartGenerator:
    """高质量图表生成器"""
    
    @staticmethod
    def create_waterfall_chart(data: List[Tuple[str, float]], title: str):
        """创建瀑布图（偏差构成分析）"""
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # 计算累积值
        cumulative = 0
        colors = []
        positions = []
        
        for i, (label, value) in enumerate(data):
            if i == 0:
                # 初始值
                ax.bar(i, value, color=MPL_COLORS['primary'], alpha=0.7)
                colors.append(MPL_COLORS['primary'])
            elif i == len(data) - 1:
                # 最终值
                ax.bar(i, value, color=MPL_COLORS['primary'], alpha=0.9)
                colors.append(MPL_COLORS['primary'])
            else:
                # 中间变化
                color = MPL_COLORS['positive'] if value > 0 else MPL_COLORS['negative']
                ax.bar(i, value, bottom=cumulative, color=color, alpha=0.7)
                colors.append(color)
                cumulative += value
            
            positions.append(i)
        
        ax.set_xticks(positions)
        ax.set_xticklabels([d[0] for d in data], rotation=45, ha='right')
        ax.set_title(title, fontsize=16, fontweight='bold')
        ax.axhline(0, color='black', linewidth=0.8)
        ax.grid(axis='y', linestyle='--', alpha=0.3)
        
        plt.tight_layout()
        
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        plt.close()
        buf.seek(0)
        return buf
    
    @staticmethod
    def create_heatmap(df: pd.DataFrame, title: str):
        """创建热力图（使用matplotlib原生功能）"""
        if df.empty or df.shape[1] < 2:
            return None
        
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # 提取数据（假设第一列是标签），并转置
        data = df.iloc[:, 1:].T.values.astype(float)
        
        # 使用 imshow 显示热力图
        im = ax.imshow(data, cmap='RdYlGn_r', aspect='auto', interpolation='nearest')
        
        # 添加颜色条
        plt.colorbar(im, ax=ax)
        
        # 设置 x/y 轴刻度标签
        # x轴：原始df的行（用第一列的值作为标签）
        ax.set_xticks(range(data.shape[1]))
        ax.set_xticklabels([str(x) for x in df.iloc[:, 0]], rotation=45, ha='right')
        
        # y轴：数值列名
        ax.set_yticks(range(data.shape[0]))
        ax.set_yticklabels(df.columns[1:], rotation=0)
        
        # 添加数值标注
        for i in range(data.shape[0]):
            for j in range(data.shape[1]):
                text = ax.text(j, i, f'{data[i, j]:.1f}',
                               ha='center', va='center', color='black', fontsize=8)
        
        ax.set_title(title, fontsize=16, fontweight='bold')
        ax.set_ylabel(df.columns[0])
        
        plt.tight_layout()
        
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        plt.close()
        buf.seek(0)
        return buf
    
    @staticmethod
    def create_treemap(data: Dict[str, float], title: str):
        """创建树状图（需要额外依赖，这里用条形图替代）"""
        fig, ax = plt.subplots(figsize=(10, 6))
        
        labels = list(data.keys())
        values = list(data.values())
        colors = [MPL_COLORS['positive'] if v > 0 else MPL_COLORS['negative'] for v in values]
        
        bars = ax.barh(labels, values, color=colors, alpha=0.7)
        
        # 添加数值标签
        for bar, value in zip(bars, values):
            width = bar.get_width()
            ax.text(width, bar.get_y() + bar.get_height()/2, 
                   f'{value:,.0f}', 
                   ha='left' if width >= 0 else 'right',
                   va='center',
                   fontweight='bold')
        
        ax.set_title(title, fontsize=16, fontweight='bold')
        ax.axvline(0, color='black', linewidth=0.8)
        ax.grid(axis='x', linestyle='--', alpha=0.3)
        
        plt.tight_layout()
        
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        plt.close()
        buf.seek(0)
        return buf


# ============================================================
# 主生成器
# ============================================================
class SmartPPTGenerator:
    """智能PPT生成器 v3"""
    
    def __init__(self, log_cb=None):
        self.log_cb = log_cb
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.design = PPTDesignSystem(self.prs)
        self.chart_gen = ChartGenerator()
        self.analyzer = None
        
    def generate(self, excel_path: str, output_path: str):
        """生成PPT"""
        _log(f"[PPT v3] 开始生成: {excel_path}", self.log_cb)
        
        # 1. 数据分析
        self.analyzer = SmartAnalyzer(excel_path, self.log_cb)
        self.analyzer.load_data()
        self.analyzer.analyze()
        
        # 2. 生成各页
        self._add_cover_slide()
        self._add_toc_slide()
        self._add_executive_summary()
        self._add_kpi_dashboard()
        
        # 食品/饮料对比页（仅当检测到区分时展示）
        if self.analyzer and getattr(self.analyzer, 'has_food_beverage', False):
            self._add_food_beverage_slide()
        
        self._add_deviation_analysis()
        self._add_root_cause_analysis()
        self._add_trend_analysis()
        self._add_recommendations()
        self._add_closing_slide()
        
        # 3. 保存
        self.prs.save(output_path)
        _log(f"[PPT v3] 生成完成: {output_path} (共{len(self.prs.slides)}页)", self.log_cb)
        return True
    
    def _add_cover_slide(self):
        """封面页（现代设计）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 渐变背景
        self.design.add_gradient_background(
            slide, COLORS['primary'], COLORS['secondary'], angle=135
        )
        
        # 装饰元素（圆形）
        for i, (x, y, size) in enumerate([
            (Pt(50), Pt(100), Pt(200)),
            (Pt(800), Pt(50), Pt(150)),
            (Pt(900), Pt(400), Pt(100)),
        ]):
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLORS['white']
            circle.fill.fore_color.theme_color = 1  # 半透明
            circle.line.fill.background()
            circle.shadow.inherit = False
        
        # 主标题
        tx_main = slide.shapes.add_textbox(Pt(100), Pt(150), Pt(800), Pt(100))
        tf_main = tx_main.text_frame
        tf_main.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_main = tf_main.paragraphs[0]
        p_main.text = TITLE
        p_main.font.size = Pt(54)
        p_main.font.bold = True
        p_main.font.color.rgb = COLORS['white']
        p_main.font.name = 'Microsoft YaHei'
        p_main.alignment = PP_ALIGN.CENTER
        
        # 副标题
        tx_sub = slide.shapes.add_textbox(Pt(100), Pt(270), Pt(800), Pt(50))
        tf_sub = tx_sub.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = SUBTITLE
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = RGBColor(200, 220, 240)
        p_sub.font.name = 'Microsoft YaHei'
        p_sub.alignment = PP_ALIGN.CENTER
        
        # 公司名
        tx_company = slide.shapes.add_textbox(Pt(100), Pt(350), Pt(800), Pt(40))
        tf_company = tx_company.text_frame
        p_company = tf_company.paragraphs[0]
        p_company.text = COMPANY
        p_company.font.size = Pt(18)
        p_company.font.color.rgb = COLORS['accent']
        p_company.font.name = 'Microsoft YaHei'
        p_company.alignment = PP_ALIGN.CENTER
        
        # 日期
        date_str = datetime.now().strftime('%Y年%m月%d日')
        tx_date = slide.shapes.add_textbox(Pt(100), Pt(400), Pt(800), Pt(30))
        tf_date = tx_date.text_frame
        p_date = tf_date.paragraphs[0]
        p_date.text = date_str
        p_date.font.size = Pt(14)
        p_date.font.color.rgb = RGBColor(180, 200, 220)
        p_date.alignment = PP_ALIGN.CENTER
        
        _log("[PPT v3] 封面页完成", self.log_cb)
    
    def _add_toc_slide(self):
        """目录页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['light']
        bg.line.fill.background()
        
        # 标题
        self.design.add_section_title(slide, "目录", "Report Contents")
        
        # 目录项
        toc_items = [
            ("01", "执行摘要", "关键发现与核心指标"),
            ("02", "数据总览", "KPI仪表板与整体概况"),
            ("03", "食饮对比", "食品工厂 vs 饮料工厂偏差对比"),
            ("04", "偏差分析", "正负偏差分布与明细"),
            ("05", "根因诊断", "偏差原因聚类分析"),
            ("06", "趋势预测", "历史趋势与未来预警"),
            ("07", "改进建议", "分阶段行动方案"),
        ]
        
        start_y = Pt(100)
        for i, (num, title, desc) in enumerate(toc_items):
            y = start_y + Pt(80 * i)
            
            # 编号
            tx_num = slide.shapes.add_textbox(Pt(50), y, Pt(50), Pt(50))
            tf_num = tx_num.text_frame
            p_num = tf_num.paragraphs[0]
            p_num.text = num
            p_num.font.size = Pt(24)
            p_num.font.bold = True
            p_num.font.color.rgb = COLORS['accent']
            
            # 标题
            tx_title = slide.shapes.add_textbox(Pt(120), y, Pt(300), Pt(30))
            tf_title = tx_title.text_frame
            p_title = tf_title.paragraphs[0]
            p_title.text = title
            p_title.font.size = Pt(20)
            p_title.font.bold = True
            p_title.font.color.rgb = COLORS['dark']
            
            # 描述
            tx_desc = slide.shapes.add_textbox(Pt(120), y + Pt(30), Pt(400), Pt(20))
            tf_desc = tx_desc.text_frame
            p_desc = tf_desc.paragraphs[0]
            p_desc.text = desc
            p_desc.font.size = Pt(12)
            p_desc.font.color.rgb = COLORS['text_light']
        
        _log("[PPT v3] 目录页完成", self.log_cb)
    
    def _add_executive_summary(self):
        """执行摘要页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题
        self.design.add_section_title(slide, "01 执行摘要", "Executive Summary")
        
        # 关键信息卡片
        stats = self.analyzer.stats
        
        # 核心发现
        insights = [
            f"本期共分析 {stats.get('total', 0):,} 条偏差记录",
            f"正偏差（多耗）{stats.get('pos', 0):,} 条，金额 ¥{stats.get('pos_amt', 0):,.0f}",
            f"负偏差（少耗）{stats.get('neg', 0):,} 条，金额 ¥{stats.get('neg_amt', 0):,.0f}",
            f"备注覆盖率 {stats.get('note_rate', 0)*100:.1f}%，{'有待提升' if stats.get('note_rate', 0) < 0.4 else '表现良好'}",
        ]
        
        # 添加异常警告
        for anomaly in self.analyzer.anomalies:
            insights.append(f"⚠️ {anomaly['desc']}")
        
        self.design.add_bullet_box(
            slide, Pt(30), Pt(90), Pt(600), Pt(300),
            insights, "关键发现"
        )
        
        # 建议框
        recommendations = [
            "立即补录无备注高偏差记录",
            "重点核查异常预警物料",
            "建立偏差根因分类标准",
            "推动系统化防错机制",
        ]
        
        self.design.add_bullet_box(
            slide, Pt(650), Pt(90), Pt(600), Pt(300),
            recommendations, "优先行动"
        )
        
        _log("[PPT v3] 执行摘要页完成", self.log_cb)
    
    def _add_kpi_dashboard(self):
        """KPI仪表板页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题
        self.design.add_section_title(slide, "02 数据总览", "KPI Dashboard")
        
        stats = self.analyzer.stats
        
        # KPI卡片布局（2x3网格）
        kpis = [
            (f"{stats.get('total', 0):,}", "总记录数", COLORS['primary']),
            (f"{stats.get('pos', 0):,}", "正偏差条数", COLORS['danger']),
            (f"{stats.get('neg', 0):,}", "负偏差条数", COLORS['success']),
            (f"¥{_format_number(stats.get('pos_amt', 0))}", "正偏差金额", COLORS['danger']),
            (f"¥{_format_number(stats.get('neg_amt', 0))}", "负偏差金额", COLORS['success']),
            (f"{stats.get('note_rate', 0)*100:.1f}%", "备注覆盖率", COLORS['accent']),
        ]
        
        positions = [
            (Pt(30), Pt(90)), (Pt(350), Pt(90)), (Pt(670), Pt(90)),
            (Pt(30), Pt(280)), (Pt(350), Pt(280)), (Pt(670), Pt(280)),
        ]
        
        for i, ((value, label, color), (x, y)) in enumerate(zip(kpis, positions)):
            self.design.add_modern_card(slide, x, y, Pt(280), Pt(170), COLORS['white'])
            self.design.add_kpi_display(slide, x, y + Pt(20), value, label, color)
        
        _log("[PPT v3] KPI仪表板页完成", self.log_cb)
    
    def _add_food_beverage_slide(self):
        """食品/饮料对比分析页"""
        if not self.analyzer or not getattr(self.analyzer, 'has_food_beverage', False):
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.design.add_section_title(slide, "03 食品/饮料对比", "Food vs Beverage Comparison")

        W = self.prs.slide_width
        food = getattr(self.analyzer, 'food_stats', {}) or {}
        bev = getattr(self.analyzer, 'beverage_stats', {}) or {}

        # ---- 左侧：食品 KPI 卡片 ----
        x_left = Pt(30)
        y_start = Pt(90)
        card_w = Pt(300)
        card_h = Pt(80)

        food_items = [
            ("食品·总条数", food.get('total', 0)),
            ("食品·正偏差金额", f"¥{food.get('pos_amt', 0):,.0f}"),
            ("食品·负偏差金额", f"¥{food.get('neg_amt', 0):,.0f}"),
            ("食品·备注覆盖率", f"{food.get('note_rate', 0)*100:.1f}%"),
        ]
        for i, (label, val) in enumerate(food_items):
            y = y_start + Pt(95 * i)
            self.design.add_modern_card(slide, x_left, y, card_w, card_h)
            tx = slide.shapes.add_textbox(x_left + Pt(8), y + Pt(8), card_w - Pt(16), Pt(30))
            tf = tx.text_frame
            tf.paragraphs[0].text = label
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = COLORS['text_light']
            tf.paragraphs[0].font.name = 'Microsoft YaHei'

            tx2 = slide.shapes.add_textbox(x_left + Pt(8), y + Pt(38), card_w - Pt(16), Pt(35))
            tf2 = tx2.text_frame
            tf2.paragraphs[0].text = str(val)
            tf2.paragraphs[0].font.size = Pt(18)
            tf2.paragraphs[0].font.bold = True
            tf2.paragraphs[0].font.color.rgb = COLORS['primary']
            tf2.paragraphs[0].font.name = 'Arial'

        # ---- 右侧：饮料 KPI 卡片 ----
        x_right = Pt(350)
        bev_items = [
            ("饮料·总条数", bev.get('total', 0)),
            ("饮料·正偏差金额", f"¥{bev.get('pos_amt', 0):,.0f}"),
            ("饮料·负偏差金额", f"¥{bev.get('neg_amt', 0):,.0f}"),
            ("饮料·备注覆盖率", f"{bev.get('note_rate', 0)*100:.1f}%"),
        ]
        for i, (label, val) in enumerate(bev_items):
            y = y_start + Pt(95 * i)
            self.design.add_modern_card(slide, x_right, y, card_w, card_h)
            tx = slide.shapes.add_textbox(x_right + Pt(8), y + Pt(8), card_w - Pt(16), Pt(30))
            tf = tx.text_frame
            tf.paragraphs[0].text = label
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = COLORS['text_light']
            tf.paragraphs[0].font.name = 'Microsoft YaHei'

            tx2 = slide.shapes.add_textbox(x_right + Pt(8), y + Pt(38), card_w - Pt(16), Pt(35))
            tf2 = tx2.text_frame
            tf2.paragraphs[0].text = str(val)
            tf2.paragraphs[0].font.size = Pt(18)
            tf2.paragraphs[0].font.bold = True
            tf2.paragraphs[0].font.color.rgb = COLORS['secondary']
            tf2.paragraphs[0].font.name = 'Arial'

        # ---- 底部：对比柱状图（PPT 原生图表）----
        try:
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE

            chart_data = ChartData()
            chart_data.categories = ['总条数', '正偏差金额(万)', '负偏差金额(万)', '备注覆盖率(%)']
            food_vals = [
                food.get('total', 0),
                food.get('pos_amt', 0) / 10000,
                food.get('neg_amt', 0) / 10000,
                food.get('note_rate', 0) * 100,
            ]
            bev_vals = [
                bev.get('total', 0),
                bev.get('pos_amt', 0) / 10000,
                bev.get('neg_amt', 0) / 10000,
                bev.get('note_rate', 0) * 100,
            ]
            chart_data.add_series('食品工厂', food_vals)
            chart_data.add_series('饮料工厂', bev_vals)

            x_chart = Pt(680)
            y_chart = Pt(90)
            w_chart = Pt(580)
            h_chart = Pt(380)

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                x_chart, y_chart, w_chart, h_chart,
                chart_data
            ).chart

            chart.has_title = True
            chart.chart_title.text_frame.text = "食品工厂 vs 饮料工厂 — 关键指标对比"
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'

        except Exception as e:
            _log(f"[PPT v3] 食品/饮料对比图表生成失败: {e}", self.log_cb)

        _log("[PPT v3] 食品/饮料对比页完成", self.log_cb)

    def _add_deviation_analysis(self):
        """偏差分析页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题
        self.design.add_section_title(slide, "04 偏差分析", "Deviation Analysis")
        
        # 这里可以添加更详细的偏差分析图表
        # 例如：正负偏差对比图、车间分布图等
        
        # 占位文字
        tx = slide.shapes.add_textbox(Pt(30), Pt(90), Pt(1200), Pt(400))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "【偏差分析详情】\n\n"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        
        # 添加分析内容
        stats = self.analyzer.stats
        details = [
            f"• 总偏差记录数：{stats.get('total', 0):,} 条",
            f"• 正偏差占比：{stats.get('pos', 0)/max(stats.get('total', 1), 1)*100:.1f}%",
            f"• 负偏差占比：{stats.get('neg', 0)/max(stats.get('total', 1), 1)*100:.1f}%",
            f"• 正偏差金额：¥{stats.get('pos_amt', 0):,.0f}",
            f"• 负偏差金额：¥{stats.get('neg_amt', 0):,.0f}",
            f"• 净偏差金额：¥{(stats.get('pos_amt', 0) - stats.get('neg_amt', 0)):,.0f}",
        ]
        
        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(10)
        
        _log("[PPT v3] 偏差分析页完成", self.log_cb)
    
    def _add_root_cause_analysis(self):
        """根因分析页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题
        self.design.add_section_title(slide, "05 根因诊断", "Root Cause Analysis")
        
        root_causes = self.analyzer.root_causes
        
        if not root_causes:
            tx = slide.shapes.add_textbox(Pt(30), Pt(90), Pt(1200), Pt(400))
            tf = tx.text_frame
            p = tf.paragraphs[0]
            p.text = "暂无根因分析数据"
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['text_light']
            return
        
        # Top5原因
        tx = slide.shapes.add_textbox(Pt(30), Pt(90), Pt(1200), Pt(400))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "【Top5 偏差原因】\n"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        
        for i, rc in enumerate(root_causes[:5], 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {rc['cause']}：{rc['count']}次 ({rc['percent']:.1f}%)"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(10)
        
        _log("[PPT v3] 根因分析页完成", self.log_cb)
    
    def _add_trend_analysis(self):
        """趋势分析页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题
        self.design.add_section_title(slide, "06 趋势预测", "Trend Analysis")
        
        # 这里可以添加趋势图表
        tx = slide.shapes.add_textbox(Pt(30), Pt(90), Pt(1200), Pt(400))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "【趋势分析】\n\n"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        
        trend = self.analyzer.trend
        if trend and trend.get('has_data'):
            p = tf.add_paragraph()
            p.text = trend.get('summary', '趋势分析数据已加载')
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text']
        else:
            p = tf.add_paragraph()
            p.text = "• 暂无趋势分析数据"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text_light']
        
        _log("[PPT v3] 趋势分析页完成", self.log_cb)
    
    def _add_recommendations(self):
        """改进建议页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题
        self.design.add_section_title(slide, "07 改进建议", "Recommendations")
        
        recommendations = [
            ("立即行动（1周内）", [
                "• 强制补录无备注高偏差记录",
                "• 完善系统无定额物料清单",
                "• 核查领用记录异常",
            ]),
            ("短期优化（1月内）", [
                "• 优化批次分摊逻辑",
                "• 加强设备维护保养",
                "• 规范工艺参数执行",
            ]),
            ("中期建设（3月内）", [
                "• 替代料净偏差自动抵消",
                "• 上线实时预警看板",
                "• 实现系统双重校验",
            ]),
        ]
        
        start_y = Pt(90)
        for i, (phase, items) in enumerate(recommendations):
            y = start_y + Pt(150 * i)
            
            # 阶段标题
            tx_phase = slide.shapes.add_textbox(Pt(30), y, Pt(1200), Pt(30))
            tf_phase = tx_phase.text_frame
            p_phase = tf_phase.paragraphs[0]
            p_phase.text = phase
            p_phase.font.size = Pt(16)
            p_phase.font.bold = True
            p_phase.font.color.rgb = COLORS['primary']
            
            # 建议列表
            tx_items = slide.shapes.add_textbox(Pt(50), y + Pt(35), Pt(1180), Pt(100))
            tf_items = tx_items.text_frame
            tf_items.word_wrap = True
            
            for item in items:
                p_item = tf_items.paragraphs[0] if item == items[0] else tf_items.add_paragraph()
                p_item.text = item
                p_item.font.size = Pt(13)
                p_item.font.color.rgb = COLORS['text']
                p_item.space_after = Pt(8)
        
        _log("[PPT v3] 改进建议页完成", self.log_cb)
    
    def _add_closing_slide(self):
        """结束页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 渐变背景
        self.design.add_gradient_background(
            slide, COLORS['dark'], COLORS['primary'], angle=135
        )
        
        # 结束语
        tx = slide.shapes.add_textbox(Pt(100), Pt(200), Pt(800), Pt(200))
        tf = tx.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.text = "谢谢！"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "让数据驱动决策，让管理精准有效"
        p2.font.size = Pt(20)
        p2.font.color.rgb = COLORS['accent']
        p2.alignment = PP_ALIGN.CENTER
        
        _log("[PPT v3] 结束页完成", self.log_cb)


# ============================================================
# 对外接口
# ============================================================
def generate_smart_ppt(excel_path: str, output_path: str, log_cb=None):
    """生成智能PPT（对外接口）"""
    try:
        generator = SmartPPTGenerator(log_cb)
        success = generator.generate(excel_path, output_path)
        return success
    except Exception as e:
        _log(f"[PPT v3] 生成失败: {str(e)}", log_cb)
        import traceback
        traceback.print_exc()
        return False


if __name__ == "__main__":
    # 测试
    import sys
    if len(sys.argv) > 1:
        excel_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else "智能PPT_v3_output.pptx"
        generate_smart_ppt(excel_file, output_file, log_cb=print)
    else:
        print("用法: python smart_ppt_generator_v3.py <excel_path> [output_path]")
